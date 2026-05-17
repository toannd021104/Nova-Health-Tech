"""Build a single Customer Scenario slide matching the existing deck theme."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(".")

# Theme colors (from existing deck)
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK  = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED     = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bordered(shape, fill_color, line_color, line_pt=1.5, radius=0.06):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_pt)
    shape.adjustments[0] = radius


def text_box(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    color = color or GRAY_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


# ── White background ──────────────────────────────────────────────────────────
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
solid(bg, WHITE)

# ── Orange top bar ────────────────────────────────────────────────────────────
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
solid(bar, ORANGE)

# ── Slide title ───────────────────────────────────────────────────────────────
text_box(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.65),
         "Customer Scenario", size=28, bold=True, color=ORANGE_DARK)

# Orange underline
ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(0.92), Inches(2.4), Inches(0.04))
solid(ul, ORANGE)

# ── Footer ────────────────────────────────────────────────────────────────────
text_box(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28),
         "Nova Health Tech  ·  Clinical GenAI Assistant  ·  AWS Singapore",
         size=9, italic=True, color=GRAY_MED, align=PP_ALIGN.RIGHT)

# ── Vertical divider ──────────────────────────────────────────────────────────
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(6.15), Inches(1.05), Inches(0.03), Inches(5.2))
solid(div, ORANGE_LIGHT)

# =============================================================================
# LEFT COLUMN  (x = 0.5 .. 5.9)
# =============================================================================

# ── Company context card ──────────────────────────────────────────────────────
ctx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(1.1), Inches(5.5), Inches(1.6))
bordered(ctx, ORANGE_LIGHT, ORANGE, line_pt=1.5)

tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(5.1), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run()
r.text = "Nova Health Tech"
r.font.name = "Calibri"; r.font.size = Pt(14); r.font.bold = True
r.font.color.rgb = ORANGE_DARK

p2 = tf.add_paragraph()
p2.space_before = Pt(4)
r2 = p2.add_run()
r2.text = ("Digital Health / Clinical Decision Support company. "
           "Flagship clinical support tool is struggling to meet physician "
           "expectations for speed and medical relevance. "
           "Executive board approves a GenAI assistant initiative.")
r2.font.name = "Calibri"; r2.font.size = Pt(10.5)
r2.font.color.rgb = GRAY_DARK

# ── Challenge statement card ──────────────────────────────────────────────────
chal = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(2.88), Inches(5.5), Inches(1.3))
bordered(chal, WHITE, ORANGE_DARK, line_pt=2)

ch_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(2.88), Inches(5.5), Inches(0.44))
solid(ch_hdr, ORANGE_DARK)

text_box(slide, Inches(0.5), Inches(2.88), Inches(5.5), Inches(0.44),
         "The Challenge", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

text_box(slide, Inches(0.7), Inches(3.38), Inches(5.1), Inches(0.75),
         ("Build a GenAI assistant for internal clinical staff and hospital clients "
          "that is fast, accurate, auditable, and compliant with healthcare regulations."),
         size=11, color=GRAY_DARK)

# ── Requirements ─────────────────────────────────────────────────────────────
text_box(slide, Inches(0.5), Inches(4.35), Inches(5.5), Inches(0.38),
         "The assistant must:", size=13, bold=True, color=ORANGE_DARK)

reqs = [
    "Answer complex medical questions in natural language",
    "Rely on internal trials, treatment protocols, WHO, PubMed",
    "Be auditable, compliant, and fast enough for diagnosis",
]
for i, req in enumerate(reqs):
    rb = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.82 + i * 0.38), Inches(5.5), Inches(0.36))
    tf = rb.text_frame
    p = tf.paragraphs[0]
    bullet = p.add_run()
    bullet.text = "\u25b8 "
    bullet.font.name = "Calibri"; bullet.font.size = Pt(12)
    bullet.font.bold = True; bullet.font.color.rgb = ORANGE
    body = p.add_run()
    body.text = req
    body.font.name = "Calibri"; body.font.size = Pt(11)
    body.font.color.rgb = GRAY_DARK

# =============================================================================
# RIGHT COLUMN  (x = 6.3 .. 12.8)
# =============================================================================

text_box(slide, Inches(6.3), Inches(1.1), Inches(6.5), Inches(0.4),
         "Use Case Scenarios", size=15, bold=True, color=ORANGE_DARK)

scenarios = [
    ("\U0001f550", "Emergency Speed",
     "Emergency care needs 2-second response time",
     RED, RED_LIGHT),
    ("\U0001f4c5", "WHO Monthly Updates",
     "WHO publishes monthly protocol updates via structured API",
     ORANGE_DARK, ORANGE_LIGHT),
    ("\U0001f512", "Patient-Sensitive Data",
     "Internal trials include patient-sensitive data (PHI)",
     RED, RED_LIGHT),
    ("\u270d\ufe0f", "Consistent Tone",
     "Users want consistent clinical tone and phrasing across answers",
     ORANGE, ORANGE_LIGHT),
    ("\U0001f4c4", "Legacy PDF Ingestion",
     "Internal trial reports in legacy PDFs with inconsistent tagging",
     ORANGE_DARK, ORANGE_LIGHT),
]

card_h = Inches(1.02)
gap    = Inches(0.08)
start_y = Inches(1.6)
card_w  = Inches(6.5)

for i, (icon, title, desc, color, bg_color) in enumerate(scenarios):
    y = start_y + i * (card_h + gap)
    x = Inches(6.3)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    bordered(card, bg_color, color, line_pt=1.5, radius=0.08)

    # Left accent bar
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
    solid(acc, color)

    # Icon
    icon_tb = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.1), Inches(0.65), Inches(0.82))
    tf = icon_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(28)

    # Title
    title_tb = slide.shapes.add_textbox(
        x + Inches(0.88), y + Inches(0.1), Inches(5.5), Inches(0.38))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"; run.font.size = Pt(13)
    run.font.bold = True; run.font.color.rgb = color

    # Description
    desc_tb = slide.shapes.add_textbox(
        x + Inches(0.88), y + Inches(0.5), Inches(5.5), Inches(0.48))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"; run.font.size = Pt(10.5)
    run.font.color.rgb = GRAY_DARK

# =============================================================================
# BOTTOM BANNER — Deliverables
# =============================================================================

banner_y = Inches(6.42)
banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.5), banner_y, Inches(12.3), Inches(0.82))
solid(banner, ORANGE)
banner.adjustments[0] = 0.1

banner_tb = slide.shapes.add_textbox(
    Inches(0.7), banner_y + Inches(0.06), Inches(11.9), Inches(0.72))
tf = banner_tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Architecture Deliverables Required"
run.font.name = "Calibri"; run.font.size = Pt(13)
run.font.bold = True; run.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = ("Data Pipeline  \u00b7  Model Orchestration  \u00b7  Security Architecture  \u00b7  "
             "Deployment Approach (RAG vs Fine-tuning, Cloud/Hybrid)  \u00b7  Performance Optimization")
run2.font.name = "Calibri"; run2.font.size = Pt(11)
run2.font.color.rgb = WHITE

# =============================================================================
# Save
# =============================================================================
out = "docs/Customer_Scenario_Slide.pptx"
prs.save(out)
size = Path(out).stat().st_size
print(f"Saved: {out}  ({size:,} bytes)")
