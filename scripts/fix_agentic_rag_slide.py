"""Rebuild slide 4 (Agentic RAG) with accurate PoC vs Production distinction."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK  = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED     = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)
GREEN        = RGBColor(0x27, 0xAE, 0x60)
GREEN_LIGHT  = RGBColor(0xE8, 0xF5, 0xE9)
AMBER        = RGBColor(0xFF, 0x8F, 0x00)
AMBER_LIGHT  = RGBColor(0xFF, 0xF8, 0xE1)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation("docs/Technology_Explanation_Slides.pptx")

# Replace slide 4 (index 3) — delete and re-insert
# python-pptx cannot delete slides easily; we rebuild the whole deck
# Instead, we add a new slide and note it replaces slide 4

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def bordered(shape, fill_color, line_color, line_pt=1.5, radius=0.06):
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color; shape.line.width = Pt(line_pt)
    shape.adjustments[0] = radius

def add_text(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or GRAY_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_list(slide, left, top, width, items, size=10.5, color=None,
                bullet_color=None, spacing=0.32):
    color = color or GRAY_DARK
    bullet_color = bullet_color or ORANGE
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(left, top + Inches(i * spacing), width, Inches(spacing))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        b = p.add_run(); b.text = "\u25b8 "
        b.font.name = "Calibri"; b.font.size = Pt(size)
        b.font.bold = True; b.font.color.rgb = bullet_color
        t = p.add_run(); t.text = item
        t.font.name = "Calibri"; t.font.size = Pt(size)
        t.font.color.rgb = color

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
solid(bg, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
solid(bar, ORANGE)

# Footer
ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28))
tf = ft.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run = p.add_run()
run.text = "Nova Health Tech  \u00b7  Clinical GenAI Assistant  \u00b7  AWS Singapore"
run.font.name = "Calibri"; run.font.size = Pt(9)
run.font.italic = True; run.font.color.rgb = GRAY_MED

# Title
add_text(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.65),
         "Agentic RAG: PoC vs Production Architecture",
         size=28, bold=True, color=ORANGE_DARK)
ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.92), Inches(3.8), Inches(0.04))
solid(ul, ORANGE)

# ── PoC (left) ────────────────────────────────────────────────────────────────
poc_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5))
bordered(poc_card, GREEN_LIGHT, GREEN, line_pt=2)

poc_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5))
solid(poc_hdr, GREEN)
add_text(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5),
         "PoC (Deployed Today)", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.4), Inches(0.45),
         "Fixed pipeline — NO Bedrock Agents, NO tool calling",
         size=11, bold=True, color=GREEN)

poc_steps = [
    ("1", "PHI mask (regex)", GREEN),
    ("2", "Cache lookup (Redis)", GREEN),
    ("3", "Lane: emergency or complex (if/else)", GREEN),
    ("4", "Router: Nova Micro via bedrock.converse()", GREEN),
    ("5", "Retrieve: Bedrock KB Retrieve API (direct call)", GREEN),
    ("6", "GraphRAG: Bedrock KB Retrieve API (direct call)", GREEN),
    ("7", "Generate: bedrock.converse_stream() directly", GREEN),
    ("8", "Cache write + SSE stream to browser", GREEN),
]
for i, (num, text, color) in enumerate(poc_steps):
    y = Inches(2.25 + i * 0.52)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.36), Inches(0.36))
    solid(badge, color)
    tb = slide.shapes.add_textbox(Inches(0.7), y, Inches(0.36), Inches(0.36))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.name = "Calibri"; run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(slide, Inches(1.18), y + Inches(0.04), Inches(4.9), Inches(0.38),
             text, size=10.5, color=GRAY_DARK)

add_text(slide, Inches(0.7), Inches(6.35), Inches(5.4), Inches(0.25),
         "No icd11_lookup, no pubmed_search in PoC. PubMed = future feature.",
         size=9.5, italic=True, color=GRAY_MED)

# ── Production (right) ────────────────────────────────────────────────────────
prod_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(5.5))
bordered(prod_card, AMBER_LIGHT, AMBER, line_pt=2)

prod_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5))
solid(prod_hdr, AMBER)
add_text(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5),
         "Production Target (Bedrock Agents)", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.8), Inches(1.72), Inches(5.8), Inches(0.45),
         "Agentic loop — LLM decides which tools to call",
         size=11, bold=True, color=AMBER)

# 4 tools
tools = [
    ("kb_retrieve", "Vector + GraphRAG KB", "Semantic search on WHO, trials, protocols"),
    ("graph_retrieve", "Neptune Analytics", "Multi-hop entity traversal"),
    ("icd11_lookup", "WHO ICD-11 API", "Live disease classification + synonyms"),
    ("pubmed_search", "NCBI E-utilities", "Real-time research literature"),
]
add_text(slide, Inches(6.8), Inches(2.25), Inches(5.8), Inches(0.35),
         "Agent tools (Bedrock action groups):", size=11, bold=True, color=AMBER)
for i, (tool_name, source, desc) in enumerate(tools):
    y = Inches(2.65 + i * 0.62)
    tc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), y, Inches(5.8), Inches(0.55))
    bordered(tc, WHITE, AMBER, line_pt=1.2, radius=0.1)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), y, Inches(0.06), Inches(0.55))
    solid(acc, AMBER)
    add_text(slide, Inches(7.0), y + Inches(0.04), Inches(1.6), Inches(0.28),
             tool_name, size=10, bold=True, color=AMBER)
    add_text(slide, Inches(7.0), y + Inches(0.3), Inches(1.6), Inches(0.22),
             source, size=9, italic=True, color=GRAY_MED)
    add_text(slide, Inches(8.7), y + Inches(0.12), Inches(3.8), Inches(0.35),
             desc, size=10, color=GRAY_DARK)

add_text(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(0.45),
         "Why not in PoC?", size=11, bold=True, color=AMBER)
reasons = [
    "Bedrock Agent InvokeAgent blocked by IAM trust chain issue",
    "converse_stream() used directly instead (works, faster to build)",
    "Production: resolve IAM + enable full agentic loop",
]
bullet_list(slide, Inches(6.8), Inches(5.68), Inches(5.8), reasons,
            size=10, bullet_color=AMBER, spacing=0.28)

# Bottom banner
banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.38))
solid(banner, ORANGE)
banner.adjustments[0] = 0.2
add_text(slide, Inches(0.7), Inches(6.69), Inches(11.9), Inches(0.32),
         "PoC proves the retrieval + generation quality. Production adds Bedrock Agents for dynamic tool selection.",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

prs.save("docs/Technology_Explanation_Slides.pptx")
size = Path("docs/Technology_Explanation_Slides.pptx").stat().st_size
print(f"Saved: {size:,} bytes, {len(prs.slides)} slides")
