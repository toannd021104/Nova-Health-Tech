"""Build the 6-minute Alibaba Cloud client presentation deck for Nova Health Tech.

Output: docs/Nova_Health_Tech_AlibabaCloud_Presentation_6min.pptx

Flow (per user feedback, 6 minutes):
1. Title (10s)
2. The Challenge (40s) - PHI focus, hook to compliance
3. Singapore Compliance & Healthcare Regulations (75s) - PDPA, HCSA, Cyber Act, AI Verify, MOH guidelines
4. Continue Challenge: Speed + Freshness (30s)
5. Solution Overview: 2-Lane + SG Intl (60s)
6. Data Pipeline (40s)
7. Model Orchestration & RAG (50s)
8. Security Architecture (40s) - PHI mask, encryption, zero-trust
9. Cost & Roadmap (50s)
10. Ready to Deploy (15s)

Total: ~6 minutes
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

REPO = Path(__file__).resolve().parent.parent
BG_IMAGE = str(REPO / "mainslide.png")
OUT_PATH = str(REPO / "docs" / "Nova_Health_Tech_AlibabaCloud_Presentation_6min.pptx")

# Alibaba Cloud uses orange too, slightly different shade
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x21, 0x96, 0xF3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_slide(prs, title, subtitle, notes):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(BG_IMAGE, 0, 0, width=SLIDE_W, height=SLIDE_H)

    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    sp = overlay.fill._xPr.find(qn('a:solidFill'))
    if sp is not None:
        srgbClr = sp.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '40000')
    overlay.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_slide(prs, title, notes):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK

    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(2), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Nova Health Tech · Clinical GenAI Assistant · Alibaba Cloud Singapore"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY_MED
    run.font.italic = True

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_simple(slide, left, top, width, height, headers, rows, *, header_color=ORANGE):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = h
        run.font.name = "Calibri"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else ORANGE_LIGHT
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_top = Inches(0.03)
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = "Calibri"
            run.font.size = Pt(9.5)
            run.font.color.rgb = GRAY_DARK
    return tbl


# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ----------------------------------------------------------------------------
# SLIDE 1: TITLE (10s)
# ----------------------------------------------------------------------------
add_title_slide(
    prs,
    title="Nova Health Tech · Clinical GenAI Assistant",
    subtitle="Alibaba Cloud Singapore International · PDPA & HCSA compliant by design",
    notes=(
        "[10s opening]\n\n"
        "Đề xuất kiến trúc trợ lý AI lâm sàng cho Nova Health Tech: "
        "Alibaba Cloud Singapore International, mô hình Qwen, "
        "tuân thủ PDPA và HCSA ngay từ thiết kế."
    )
)
print("  Slide 1 (title): added")

# ----------------------------------------------------------------------------
# SLIDE 2: THE CHALLENGE - PHI focus (40s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "The Challenge",
    "[40s]\n\n"
    "Nova Health Tech là công ty Healthcare Tech cung cấp công cụ hỗ trợ quyết định lâm sàng "
    "cho các bệnh viện Đông Nam Á. Có 3 thách thức lớn nhất.\n\n"
    "Quan trọng nhất: dữ liệu thử nghiệm nội bộ và bệnh án chứa thông tin bệnh nhân nhạy cảm. "
    "Đây không chỉ là yêu cầu kỹ thuật, mà là nghĩa vụ pháp lý nghiêm ngặt theo luật y tế "
    "Singapore. PHI phải được mask trước khi đến model, audit trail đầy đủ, và data residency "
    "phải đảm bảo. Trước khi nói tiếp 2 thách thức còn lại, để tôi đi sâu vào tuân thủ Singapore."
)

# 3 pain cards, PHI prominent
pains = [
    ("⚡", "Speed", "Emergency response\np95 ≤ 2 seconds", ORANGE, False),
    ("🔒", "Patient PHI", "Healthcare data protection\nobligations are strict", RED, True),
    ("📅", "Freshness", "WHO monthly · ICD-11 daily\nLegacy PDFs everywhere", ORANGE, False),
]

card_w = Inches(3.8)
card_h = Inches(3.5)
gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) / 2

for i, (icon, title, desc, color, highlight) in enumerate(pains):
    x = start_x + i * (card_w + gap)
    y = Inches(1.5)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    if highlight:
        card.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
        card.line.color.rgb = color
        card.line.width = Pt(4)
    else:
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GRAY_LIGHT
        card.line.width = Pt(1.5)
    card.adjustments[0] = 0.06
    
    icon_tb = slide.shapes.add_textbox(x, y + Inches(0.4), card_w, Inches(1.0))
    tf = icon_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(56)
    
    title_tb = slide.shapes.add_textbox(x, y + Inches(1.6), card_w, Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(22) if not highlight else Pt(26)
    run.font.bold = True
    run.font.color.rgb = color if highlight else GRAY_DARK
    
    desc_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(2.3), card_w - Inches(0.4), Inches(1.1))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY_DARK if not highlight else color
    run.font.bold = highlight

# Arrow
arrow_text = slide.shapes.add_textbox(Inches(4.2), Inches(5.2), Inches(5.0), Inches(0.5))
tf = arrow_text.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "↑ Singapore healthcare regulation makes this non-negotiable"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.italic = True
run.font.color.rgb = RED

# Bottom: data sources
note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0))
note_box.fill.solid()
note_box.fill.fore_color.rgb = ORANGE_LIGHT
note_box.line.color.rgb = ORANGE
note_box.line.width = Pt(1)
note_box.adjustments[0] = 0.2

note_tb = slide.shapes.add_textbox(Inches(1.2), Inches(5.95), Inches(10.9), Inches(0.85))
tf = note_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Data sources: WHO guidelines · ICD-11 API · internal trial PDFs · treatment protocols · PubMed · EHR (FHIR R4)"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Workload: 500 physicians · 40 queries/day · 30/70 emergency-to-complex split · 6-year audit retention"
run2.font.name = "Calibri"
run2.font.size = Pt(11)
run2.font.italic = True
run2.font.color.rgb = GRAY_MED

print("  Slide 2 (challenge - PHI focus): added")

# ----------------------------------------------------------------------------
# SLIDE 3: SINGAPORE COMPLIANCE & HEALTHCARE REGULATIONS (75s) - KEY SLIDE
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Singapore Healthcare Compliance & Regulations",
    "[75s, KEY SLIDE]\n\n"
    "Singapore có hệ thống regulation chặt chẽ áp dụng cho hệ thống AI lâm sàng. "
    "Bốn quy định chính cần lưu ý.\n\n"
    "PDPA do PDPC quản, bảo vệ dữ liệu cá nhân, yêu cầu Data Protection Officer, "
    "vi phạm phải báo trong 72 giờ. HCSA 2020 do MOH ban hành thay thế PHMC 1980, bắt buộc "
    "license cho dịch vụ y tế bao gồm clinical decision support. Cybersecurity Act 2018 "
    "bảo vệ Critical Information Infrastructure trong y tế. IMDA AI Verify framework cho "
    "AI có trách nhiệm.\n\n"
    "Cộng thêm các quy định ngành: MOH HIE Act, Clinical Trial Regulations, HSA Therapeutic Products. "
    "Triển khai của chúng tôi: 100 phần trăm tại ap-southeast-1 Singapore International, "
    "PHI mask trước khi đến model qua DataWorks SDDP, KMS BYOK, ActionTrail audit 6 năm WORM. "
    "Alibaba Cloud có ISO 27001, ISO 27017, ISO 27018, ISO 27701, SOC 1/2/3, PCI-DSS, NIST 800-53 R5."
)

# 4 SG regulations as cards
sg_regs = [
    {
        "title": "PDPA",
        "subtitle": "Personal Data Protection Act",
        "authority": "PDPC",
        "scope": "Patient data, consent, breach in 72h, DPO required",
        "our_response": "DataWorks SDDP PHI mask, KMS BYOK, ActionTrail audit",
        "color": ORANGE_DARK,
    },
    {
        "title": "HCSA 2020",
        "subtitle": "Healthcare Services Act",
        "authority": "MOH",
        "scope": "Clinical decision support license (replaces PHMC 1980)",
        "our_response": "Decision support only, human-in-the-loop, audit ready",
        "color": ORANGE,
    },
    {
        "title": "Cyber Act",
        "subtitle": "Cybersecurity Act 2018",
        "authority": "CSA",
        "scope": "Healthcare CII protection, incident reporting",
        "our_response": "Anti-DDoS, WAF, Security Center, Content Moderation 2.0",
        "color": ORANGE_DARK,
    },
    {
        "title": "AI Verify",
        "subtitle": "IMDA Responsible AI",
        "authority": "IMDA",
        "scope": "Transparency, fairness, accountability for AI",
        "our_response": "ARMS LLM Trace, citation validator, grounding score ≥ 0.7",
        "color": ORANGE,
    },
]

cw = Inches(3.0)
ch = Inches(3.3)
gap = Inches(0.1)
total_w = cw * 4 + gap * 3
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.2)

for i, reg in enumerate(sg_regs):
    x = start_x + i * (cw + gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = reg["color"]
    card.line.width = Pt(2)
    card.adjustments[0] = 0.05
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.75))
    header.fill.solid()
    header.fill.fore_color.rgb = reg["color"]
    header.line.fill.background()
    
    h_tb = slide.shapes.add_textbox(x, y + Inches(0.05), cw, Inches(0.75))
    tf = h_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = reg["title"]
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = reg["subtitle"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(9)
    run2.font.color.rgb = WHITE
    
    body_y = y + Inches(0.85)
    
    auth_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y, cw - Inches(0.3), Inches(0.45))
    tf = auth_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Authority"
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["authority"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = GRAY_DARK
    
    scope_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y + Inches(0.7), cw - Inches(0.3), Inches(0.95))
    tf = scope_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Scope"
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["scope"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(9.5)
    run2.font.color.rgb = GRAY_DARK
    
    resp_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y + Inches(1.65), cw - Inches(0.3), Inches(0.75))
    tf = resp_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Our Implementation"
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.bold = True
    run.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["our_response"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(9.5)
    run2.font.color.rgb = GRAY_DARK

# Additional sectoral regulations
add_y = Inches(4.65)
add_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), add_y, Inches(12.3), Inches(1.0))
add_box.fill.solid()
add_box.fill.fore_color.rgb = ORANGE_LIGHT
add_box.line.color.rgb = ORANGE
add_box.line.width = Pt(1)
add_box.adjustments[0] = 0.1

add_tb = slide.shapes.add_textbox(Inches(0.7), add_y + Inches(0.1), Inches(11.9), Inches(0.85))
tf = add_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Sectoral Regulations Also Considered"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = (
    "MOH HIE Act (Health Information Exchange)  ·  Clinical Trials Regulations  ·  HSA Therapeutic Products  ·  "
    "Singapore Standard SS 564 (Healthcare IT)"
)
run2.font.name = "Calibri"
run2.font.size = Pt(11)
run2.font.color.rgb = GRAY_DARK

# Alibaba certifications + data residency
res_y = Inches(5.85)
res_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), res_y, Inches(12.3), Inches(1.05))
res_box.fill.solid()
res_box.fill.fore_color.rgb = WHITE
res_box.line.color.rgb = ORANGE_DARK
res_box.line.width = Pt(2)
res_box.adjustments[0] = 0.1

res_tb = slide.shapes.add_textbox(Inches(0.7), res_y + Inches(0.1), Inches(11.9), Inches(0.85))
tf = res_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Alibaba Cloud Certifications + Data Residency Commitment"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = (
    "ISO 27001/27017/27018/27701  ·  SOC 1/2/3  ·  PCI-DSS  ·  NIST 800-53 R5  ·  "
    "100% deployment in ap-southeast-1 Singapore International  ·  Zero cross-border at runtime  ·  6-year WORM audit"
)
run2.font.name = "Calibri"
run2.font.size = Pt(11)
run2.font.color.rgb = GRAY_DARK

print("  Slide 3 (compliance): added")

# ----------------------------------------------------------------------------
# SLIDE 4: BACK TO CHALLENGE - Speed + Freshness (30s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Continuing the Challenge",
    "[30s]\n\n"
    "Quay lại 2 thách thức còn lại. Speed: cấp cứu cần dưới 2 giây, mỗi giây có thể là sống chết. "
    "Sepsis, STEMI, anaphylaxis đều cần quyết định ngay. Freshness: WHO cập nhật hàng tháng, "
    "ICD-11 hàng ngày, hệ thống phải tự đồng bộ. Internal trial PDFs có legacy format với "
    "tagging không nhất quán, parsing khó.\n\n"
    "Cộng thêm yêu cầu tone consistency: phải có giọng văn lâm sàng nhất quán giữa các câu trả lời, "
    "không phải mỗi câu một kiểu."
)

cards = [
    {
        "icon": "⚡",
        "title": "Speed",
        "headline": "Emergency p95 ≤ 2 seconds",
        "details": [
            "Every second matters in acute care",
            "Sepsis, STEMI, anaphylaxis decisions",
            "Time-to-first-token (TTFT) is the key metric",
            "Streaming response for perceived speed",
        ],
        "color": RED,
    },
    {
        "icon": "📅",
        "title": "Freshness + Tone",
        "headline": "WHO monthly · ICD-11 daily · consistent tone",
        "details": [
            "WHO publishes guideline revisions monthly",
            "ICD-11 catalog refreshes daily",
            "Internal trial PDFs uploaded ad hoc",
            "Tone consistency: SFT on Nova-approved answers",
        ],
        "color": ORANGE_DARK,
    },
]

cw = Inches(5.8)
ch = Inches(4.5)
gap = Inches(0.4)
total_w = cw * 2 + gap
start_x = (SLIDE_W - total_w) / 2

for i, c in enumerate(cards):
    x = start_x + i * (cw + gap)
    y = Inches(1.5)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = c["color"]
    card.line.width = Pt(2.5)
    card.adjustments[0] = 0.04
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = c["color"]
    header.line.fill.background()
    
    h_tb = slide.shapes.add_textbox(x, y, cw, Inches(0.7))
    tf = h_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{c['icon']}  {c['title']}"
    run.font.name = "Calibri"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    headline_tb = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(0.95), cw - Inches(0.8), Inches(0.7))
    tf = headline_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = c["headline"]
    run.font.name = "Calibri"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = c["color"]
    
    detail_tb = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1.85), cw - Inches(1.0), Inches(2.5))
    tf = detail_tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(c["details"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        bullet_run = p.add_run()
        bullet_run.text = "▸ "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(14)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = c["color"]
        text_run = p.add_run()
        text_run.text = line
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(13)
        text_run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(8)

# Bottom: Functional + non-functional summary
nfr_y = Inches(6.3)
nfr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), nfr_y, Inches(11.3), Inches(0.7))
nfr_box.fill.solid()
nfr_box.fill.fore_color.rgb = ORANGE_LIGHT
nfr_box.line.color.rgb = ORANGE
nfr_box.line.width = Pt(1)
nfr_box.adjustments[0] = 0.2

nfr_tb = slide.shapes.add_textbox(Inches(1.2), nfr_y + Inches(0.1), Inches(10.9), Inches(0.55))
tf = nfr_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "NFR targets: 99.9% uptime · RPO ≤ 1h · RTO ≤ 4h · accuracy ≥ 95% vs gold-standard · zero PHI leakage tolerance"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.italic = True
run.font.color.rgb = ORANGE_DARK

print("  Slide 4 (challenge - speed + freshness): added")

# ----------------------------------------------------------------------------
# SLIDE 5: SOLUTION OVERVIEW (60s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Solution: Two-Lane Architecture, Singapore International",
    "[60s]\n\n"
    "Giải pháp là kiến trúc 2 luồng trên Alibaba Cloud Singapore International. "
    "Cấp cứu chạy Workflow Application deterministic với Qwen3.5-Flash, retrieve và stream "
    "trong 2 giây.\n\n"
    "Phức tạp dùng Router Agent Qwen3.5-Flash phân loại 12 phòng ban, rồi 60 phần trăm traffic "
    "đi đến Qwen3-8B student fine-tuned trên PAI-EAS, 40 phần trăm đến Qwen3.5-Plus teacher. "
    "Routing đặc biệt: ảnh thì force qua Qwen3-VL-Plus Radiology, prescribing thì invoke "
    "Clinical Pharmacy side-channel.\n\n"
    "Hai loại RAG: OpenSearch Vector Search Edition cho hybrid BM25+kNN, AnalyticDB PG GraphRAG "
    "cho multi-hop reasoning. Embedding text-embedding-v4 và tongyi-embedding-vision-plus, "
    "rerank qwen3-rerank. Tất cả đều ở Singapore International, zero cross-region hop runtime.\n\n"
    "Kiến trúc chi tiết xem AlibabaCloud_SA_proposal_technical_architecture.docx mục 3.1"
)

# Two-lane visualization
lane_y = Inches(1.3)
lane_w = Inches(6.0)
lane_h = Inches(2.7)

em_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), lane_y, lane_w, lane_h)
em_card.fill.solid()
em_card.fill.fore_color.rgb = WHITE
em_card.line.color.rgb = RED
em_card.line.width = Pt(2.5)
em_card.adjustments[0] = 0.04

em_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), lane_y, lane_w, Inches(0.55))
em_h.fill.solid()
em_h.fill.fore_color.rgb = RED
em_h.line.fill.background()

em_h_tb = slide.shapes.add_textbox(Inches(0.5), lane_y, lane_w, Inches(0.55))
tf = em_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "EMERGENCY LANE  ·  p95 ≤ 2s"
run.font.name = "Calibri"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = WHITE

em_lines = [
    "▸  Qwen3.5-Flash (1M context, streaming)",
    "▸  Workflow Application (deterministic DAG)",
    "▸  Hybrid retrieval + qwen3-rerank top-5",
    "▸  No router (direct fast lane)",
    "▸  3-layer cache: Tair + Context Cache + PTU",
]
em_tb = slide.shapes.add_textbox(Inches(0.8), lane_y + Inches(0.7), lane_w - Inches(0.5), lane_h - Inches(0.8))
tf = em_tb.text_frame
tf.word_wrap = True
for i, line in enumerate(em_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(4)

cx_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), lane_y, lane_w, lane_h)
cx_card.fill.solid()
cx_card.fill.fore_color.rgb = WHITE
cx_card.line.color.rgb = ORANGE_DARK
cx_card.line.width = Pt(2.5)
cx_card.adjustments[0] = 0.04

cx_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), lane_y, lane_w, Inches(0.55))
cx_h.fill.solid()
cx_h.fill.fore_color.rgb = ORANGE_DARK
cx_h.line.fill.background()

cx_h_tb = slide.shapes.add_textbox(Inches(6.8), lane_y, lane_w, Inches(0.55))
tf = cx_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "COMPLEX LANE  ·  p95 ≤ 6s"
run.font.name = "Calibri"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = WHITE

cx_lines = [
    "▸  Router Agent (Qwen3.5-Flash JSON mode)",
    "▸  60% Qwen3-8B student (PAI-EAS, A10)",
    "▸  40% Qwen3.5-Plus teacher (Model Studio)",
    "▸  Vision: Qwen3-VL-Plus, Radiology",
    "▸  Side-channel: Clinical Pharmacy auto",
]
cx_tb = slide.shapes.add_textbox(Inches(7.1), lane_y + Inches(0.7), lane_w - Inches(0.5), lane_h - Inches(0.8))
tf = cx_tb.text_frame
tf.word_wrap = True
for i, line in enumerate(cx_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(4)

# Bottom: shared infrastructure
arch_y = Inches(4.2)
arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), arch_y, Inches(12.3), Inches(2.4))
arch_box.fill.solid()
arch_box.fill.fore_color.rgb = ORANGE_LIGHT
arch_box.line.color.rgb = ORANGE
arch_box.line.width = Pt(1.5)
arch_box.adjustments[0] = 0.04

arch_label = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(0.1), Inches(11.9), Inches(0.4))
tf = arch_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Shared Infrastructure (100% Singapore International, ap-southeast-1)"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

# Component grid 2x3
shared_components = [
    ("Retrieval", "OpenSearch Vector Search Edition · AnalyticDB PG GraphRAG · qwen3-rerank"),
    ("Embeddings", "text-embedding-v4 (1024-dim) · tongyi-embedding-vision-plus (1152-dim)"),
    ("Compute", "Function Compute 3.0 · Function Workflow · SAE · PAI-EAS"),
    ("Security", "DataWorks SDDP · Content Moderation 2.0 · KMS BYOK · Credentials Manager"),
    ("Caching", "Tair + TairVector (L1) · Qwen Context Cache (L2) · Qwen PTU (L3 reserved)"),
    ("Network", "VPC + IPsec VPN · PrivateLink · CDN + Anti-DDoS + WAF"),
]

for i, (label, components) in enumerate(shared_components):
    row = i // 2
    col = i % 2
    x = Inches(0.7 + col * 5.95)
    y = arch_y + Inches(0.55 + row * 0.6)
    
    label_tb = slide.shapes.add_textbox(x, y, Inches(1.3), Inches(0.3))
    tf = label_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label + ":"
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    comp_tb = slide.shapes.add_textbox(x + Inches(1.3), y, Inches(4.3), Inches(0.3))
    tf = comp_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = components
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_DARK

# Architecture diagram path
path_label = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(2.0), Inches(11.9), Inches(0.35))
tf = path_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "📐  Full architecture: AlibabaCloud_SA_proposal_technical_architecture.docx §3.1"
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.italic = True
run.font.color.rgb = ORANGE_DARK

print("  Slide 5 (solution overview): added")

# ----------------------------------------------------------------------------
# SLIDE 6: DATA PIPELINE (40s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Data Pipeline: Ingest, Parse, Embed, Index",
    "[40s]\n\n"
    "Pipeline dữ liệu xử lý 3 nguồn với cadence riêng. WHO PDF cập nhật hàng tháng qua HTTP "
    "download và RSS webhook. ICD-11 API hàng ngày qua OAuth2 client. SharePoint qua Graph "
    "subscription với webhook real-time, plus reconciliation hàng tuần.\n\n"
    "Mọi file vào OSS đều qua Security Center quét malware và DataWorks SDDP quét PHI. "
    "Match thì quarantine và alert admin. DocMind parse PDF, Qwen-VL-Max xử lý complex pages "
    "với table và figure. Hierarchical chunker 1500/300 với 15 phần trăm overlap, "
    "respect section và table boundaries.\n\n"
    "Pipeline idempotent dựa trên document_id và revision hash. Chunk metadata: source, page, "
    "publication_date, review_date, evidence_grade, specialty, tenant_id."
)

# Pipeline stages
stages = [
    ("Sources", "WHO PDFs\nICD-11 API\nSharePoint\nUpload Portal", ORANGE),
    ("OSS Raw", "Object Created\nevent triggers\nFunction Workflow", ORANGE_DARK),
    ("Security", "Security Center\nmalware scan\n+ SDDP PHI scan", ORANGE),
    ("Parse", "DocMind +\nQwen-VL-Max\n(complex pages)", ORANGE_DARK),
    ("Chunk + Embed", "Hierarchical 1500/300\ntext-embedding-v4\n+ multimodal v-plus", ORANGE),
    ("Index", "OpenSearch HA\n+ AnalyticDB PG\nGraphRAG", ORANGE_DARK),
]

stage_w = Inches(1.85)
stage_h = Inches(2.0)
gap_x = Inches(0.15)
total_w = stage_w * 6 + gap_x * 5
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.4)

for i, (name, desc, color) in enumerate(stages):
    x = start_x + i * (stage_w + gap_x)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, stage_w, stage_h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.adjustments[0] = 0.08
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, stage_w, Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    
    header_tb = slide.shapes.add_textbox(x, y, stage_w, Inches(0.5))
    tf = header_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    body_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), stage_w - Inches(0.2), Inches(1.3))
    tf = body_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_DARK
    
    # Arrow
    if i < len(stages) - 1:
        from pptx.oxml.ns import qn
        line = slide.shapes.add_connector(1, x + stage_w, y + stage_h / 2, x + stage_w + gap_x, y + stage_h / 2)
        line.line.color.rgb = ORANGE
        line.line.width = Pt(2)
        ln = line.line._get_or_add_ln()
        tailEnd = ln.find(qn('a:tailEnd'))
        if tailEnd is None:
            tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('h', 'med')

# Bottom: refresh schedule + invalidation
bottom_y = Inches(3.7)
add_table_simple(slide, Inches(0.5), bottom_y, Inches(12.3), Inches(2.6),
    ["Source", "Cadence", "Trigger Mechanism"],
    [
        ["WHO ICD-11 API", "Daily 02:00 SGT", "CloudOps Scheduler cron, OAuth2 client"],
        ["WHO guideline PDFs", "Monthly day 1 02:30 SGT + RSS", "Cron + API Gateway webhook"],
        ["Internal trials/protocols", "Weekly Sun 03:00 SGT + Graph subscription", "Cron + Microsoft Graph webhook"],
        ["Manual upload", "Any time", "Upload Portal over IPsec VPN"],
        ["Full reconciliation", "Monthly day 1 04:00 SGT", "Function Workflow"],
    ])

# Cache invalidation note
note_y = Inches(6.4)
note_tb = slide.shapes.add_textbox(Inches(0.5), note_y, Inches(12.3), Inches(0.5))
tf = note_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Cache invalidation: tag-based (source:document_id) on every successful upsert. Tair flushes only matching keys."
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.italic = True
run.font.color.rgb = GRAY_MED

print("  Slide 6 (data pipeline): added")

# ----------------------------------------------------------------------------
# SLIDE 7: MODEL ORCHESTRATION & RAG (50s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Model Orchestration & RAG",
    "[50s]\n\n"
    "RAG dùng cho factual grounding, fine-tuning chỉ cho tone và latency. Không bao giờ "
    "train trên PHI: training data phải qua DataWorks SDDP de-identification trước.\n\n"
    "Hai store: OpenSearch Vector Search HA dual-zone với HNSW M=16, efConstruction=200, "
    "efSearch=80. AnalyticDB for PostgreSQL 7.0 với extension adbpg_graphrag, 4-core 32GB, "
    "3 zones. Hybrid search: BM25 + HNSW fuse qua Reciprocal Rank Fusion. Query expansion: "
    "khi detect disease mention thì icd11_expand_query inject synonyms và codes.\n\n"
    "Fine-tuning trên PAI: SFT+LoRA primary, DPO monthly trên clinician preference pairs, "
    "GRPO ad-hoc cho tool-calling regression. Chi phí 15-40 USD mỗi run, 2-4 GPU-hours trên A10. "
    "Eval harness dùng Qwen3.5-Plus judge: accuracy, citation coverage, PHI leakage, tone, "
    "emergency fit. Promote to PAI-EAS qua feature flag với 5 phần trăm canary 72 giờ."
)

# 8-step workflow
steps_y = Inches(1.3)
add_text_box_y = steps_y
title_label = slide.shapes.add_textbox(Inches(0.5), steps_y, Inches(12.3), Inches(0.4))
tf = title_label.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Query Flow (8 stages)"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

steps = [
    ("1", "PHI Mask", "DataWorks SDDP\ntokenize"),
    ("2", "Cache Lookup", "Tair semantic\nL1 hit ~30-45%"),
    ("3", "Lane Pick", "Emergency or\nComplex (if/else)"),
    ("4", "Route", "Router Agent\nor direct"),
    ("5", "Retrieve", "Vector + Graph\n+ rerank"),
    ("6", "Generate", "Qwen Flash/Plus/\nVL/Student"),
    ("7", "Validate", "Moderation +\ngrounding ≥ 0.7"),
    ("8", "Stream", "SSE to UI\n+ audit log"),
]

step_w = Inches(1.46)
step_h = Inches(1.5)
gap_x = Inches(0.1)
total_w = step_w * 8 + gap_x * 7
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.85)

for i, (num, name, desc) in enumerate(steps):
    x = start_x + i * (step_w + gap_x)
    
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h)
    box.fill.solid()
    box.fill.fore_color.rgb = ORANGE_LIGHT if i % 2 == 0 else WHITE
    box.line.color.rgb = ORANGE
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.1
    
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), y + Inches(0.1), Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    
    num_tb = slide.shapes.add_textbox(x, y + Inches(0.1), step_w, Inches(0.45))
    tf = num_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    name_tb = slide.shapes.add_textbox(x, y + Inches(0.65), step_w, Inches(0.3))
    tf = name_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(0.95), step_w - Inches(0.1), Inches(0.55))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(8.5)
    run.font.color.rgb = GRAY_DARK

# Bottom: model lineup
model_y = Inches(3.6)
add_text_box_y2 = slide.shapes.add_textbox(Inches(0.5), model_y, Inches(12.3), Inches(0.4))
tf = add_text_box_y2.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Model Lineup"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

add_table_simple(slide, Inches(0.5), Inches(4.05), Inches(12.3), Inches(2.85),
    ["Role", "Model", "Cost / Notes"],
    [
        ["Emergency fast lane", "Qwen3.5-Flash (1M ctx, streaming)", "$0.10/1M in, $0.40/1M out · TTFT ~300ms with cache"],
        ["Complex lane teacher", "Qwen3.5-Plus", "$0.40/1M in, $2.40/1M out · 1M ctx, multimodal"],
        ["Complex lane student", "Qwen3-8B on PAI-EAS (60% traffic)", "SFT+LoRA distilled · single A10 GPU · ~$0.0003/call"],
        ["Vision specialist", "Qwen3-VL-Plus", "Forced on has_image=true (Radiology)"],
        ["Router agent", "Qwen3.5-Flash JSON mode", "150-200ms p95, structured output"],
    ])

print("  Slide 7 (orchestration): added")

# ----------------------------------------------------------------------------
# SLIDE 8: SECURITY ARCHITECTURE (40s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Security Architecture",
    "[40s]\n\n"
    "Bảo mật theo defense-in-depth. Edge có CDN, Anti-DDoS, WAF. API Gateway xác thực JWT "
    "từ IDaaS EIAM 2.0. Compute là Function Compute serverless không expose public.\n\n"
    "PHI handling 2 cấp: ingest time DataWorks SDDP scan với healthcare PHI rule pack, "
    "match thì quarantine. Runtime FC chạy SDDP preflight, PHI biến thành reversible KMS tokens "
    "kiểu NAME_0, MRN_0, DOB_0. LLM chỉ thấy tokens, de-tokenize trong UI thôi. "
    "Audit log lưu tokenized form.\n\n"
    "Encryption: TLS 1.3 mọi nơi, KMS BYOK at-rest. Network zero-trust với VPC default-deny, "
    "PrivateLink cho service calls, NAT Gateway egress allow-list chỉ WHO và PubMed. "
    "Audit: ActionTrail + SLS + OSS WORM 6 năm, Object Lock không xóa được kể cả admin."
)

# 4 security pillars
pillars = [
    ("🛡️", "Edge Security", [
        "CDN + Anti-DDoS + WAF",
        "TLS 1.3 end-to-end",
        "Per-tenant IP allow-list",
        "API Gateway rate limits",
    ]),
    ("🔐", "PHI Protection", [
        "DataWorks SDDP scan",
        "Reversible KMS tokens",
        "Model never sees raw PHI",
        "Audit logs tokenized only",
    ]),
    ("🔑", "Identity & Access", [
        "IDaaS EIAM 2.0 (clinicians)",
        "Cloud SSO + RAM (staff)",
        "Per-tenant ABAC isolation",
        "Step-up MFA on admin",
    ]),
    ("📜", "Audit & Compliance", [
        "ActionTrail + SLS + OSS",
        "Object Lock WORM 6 years",
        "ARMS LLM Trace Explorer",
        "SIEM export nightly",
    ]),
]

cw = Inches(3.0)
ch = Inches(3.5)
gap = Inches(0.1)
total_w = cw * 4 + gap * 3
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.3)

for i, (icon, title, items) in enumerate(pillars):
    x = start_x + i * (cw + gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(2)
    card.adjustments[0] = 0.05
    
    icon_tb = slide.shapes.add_textbox(x, y + Inches(0.2), cw, Inches(0.7))
    tf = icon_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(36)
    
    title_tb = slide.shapes.add_textbox(x, y + Inches(1.0), cw, Inches(0.4))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    items_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.5), cw - Inches(0.4), Inches(1.9))
    tf = items_tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        bullet_run = p.add_run()
        bullet_run.text = "▸ "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(11)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = ORANGE
        text_run = p.add_run()
        text_run.text = item
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(11)
        text_run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(5)

# Bottom: 5-gate validation pipeline
gate_y = Inches(5.0)
gate_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), gate_y, Inches(12.3), Inches(1.9))
gate_box.fill.solid()
gate_box.fill.fore_color.rgb = ORANGE_LIGHT
gate_box.line.color.rgb = ORANGE_DARK
gate_box.line.width = Pt(2)
gate_box.adjustments[0] = 0.08

gate_label = slide.shapes.add_textbox(Inches(0.7), gate_y + Inches(0.1), Inches(11.9), Inches(0.4))
tf = gate_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "5-Gate Response Validation Pipeline"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

gates = [
    "1. Content Moderation 2.0 (jailbreak, self-harm, hate, medical misinformation)",
    "2. Citation validator (every [n] maps to a retrieved chunk)",
    "3. Grounding score ≥ 0.7 (or block answer)",
    "4. Last-mile PHI filter (regex + ML on MRN, NRIC, DOB, phone, email)",
    "5. Emergency disclaimer (auto-prepended on emergency-lane answers)",
]
gate_tb = slide.shapes.add_textbox(Inches(0.7), gate_y + Inches(0.55), Inches(11.9), Inches(1.3))
tf = gate_tb.text_frame
tf.word_wrap = True
for i, g in enumerate(gates):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = g
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(2)

print("  Slide 8 (security): added")

# ----------------------------------------------------------------------------
# SLIDE 9: COST & ROADMAP (50s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Cost & Implementation Roadmap",
    "[50s]\n\n"
    "Chi phí monthly cho 500 physicians, 40 queries mỗi ngày, split 30/70 emergency-complex. "
    "Total launch-day khoảng 2,280 đến 3,060 USD. Khoản lớn nhất là PAI-EAS A10 always-on cho "
    "student model, 720 đến 1,500 USD. AnalyticDB PG GraphRAG 300 USD, OpenSearch HA 180 USD. "
    "Per-call cost: emergency với cache hit chỉ 0.0008 USD, complex với cache 0.0026, "
    "vision không cache 0.004 USD.\n\n"
    "Lộ trình triển khai 6 đến 10 tuần. Foundation 1-2 tuần: provisioning, IDaaS, OpenSearch HA, "
    "AnalyticDB PG. Data pipeline và RAG 1-4 tuần. Model fine-tuning 3-5 tuần. "
    "Orchestration multi-agent 3-6 tuần. Clinical embedding và security 5-7 tuần. "
    "Performance compliance hardening 6-9 tuần. Clinical pilot và cut-over 9-10 tuần."
)

# Cost table
cost_y = Inches(1.3)
add_text_box_y3 = slide.shapes.add_textbox(Inches(0.5), cost_y, Inches(6.0), Inches(0.4))
tf = add_text_box_y3.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Monthly Cost (USD, list price early 2026)"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

add_table_simple(slide, Inches(0.5), Inches(1.75), Inches(6.0), Inches(4.3),
    ["Item", "Monthly"],
    [
        ["Qwen3.5-Flash (emergency, 180k post-cache)", "$47"],
        ["Qwen3.5-Plus (40% complex, teacher)", "$440"],
        ["PAI-EAS A10 always-on (student, 60% complex)", "$720-1,500"],
        ["SFT+LoRA training amortized quarterly", "$5-15"],
        ["text-embedding-v4 + multimodal + rerank", "$135"],
        ["Content Moderation 2.0", "$50"],
        ["OpenSearch Vector Search HA", "$180"],
        ["AnalyticDB PG GraphRAG (4-core 32GB)", "$300"],
        ["DataWorks SDDP", "$120"],
        ["FC + API GW + CDN + WAF + OSS + Tair", "$220"],
        ["IPsec VPN Gateway (data plane)", "$110-150"],
        ["TOTAL launch-day (per tenant)", "$2,280-3,060"],
    ])

# Roadmap on right
rm_y = Inches(1.3)
rm_label = slide.shapes.add_textbox(Inches(7.0), rm_y, Inches(5.8), Inches(0.4))
tf = rm_label.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Implementation Roadmap (6-10 weeks)"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

phases = [
    ("Wk 1-2", "Foundation", "VPC, KMS, IDaaS, OpenSearch HA, AnalyticDB PG"),
    ("Wk 1-4", "Data Pipeline", "WHO + ICD-11 ingest, DocMind parse, embed"),
    ("Wk 3-5", "Model + FT", "Qwen3-8B SFT+LoRA, eval harness, PAI-EAS"),
    ("Wk 3-6", "Orchestration", "12 Agent apps + Workflow, 4 tools, vision"),
    ("Wk 5-7", "Integration", "EHR FHIR, IDaaS federation, IPsec VPN"),
    ("Wk 6-9", "Hardening", "Red team 200+ prompts, PTU sizing, DR test"),
    ("Wk 9-10", "Cut-over", "Clinical pilot, sign-off, production live"),
]

phase_y = Inches(1.75)
phase_h = Inches(0.55)
for i, (week, name, desc) in enumerate(phases):
    y = phase_y + i * (phase_h + Inches(0.05))
    
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), y, Inches(1.0), phase_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    
    badge_tb = slide.shapes.add_textbox(Inches(7.0), y, Inches(1.0), phase_h)
    tf = badge_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = week
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    name_tb = slide.shapes.add_textbox(Inches(8.1), y, Inches(1.6), phase_h)
    tf = name_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(Inches(9.7), y, Inches(3.2), phase_h)
    tf = desc_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_DARK

# Per-call cost summary
pc_y = Inches(6.2)
pc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), pc_y, Inches(12.3), Inches(0.7))
pc_box.fill.solid()
pc_box.fill.fore_color.rgb = ORANGE_LIGHT
pc_box.line.color.rgb = ORANGE
pc_box.line.width = Pt(1)
pc_box.adjustments[0] = 0.2

pc_tb = slide.shapes.add_textbox(Inches(0.7), pc_y + Inches(0.1), Inches(11.9), Inches(0.55))
tf = pc_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Per-call cost: Emergency w/ cache hit $0.0008  ·  Student amortized $0.0003  ·  Complex w/ cache $0.0026  ·  Vision $0.004"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

print("  Slide 9 (cost + roadmap): added")

# ----------------------------------------------------------------------------
# SLIDE 10: READY TO DEPLOY (15s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Ready to Deploy",
    "[15s closing]\n\n"
    "Tóm lại: kiến trúc Singapore International, tuân thủ PDPA HCSA Cybersecurity Act, "
    "chi phí từ 2,280 USD per tenant. Lộ trình 6-10 tuần. Sẵn sàng PoC khi Nova approve. "
    "Xin cảm ơn."
)

# Key takeaways
take_y = Inches(1.3)
takeaways = [
    ("✓", "Singapore-Native", "100% in ap-southeast-1 Singapore International, zero cross-region runtime"),
    ("✓", "Compliance Built-In", "PDPA, HCSA 2020, Cybersecurity Act, IMDA AI Verify, ISO 27001/27017/27018/27701"),
    ("✓", "Proven Architecture", "Two-lane pattern with managed services: Model Studio, GraphRAG, OpenSearch, Tair"),
    ("✓", "Cost-Efficient", "$2,280-3,060 per tenant per month; 60% traffic on $0.0003/call student model"),
    ("✓", "Fast to Deploy", "6-10 week build window; one product, no phases, all features day-one"),
]

t_y = Inches(1.5)
for i, (mark, title, desc) in enumerate(takeaways):
    y = t_y + i * Inches(0.7)
    
    # Check mark
    mark_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.6), Inches(0.5))
    tf = mark_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = mark
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = GREEN
    
    title_tb = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(Inches(4.7), y + Inches(0.1), Inches(8.1), Inches(0.5))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK

# Q&A callout
qa_y = Inches(5.5)
qa_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), qa_y, Inches(8.3), Inches(1.4))
qa_box.fill.solid()
qa_box.fill.fore_color.rgb = ORANGE
qa_box.line.fill.background()
qa_box.adjustments[0] = 0.1

qa_tb = slide.shapes.add_textbox(Inches(2.7), qa_y + Inches(0.2), Inches(7.9), Inches(1.0))
tf = qa_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Questions & Discussion"
run.font.name = "Calibri"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Reference: AlibabaCloud_SA_proposal_technical_architecture.docx (35 pages)"
run2.font.name = "Calibri"
run2.font.size = Pt(12)
run2.font.italic = True
run2.font.color.rgb = WHITE

print("  Slide 10 (ready to deploy): added")

prs.save(OUT_PATH)
print(f"\nSaved: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
