"""Fix slide 3 (Multi-Agent) to accurately reflect 12 departments in PoC, not 40."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK  = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED     = RGBColor(0x64, 0x64, 0x64)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x27, 0xAE, 0x60)
GREEN_LIGHT  = RGBColor(0xE8, 0xF5, 0xE9)
AMBER        = RGBColor(0xFF, 0x8F, 0x00)
AMBER_LIGHT  = RGBColor(0xFF, 0xF8, 0xE1)
PURPLE       = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)
RED          = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation("docs/Technology_Explanation_Slides.pptx")
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def bordered(shape, fill_color, line_color, line_pt=1.5, radius=0.06):
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color; shape.line.width = Pt(line_pt)
    shape.adjustments[0] = radius

def add_text(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False):
    color = color or GRAY_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_list(slide, left, top, width, items, size=10.5, color=None,
                bullet_color=None, spacing=0.32):
    color = color or GRAY_DARK
    bullet_color = bullet_color or ORANGE
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(left, top + Inches(i * spacing), width, Inches(spacing))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        b = p.add_run(); b.text = "\u25b8 "
        b.font.name = "Calibri"; b.font.size = Pt(size)
        b.font.bold = True; b.font.color.rgb = bullet_color
        t = p.add_run(); t.text = item
        t.font.name = "Calibri"; t.font.size = Pt(size)
        t.font.color.rgb = color

# Background + bar + footer
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
solid(bg, WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
solid(bar, ORANGE)
ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28))
tf = ft.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run = p.add_run()
run.text = "Nova Health Tech  \u00b7  Clinical GenAI Assistant  \u00b7  AWS Singapore"
run.font.name = "Calibri"; run.font.size = Pt(9)
run.font.italic = True; run.font.color.rgb = GRAY_MED

# Title
add_text(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.65),
         "Multi-Agent Architecture: 12 Specialty Agents (PoC)",
         size=28, bold=True, color=ORANGE_DARK)
ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.92), Inches(4.2), Inches(0.04))
solid(ul, ORANGE)

# ── Left: How it works ────────────────────────────────────────────────────────
left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5))
bordered(left_card, ORANGE_LIGHT, ORANGE, line_pt=1.5)

left_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5))
solid(left_hdr, ORANGE_DARK)
add_text(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5),
         "How It Works", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.4), Inches(0.55),
         "Each department is a specialized agent: its own system prompt, "
         "clinical scope, and model. The router picks the right one.",
         size=11, color=GRAY_DARK)

# Flow diagram
flow_items = [
    ("Doctor's question", ORANGE, False),
    ("PHI mask + cache check", ORANGE, False),
    ("Emergency toggle ON?", RED, True),
    ("YES: Emergency agent (Haiku 4.5)", RGBColor(0xC0,0x39,0x2B), False),
    ("NO: Router (Nova Micro, JSON mode)", ORANGE_DARK, False),
    ("Routes to 1 of 12 specialty agents", ORANGE, False),
    ("Specialty agent (Sonnet 4.5) generates answer", ORANGE_DARK, False),
    ("Citations validated + streamed to browser", GREEN, False),
]
for i, (text, color, is_decision) in enumerate(flow_items):
    y = Inches(2.38 + i * 0.46)
    if is_decision:
        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(0.7), y, Inches(5.2), Inches(0.38))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(5.2), Inches(0.38))
        shape.adjustments[0] = 0.15
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.7), y, Inches(5.2), Inches(0.38))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(9.5)
    run.font.bold = True; run.font.color.rgb = WHITE

# ── Right: 12 departments grid ────────────────────────────────────────────────
right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(5.5))
bordered(right_card, WHITE, ORANGE, line_pt=1.5)

right_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5))
solid(right_hdr, ORANGE)
add_text(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5),
         "12 Deployed Department Agents (PoC)", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 12 departments in 3 columns x 4 rows
depts = [
    ("Emergency", "Haiku 4.5", RGBColor(0xC0,0x39,0x2B)),
    ("Cardiology", "Sonnet 4.5", ORANGE_DARK),
    ("Pulmonology", "Sonnet 4.5", ORANGE_DARK),
    ("Gastroenterology", "Sonnet 4.5", ORANGE_DARK),
    ("Nephrology", "Sonnet 4.5", ORANGE_DARK),
    ("Endocrinology", "Sonnet 4.5", ORANGE_DARK),
    ("Neurology", "Sonnet 4.5", ORANGE_DARK),
    ("Infectious Disease", "Sonnet 4.5", ORANGE_DARK),
    ("Oncology", "Sonnet 4.5", ORANGE_DARK),
    ("Obstetrics & Gyn", "Sonnet 4.5", ORANGE_DARK),
    ("Pediatrics", "Sonnet 4.5", ORANGE_DARK),
    ("Radiology", "Sonnet 4.5 + Vision", ORANGE),
]
col_w = Inches(1.9)
row_h = Inches(0.72)
gap_x = Inches(0.1)
gap_y = Inches(0.1)
start_x = Inches(6.7)
start_y = Inches(1.72)

for i, (name, model, color) in enumerate(depts):
    col = i % 3
    row = i // 3
    x = start_x + col * (col_w + gap_x)
    y = start_y + row * (row_h + gap_y)
    
    dc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
    bordered(dc, ORANGE_LIGHT if i % 2 == 0 else WHITE, color, line_pt=1.5, radius=0.1)
    
    add_text(slide, x + Inches(0.08), y + Inches(0.08), col_w - Inches(0.16), Inches(0.35),
             name, size=10, bold=True, color=color)
    add_text(slide, x + Inches(0.08), y + Inches(0.42), col_w - Inches(0.16), Inches(0.25),
             model, size=8.5, italic=True, color=GRAY_MED)

# Note about production
note_y = Inches(6.55)
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), note_y, Inches(12.3), Inches(0.42))
solid(note, ORANGE)
note.adjustments[0] = 0.2
add_text(slide, Inches(0.7), note_y + Inches(0.06), Inches(11.9), Inches(0.32),
         "PoC: 12 departments, system prompt per agent, bedrock.converse() directly.  "
         "Production: 40 sub-specialties, Bedrock Agents service with tool calling.",
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

prs.save("docs/Technology_Explanation_Slides.pptx")
size = Path("docs/Technology_Explanation_Slides.pptx").stat().st_size
print(f"Saved: {size:,} bytes, {len(prs.slides)} slides")

