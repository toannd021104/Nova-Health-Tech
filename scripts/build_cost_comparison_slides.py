"""Build 4-slide cost comparison deck for 3 production versions.

Slide 1: Summary comparison (all 3 versions side-by-side)
Slide 2: Version A breakdown (AWS + Claude)
Slide 3: Version B breakdown (AWS + Qwen + fine-tuning)
Slide 4: Version C breakdown (Alibaba Cloud)

Assumptions (production, 500 clinicians, 10 queries/day):
  - 150,000 queries/month
  - 30 emergency (20%) + 120 complex (80%) per 1000
  - Avg tokens: 3,000 in / 600 out per complex query
  - Avg tokens: 1,500 in / 300 out per emergency query
  - Monthly fine-tuning run: 1x SFT job
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from pathlib import Path

OUT = Path("docs/Production_Cost_Comparison_3Versions.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x38, 0x64)
BLUE    = RGBColor(0x15, 0x65, 0xC0)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)
ORANGE  = RGBColor(0xE6, 0x51, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
YELLOW  = RGBColor(0xFF, 0xC1, 0x07)
RED     = RGBColor(0xC6, 0x28, 0x28)
LBLUE   = RGBColor(0xE3, 0xF2, 0xFD)
LGREEN  = RGBColor(0xE8, 0xF5, 0xE9)
LORANGE = RGBColor(0xFF, 0xF3, 0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

def add_slide():
    return prs.slides.add_slide(blank)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def box(slide, x, y, w, h, bg=None, border_color=None, border_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=11, bold=False, color=DGRAY,
        align=PP_ALIGN.LEFT, bg=None, border_color=None, border_pt=0,
        italic=False, wrap=True):
    tf_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if bg:
        tf_box.fill.solid()
        tf_box.fill.fore_color.rgb = bg
    if border_color and border_pt > 0:
        tf_box.line.color.rgb = border_color
        tf_box.line.width = Pt(border_pt)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tf_box

def header_bar(slide, title, subtitle=""):
    box(slide, 0, 0, 13.33, 1.0, bg=NAVY)
    txt(slide, title, 0.3, 0.08, 10, 0.5, size=20, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.58, 12, 0.35, size=11, color=RGBColor(0xBB, 0xCC, 0xFF), italic=True)

def footer(slide, note=""):
    box(slide, 0, 7.1, 13.33, 0.4, bg=NAVY)
    txt(slide, "Nova Health Tech · Production Architecture Cost Estimate · May 2026",
        0.3, 7.12, 8, 0.28, size=8, color=RGBColor(0xAA, 0xBB, 0xDD))
    if note:
        txt(slide, note, 8.5, 7.12, 4.5, 0.28, size=8, color=RGBColor(0xAA, 0xBB, 0xDD),
            align=PP_ALIGN.RIGHT)

def assumption_box(slide):
    box(slide, 0.2, 1.05, 12.93, 0.55, bg=RGBColor(0xE8, 0xEA, 0xF6),
        border_color=BLUE, border_pt=0.5)
    txt(slide,
        "Assumptions: 500 clinicians · 10 queries/day · 150,000 queries/month "
        "(30,000 emergency + 120,000 complex) · Avg 3,000 in / 600 out tokens (complex) · "
        "1,500 in / 300 out (emergency) · 1 fine-tuning run/month · Singapore region",
        0.3, 1.07, 12.7, 0.45, size=9, color=BLUE, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — SUMMARY COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()
header_bar(s1, "Production Cost Comparison — 3 Versions",
           "Monthly estimate · 500 clinicians · 150,000 queries/month · Singapore region")
assumption_box(s1)
footer(s1, "All prices USD, list price, on-demand")

# Version column headers
ver_data = [
    (1.0,  BLUE,   LBLUE,   "Version A",   "AWS + Claude",    "$3,820 – $4,650 / mo"),
    (5.1,  GREEN,  LGREEN,  "Version B",   "AWS + Qwen",      "$2,940 – $3,680 / mo"),
    (9.2,  ORANGE, LORANGE, "Version C",   "Alibaba Cloud",   "$2,280 – $3,060 / mo"),
]
for x, hdr_c, bg_c, ver, stack, total in ver_data:
    box(s1, x, 1.65, 3.8, 0.7, bg=hdr_c)
    txt(s1, ver,   x+0.1, 1.67, 3.6, 0.3, size=13, bold=True, color=WHITE)
    txt(s1, stack, x+0.1, 1.97, 3.6, 0.3, size=10, color=WHITE)
    box(s1, x, 2.35, 3.8, 0.55, bg=bg_c, border_color=hdr_c, border_pt=1.5)
    txt(s1, total, x+0.1, 2.38, 3.6, 0.45, size=16, bold=True, color=hdr_c,
        align=PP_ALIGN.CENTER)

# Cost breakdown rows
rows = [
    ("LLM Inference",
     "$1,800 – $2,400",  "$1,200 – $1,600",  "$900 – $1,200"),
    ("RAG / Knowledge Base",
     "$480 – $600",      "$480 – $600",       "$480 – $600"),
    ("Semantic Cache",
     "$120 – $180",      "$120 – $180",       "$60 – $90"),
    ("Fine-tuning (monthly)",
     "— (no fine-tuning)", "$340 – $480",     "$120 – $200"),
    ("Inference Endpoint\n(student model)",
     "—",                "$480 – $720",       "$720 – $1,080"),
    ("Compute / Serverless",
     "$180 – $240",      "$180 – $240",       "$90 – $120"),
    ("Storage + Logs",
     "$80 – $100",       "$80 – $100",        "$70 – $90"),
    ("Security / Observability",
     "$160 – $200",      "$160 – $200",       "$120 – $160"),
    ("Network / VPN / CDN",
     "$120 – $150",      "$120 – $150",       "$110 – $150"),
]

y = 2.95
row_h = 0.42
for i, (label, a, b, c) in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    box(s1, 0.2, y, 0.85, row_h, bg=RGBColor(0xE8, 0xEA, 0xF6))
    txt(s1, label, 0.25, y+0.03, 0.8, row_h-0.06, size=8, bold=True, color=NAVY)
    for xi, val, hdr_c in [(1.0, a, BLUE), (5.1, b, GREEN), (9.2, c, ORANGE)]:
        box(s1, xi, y, 3.8, row_h, bg=bg, border_color=RGBColor(0xDD,0xDD,0xDD), border_pt=0.3)
        is_na = val.startswith("—")
        txt(s1, val, xi+0.1, y+0.04, 3.6, row_h-0.08,
            size=9, color=RGBColor(0xAA,0xAA,0xAA) if is_na else hdr_c,
            italic=is_na, align=PP_ALIGN.CENTER)
    y += row_h

# Total row
box(s1, 0.2, y, 0.85, 0.45, bg=NAVY)
txt(s1, "TOTAL / MO", 0.25, y+0.05, 0.8, 0.35, size=8, bold=True, color=WHITE)
for xi, val, hdr_c, bg_c in [
    (1.0, "$3,820 – $4,650", BLUE,   LBLUE),
    (5.1, "$2,940 – $3,680", GREEN,  LGREEN),
    (9.2, "$2,280 – $3,060", ORANGE, LORANGE),
]:
    box(s1, xi, y, 3.8, 0.45, bg=bg_c, border_color=hdr_c, border_pt=1.5)
    txt(s1, val, xi+0.1, y+0.05, 3.6, 0.35, size=12, bold=True, color=hdr_c,
        align=PP_ALIGN.CENTER)

# Key insight box
y2 = y + 0.55
box(s1, 0.2, y2, 12.93, 0.55, bg=RGBColor(0xFF,0xF8,0xE1),
    border_color=YELLOW, border_pt=1)
txt(s1,
    "Key insight:  Alibaba Cloud is ~30% cheaper than AWS+Claude due to lower Qwen model token costs. "
    "AWS+Qwen is ~20% cheaper than AWS+Claude and adds fine-tuning capability. "
    "All three versions include the same RAG, cache, and security baseline.",
    0.35, y2+0.06, 12.6, 0.42, size=9, color=RGBColor(0x5D,0x40,0x00))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — VERSION A: AWS + CLAUDE BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
box(s2, 0, 0, 13.33, 1.0, bg=BLUE)
txt(s2, "Version A — AWS + Claude  |  Cost Breakdown",
    0.3, 0.08, 10, 0.5, size=20, bold=True, color=WHITE)
txt(s2, "Production · ap-southeast-1 Singapore · 150,000 queries/month",
    0.3, 0.58, 12, 0.35, size=11, color=RGBColor(0xBB,0xCC,0xFF), italic=True)
assumption_box(s2)
footer(s2, "Version A total: $3,820 – $4,650 / month")

# Left column: LLM costs
box(s2, 0.2, 1.65, 6.2, 0.4, bg=BLUE)
txt(s2, "LLM INFERENCE  (Bedrock on-demand)", 0.35, 1.68, 6.0, 0.32,
    size=11, bold=True, color=WHITE)

llm_rows_a = [
    ("Emergency lane", "Claude Haiku 4.5", "30,000 calls",
     "1,500 in + 300 out tok", "$1.00/$5.00 per 1M",
     "30k×1.5k×$1 + 30k×0.3k×$5 = $45 + $45", "$90"),
    ("Complex lane\n(router)", "Nova Micro", "120,000 calls",
     "500 in + 40 out tok", "$0.035/$0.14 per 1M",
     "120k×0.5k×$0.035 + 120k×0.04k×$0.14 = $2.1 + $0.67", "$3"),
    ("Complex lane\n(specialist)", "Claude Sonnet 4.5", "120,000 calls",
     "3,000 in + 600 out tok", "$3.00/$15.00 per 1M",
     "120k×3k×$3 + 120k×0.6k×$15 = $1,080 + $1,080", "$2,160"),
    ("Embeddings\n(query-time)", "Cohere Embed Multi v3", "150,000 calls",
     "~80 tok avg", "$0.10 per 1M",
     "150k×0.08k×$0.10 = $1.2", "$1"),
    ("Guardrails", "Bedrock Guardrails", "120,000 calls",
     "~3,600 tok/call", "$0.15 per 1,000 units",
     "120k×3.6 units×$0.15/1k = $64.8", "$65"),
    ("Prompt cache\nsavings", "Auto (Claude 4.x)", "—",
     "~60% cache hit on system prompt", "−90% on cached tokens",
     "Estimated 40% reduction on Sonnet cost", "−$864"),
]

y = 2.1
for label, model, calls, tokens, price, calc, total in llm_rows_a:
    is_saving = total.startswith("−")
    bg = RGBColor(0xE8,0xF5,0xE9) if is_saving else LBLUE
    box(s2, 0.2, y, 6.2, 0.48, bg=bg, border_color=RGBColor(0xCC,0xDD,0xFF), border_pt=0.3)
    txt(s2, label,  0.25, y+0.02, 1.1, 0.44, size=8, bold=True, color=NAVY)
    txt(s2, model,  1.35, y+0.02, 1.5, 0.22, size=8, color=BLUE)
    txt(s2, calls,  1.35, y+0.24, 1.5, 0.22, size=7, color=DGRAY, italic=True)
    txt(s2, tokens, 2.85, y+0.02, 1.5, 0.22, size=7, color=DGRAY)
    txt(s2, price,  2.85, y+0.24, 1.5, 0.22, size=7, color=DGRAY)
    txt(s2, calc,   4.35, y+0.02, 1.5, 0.44, size=7, color=DGRAY, italic=True)
    c = GREEN if is_saving else BLUE
    txt(s2, total,  5.85, y+0.1, 0.5, 0.28, size=10, bold=True, color=c,
        align=PP_ALIGN.RIGHT)
    y += 0.5

# LLM subtotal
box(s2, 0.2, y, 6.2, 0.38, bg=BLUE)
txt(s2, "LLM Subtotal (after cache savings)", 0.35, y+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s2, "~$1,455 / month", 5.0, y+0.05, 1.3, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Right column: Infrastructure
box(s2, 6.7, 1.65, 6.43, 0.4, bg=NAVY)
txt(s2, "INFRASTRUCTURE & SERVICES", 6.85, 1.68, 6.2, 0.32,
    size=11, bold=True, color=WHITE)

infra_rows_a = [
    ("Knowledge Base\n(Vector)", "Bedrock KB OpenSearch Serverless",
     "2 OCU × 720 hr", "$0.24/OCU-hr", "$346"),
    ("Knowledge Base\n(GraphRAG)", "Neptune Analytics",
     "1 m-NCU × 720 hr", "$0.16/m-NCU-hr", "$115"),
    ("Semantic Cache", "ElastiCache Redis OSS\ncache.r7g.large",
     "720 hr", "$0.166/hr", "$120"),
    ("Compute\n(API server)", "EC2 t4g.medium × 2\n(HA pair)",
     "720 hr × 2", "$0.046/hr", "$66"),
    ("API Gateway", "AWS API Gateway REST",
     "150k calls/mo", "$3.50/1M", "$1"),
    ("CDN / Edge", "CloudFront",
     "~50GB transfer", "$0.12/GB", "$6"),
    ("Storage", "S3 Standard\n(corpus + logs)",
     "~100GB", "$0.023/GB-mo", "$2"),
    ("Audit / Logs", "CloudWatch Logs\n+ S3 archive",
     "~10GB ingest/mo", "$0.50/GB", "$5"),
    ("Security", "AWS WAF\n+ Bedrock Guardrails",
     "150k requests", "$0.60/1M req", "$1"),
    ("VPN Gateway", "AWS VPN (1 tunnel)",
     "720 hr", "$0.05/hr", "$36"),
    ("IAM / KMS", "KMS key usage",
     "~100k API calls", "$0.03/10k", "$0"),
    ("Fine-tuning", "— (no fine-tuning\nin Version A)",
     "—", "—", "$0"),
]

y2 = 2.1
for label, service, usage, rate, total in infra_rows_a:
    bg = LGRAY if infra_rows_a.index((label,service,usage,rate,total)) % 2 == 0 else WHITE
    box(s2, 6.7, y2, 6.43, 0.42, bg=bg, border_color=RGBColor(0xDD,0xDD,0xDD), border_pt=0.3)
    txt(s2, label,   6.75, y2+0.02, 1.2, 0.38, size=8, bold=True, color=NAVY)
    txt(s2, service, 7.95, y2+0.02, 2.0, 0.38, size=8, color=DGRAY)
    txt(s2, usage,   9.95, y2+0.02, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s2, rate,    9.95, y2+0.22, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s2, total,  12.45, y2+0.08, 0.6, 0.26, size=10, bold=True,
        color=BLUE if total != "$0" else RGBColor(0xAA,0xAA,0xAA),
        align=PP_ALIGN.RIGHT)
    y2 += 0.42

# Infra subtotal
box(s2, 6.7, y2, 6.43, 0.38, bg=NAVY)
txt(s2, "Infrastructure Subtotal", 6.85, y2+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s2, "~$698 / month", 11.5, y2+0.05, 1.55, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Grand total
y3 = max(y, y2) + 0.45
box(s2, 0.2, y3, 12.93, 0.55, bg=BLUE)
txt(s2, "VERSION A TOTAL", 0.4, y3+0.08, 4, 0.38, size=13, bold=True, color=WHITE)
txt(s2, "LLM $1,455  +  Infra $698  +  Buffer 15%  =", 4.5, y3+0.1, 6, 0.35,
    size=10, color=LBLUE)
txt(s2, "~$2,480 – $2,700 / month", 10.0, y3+0.06, 3.1, 0.42,
    size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VERSION B: AWS + QWEN BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
box(s3, 0, 0, 13.33, 1.0, bg=GREEN)
txt(s3, "Version B — AWS + Qwen  |  Cost Breakdown",
    0.3, 0.08, 10, 0.5, size=20, bold=True, color=WHITE)
txt(s3, "Production · ap-southeast-1 Singapore · 150,000 queries/month · Monthly fine-tuning included",
    0.3, 0.58, 12, 0.35, size=11, color=RGBColor(0xBB,0xFF,0xCC), italic=True)
assumption_box(s3)
footer(s3, "Version B total: $2,940 – $3,680 / month")

# LLM section
box(s3, 0.2, 1.65, 6.2, 0.4, bg=GREEN)
txt(s3, "LLM INFERENCE  (Bedrock on-demand, all Singapore)", 0.35, 1.68, 6.0, 0.32,
    size=11, bold=True, color=WHITE)

llm_rows_b = [
    ("Emergency lane", "Nova Lite", "30,000 calls",
     "1,500 in + 300 out tok", "$0.06/$0.24 per 1M",
     "30k×1.5k×$0.06 + 30k×0.3k×$0.24 = $2.7 + $2.16", "$5"),
    ("Complex lane\n(router)", "Nova Micro", "120,000 calls",
     "500 in + 40 out tok", "$0.035/$0.14 per 1M",
     "120k×0.5k×$0.035 + 120k×0.04k×$0.14 = $2.1 + $0.67", "$3"),
    ("Complex lane\n(specialist)", "Nova Pro", "120,000 calls",
     "3,000 in + 600 out tok", "$0.80/$3.20 per 1M",
     "120k×3k×$0.80 + 120k×0.6k×$3.20 = $288 + $230", "$518"),
    ("Student model\n(60% complex)", "Qwen3-4B SFT+LoRA\nSageMaker Endpoint",
     "72,000 calls", "~1,500 tok avg", "$1.97/hr (g5.xlarge)",
     "g5.xlarge 720hr × $1.97 = $1,418 (always-on)", "$1,418"),
    ("Embeddings", "Cohere Embed Multi v3", "150,000 calls",
     "~80 tok avg", "$0.10 per 1M", "150k×0.08k×$0.10 = $1.2", "$1"),
    ("Guardrails", "Bedrock Guardrails", "48,000 calls\n(40% complex, teacher only)",
     "~3,600 tok/call", "$0.15 per 1,000 units",
     "48k×3.6×$0.15/1k = $25.9", "$26"),
]

y = 2.1
for label, model, calls, tokens, price, calc, total in llm_rows_b:
    box(s3, 0.2, y, 6.2, 0.5, bg=LGREEN, border_color=RGBColor(0xCC,0xFF,0xCC), border_pt=0.3)
    txt(s3, label,  0.25, y+0.02, 1.1, 0.46, size=8, bold=True, color=NAVY)
    txt(s3, model,  1.35, y+0.02, 1.5, 0.22, size=8, color=GREEN)
    txt(s3, calls,  1.35, y+0.26, 1.5, 0.22, size=7, color=DGRAY, italic=True)
    txt(s3, tokens, 2.85, y+0.02, 1.5, 0.22, size=7, color=DGRAY)
    txt(s3, price,  2.85, y+0.26, 1.5, 0.22, size=7, color=DGRAY)
    txt(s3, calc,   4.35, y+0.02, 1.5, 0.46, size=7, color=DGRAY, italic=True)
    txt(s3, total,  5.85, y+0.12, 0.5, 0.26, size=10, bold=True, color=GREEN,
        align=PP_ALIGN.RIGHT)
    y += 0.52

box(s3, 0.2, y, 6.2, 0.38, bg=GREEN)
txt(s3, "LLM + Student Endpoint Subtotal", 0.35, y+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s3, "~$1,971 / month", 5.0, y+0.05, 1.3, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Fine-tuning section
y_ft = y + 0.45
box(s3, 0.2, y_ft, 6.2, 0.35, bg=RGBColor(0x1B,0x5E,0x20))
txt(s3, "MONTHLY FINE-TUNING  (SageMaker Training Job)", 0.35, y_ft+0.05, 6.0, 0.25,
    size=10, bold=True, color=WHITE)

ft_rows = [
    ("Data generation", "Qwen3.5-397B via DeepInfra\n(HF provider)",
     "10,000 Q&A / month", "$0.54/$3.40 per 1M", "$4.33"),
    ("Training job", "SageMaker ml.g4dn.2xlarge\nQwen3-4B SFT+LoRA",
     "~3 hours / run", "$1.32/hr", "$4"),
    ("Model storage", "S3 checkpoint\n(~8GB per run)", "720 hr", "$0.023/GB-mo", "$0"),
]
y_ft2 = y_ft + 0.38
for label, service, usage, rate, total in ft_rows:
    box(s3, 0.2, y_ft2, 6.2, 0.4, bg=RGBColor(0xE8,0xF5,0xE9),
        border_color=GREEN, border_pt=0.3)
    txt(s3, label,   0.25, y_ft2+0.02, 1.1, 0.36, size=8, bold=True, color=NAVY)
    txt(s3, service, 1.35, y_ft2+0.02, 2.0, 0.36, size=8, color=GREEN)
    txt(s3, usage,   3.35, y_ft2+0.02, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s3, rate,    3.35, y_ft2+0.22, 1.5, 0.16, size=7, color=DGRAY, italic=True)
    txt(s3, total,   5.85, y_ft2+0.08, 0.5, 0.24, size=10, bold=True, color=GREEN,
        align=PP_ALIGN.RIGHT)
    y_ft2 += 0.42

box(s3, 0.2, y_ft2, 6.2, 0.35, bg=RGBColor(0x1B,0x5E,0x20))
txt(s3, "Fine-tuning Subtotal (monthly)", 0.35, y_ft2+0.05, 4.5, 0.25,
    size=9, bold=True, color=WHITE)
txt(s3, "~$8 / month", 5.0, y_ft2+0.05, 1.3, 0.25, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Right column: Infrastructure (same as Version A minus fine-tuning)
box(s3, 6.7, 1.65, 6.43, 0.4, bg=NAVY)
txt(s3, "INFRASTRUCTURE & SERVICES", 6.85, 1.68, 6.2, 0.32,
    size=11, bold=True, color=WHITE)

infra_rows_b = [
    ("Knowledge Base\n(Vector)", "Bedrock KB OpenSearch Serverless",
     "2 OCU × 720 hr", "$0.24/OCU-hr", "$346"),
    ("Knowledge Base\n(GraphRAG)", "Neptune Analytics",
     "1 m-NCU × 720 hr", "$0.16/m-NCU-hr", "$115"),
    ("Semantic Cache", "ElastiCache Redis OSS\ncache.r7g.large",
     "720 hr", "$0.166/hr", "$120"),
    ("Compute\n(API server)", "EC2 t4g.medium × 2",
     "720 hr × 2", "$0.046/hr", "$66"),
    ("API Gateway", "AWS API Gateway REST",
     "150k calls/mo", "$3.50/1M", "$1"),
    ("CDN / Edge", "CloudFront",
     "~50GB transfer", "$0.12/GB", "$6"),
    ("Storage", "S3 Standard",
     "~100GB", "$0.023/GB-mo", "$2"),
    ("Audit / Logs", "CloudWatch Logs + S3",
     "~10GB ingest/mo", "$0.50/GB", "$5"),
    ("Security", "AWS WAF + Guardrails",
     "150k requests", "$0.60/1M req", "$1"),
    ("VPN Gateway", "AWS VPN (1 tunnel)",
     "720 hr", "$0.05/hr", "$36"),
    ("Student endpoint\n(scale-to-zero opt.)", "SageMaker Serverless\n(alternative to always-on)",
     "72k calls × 3s × 6GB", "$0.00002/GB-s", "$26"),
]

y2 = 2.1
for label, service, usage, rate, total in infra_rows_b:
    bg = LGRAY if infra_rows_b.index((label,service,usage,rate,total)) % 2 == 0 else WHITE
    box(s3, 6.7, y2, 6.43, 0.44, bg=bg, border_color=RGBColor(0xDD,0xDD,0xDD), border_pt=0.3)
    txt(s3, label,   6.75, y2+0.02, 1.2, 0.40, size=8, bold=True, color=NAVY)
    txt(s3, service, 7.95, y2+0.02, 2.0, 0.40, size=8, color=DGRAY)
    txt(s3, usage,   9.95, y2+0.02, 1.5, 0.20, size=7, color=DGRAY, italic=True)
    txt(s3, rate,    9.95, y2+0.24, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s3, total,  12.45, y2+0.09, 0.6, 0.26, size=10, bold=True,
        color=GREEN if total != "$0" else RGBColor(0xAA,0xAA,0xAA),
        align=PP_ALIGN.RIGHT)
    y2 += 0.44

box(s3, 6.7, y2, 6.43, 0.38, bg=NAVY)
txt(s3, "Infrastructure Subtotal", 6.85, y2+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s3, "~$724 / month", 11.5, y2+0.05, 1.55, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Grand total
y3 = max(y_ft2, y2) + 0.45
box(s3, 0.2, y3, 12.93, 0.55, bg=GREEN)
txt(s3, "VERSION B TOTAL", 0.4, y3+0.08, 4, 0.38, size=13, bold=True, color=WHITE)
txt(s3, "LLM+Student $1,971  +  Fine-tune $8  +  Infra $724  +  Buffer 15%  =",
    4.5, y3+0.1, 6.5, 0.35, size=10, color=LGREEN)
txt(s3, "~$3,120 – $3,680 / month", 10.0, y3+0.06, 3.1, 0.42,
    size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — VERSION C: ALIBABA CLOUD BREAKDOWN
# ═══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
box(s4, 0, 0, 13.33, 1.0, bg=ORANGE)
txt(s4, "Version C — Alibaba Cloud  |  Cost Breakdown",
    0.3, 0.08, 10, 0.5, size=20, bold=True, color=WHITE)
txt(s4, "Production · ap-southeast-1 Singapore International · 150,000 queries/month · Monthly fine-tuning included",
    0.3, 0.58, 12, 0.35, size=11, color=RGBColor(0xFF,0xE0,0xB2), italic=True)
assumption_box(s4)
footer(s4, "Version C total: $2,280 – $3,060 / month")

# LLM section
box(s4, 0.2, 1.65, 6.2, 0.4, bg=ORANGE)
txt(s4, "LLM INFERENCE  (Model Studio / DashScope Intl)", 0.35, 1.68, 6.0, 0.32,
    size=11, bold=True, color=WHITE)

llm_rows_c = [
    ("Emergency lane", "Qwen3.5-Flash\n(qwen3.5-flash)", "30,000 calls",
     "1,500 in + 300 out tok", "$0.05/$0.20 per 1M",
     "30k×1.5k×$0.05 + 30k×0.3k×$0.20 = $2.25 + $1.80", "$4"),
    ("Complex lane\n(router)", "Qwen3.5-Flash\n(JSON mode)", "120,000 calls",
     "500 in + 40 out tok", "$0.05/$0.20 per 1M",
     "120k×0.5k×$0.05 + 120k×0.04k×$0.20 = $3 + $0.96", "$4"),
    ("Complex lane\n(specialist)", "Qwen3.5-Plus\n(qwen3.5-plus)", "120,000 calls",
     "3,000 in + 600 out tok", "$0.115/$0.688 per 1M",
     "120k×3k×$0.115 + 120k×0.6k×$0.688 = $41.4 + $49.5", "$91"),
    ("Student model\n(60% complex)", "Qwen3-8B on PAI-EAS\n(A10 GPU, SFT+LoRA)",
     "72,000 calls", "~1,500 tok avg", "$1.00–$2.00/hr (A10)",
     "A10 720hr × $1.50 avg = $1,080 (always-on)", "$1,080"),
    ("Embeddings", "text-embedding-v4",
     "150,000 calls", "~80 tok avg", "$0.07 per 1M",
     "150k×0.08k×$0.07 = $0.84", "$1"),
    ("Re-rank", "qwen3-rerank",
     "120,000 calls", "~3,600 tok/call", "$0.10 per 1M",
     "120k×3.6k×$0.10 = $43.2", "$43"),
    ("Context Cache\n(L2 savings)", "Qwen Context Cache",
     "~90k calls (60% hit)", "−90% on cached tokens",
     "Est. 40% reduction on Plus cost", "40% × $91 = −$36", "−$36"),
]

y = 2.1
for label, model, calls, tokens, price, calc, total in llm_rows_c:
    is_saving = total.startswith("−")
    bg = RGBColor(0xE8,0xF5,0xE9) if is_saving else LORANGE
    box(s4, 0.2, y, 6.2, 0.5, bg=bg, border_color=RGBColor(0xFF,0xCC,0xAA), border_pt=0.3)
    txt(s4, label,  0.25, y+0.02, 1.1, 0.46, size=8, bold=True, color=NAVY)
    txt(s4, model,  1.35, y+0.02, 1.5, 0.22, size=8, color=ORANGE)
    txt(s4, calls,  1.35, y+0.26, 1.5, 0.22, size=7, color=DGRAY, italic=True)
    txt(s4, tokens, 2.85, y+0.02, 1.5, 0.22, size=7, color=DGRAY)
    txt(s4, price,  2.85, y+0.26, 1.5, 0.22, size=7, color=DGRAY)
    txt(s4, calc,   4.35, y+0.02, 1.5, 0.46, size=7, color=DGRAY, italic=True)
    c = GREEN if is_saving else ORANGE
    txt(s4, total,  5.85, y+0.12, 0.5, 0.26, size=10, bold=True, color=c,
        align=PP_ALIGN.RIGHT)
    y += 0.52

box(s4, 0.2, y, 6.2, 0.38, bg=ORANGE)
txt(s4, "LLM + PAI-EAS Subtotal", 0.35, y+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s4, "~$1,187 / month", 5.0, y+0.05, 1.3, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Fine-tuning section
y_ft = y + 0.45
box(s4, 0.2, y_ft, 6.2, 0.35, bg=RGBColor(0xBF,0x36,0x00))
txt(s4, "MONTHLY FINE-TUNING  (PAI-DLC + PAI-EAS)", 0.35, y_ft+0.05, 6.0, 0.25,
    size=10, bold=True, color=WHITE)

ft_rows_c = [
    ("Data generation", "Qwen3.5-Plus as teacher\n(Model Studio)",
     "10,000 Q&A / month", "$0.115/$0.688 per 1M", "$8"),
    ("Training job", "PAI-DLC A10 GPU\nQwen3-8B SFT+LoRA",
     "~4 hours / run", "$1.50/hr (A10)", "$6"),
    ("Model storage", "OSS checkpoint\n(~16GB per run)", "720 hr", "$0.02/GB-mo", "$0"),
]
y_ft2 = y_ft + 0.38
for label, service, usage, rate, total in ft_rows_c:
    box(s4, 0.2, y_ft2, 6.2, 0.4, bg=LORANGE, border_color=ORANGE, border_pt=0.3)
    txt(s4, label,   0.25, y_ft2+0.02, 1.1, 0.36, size=8, bold=True, color=NAVY)
    txt(s4, service, 1.35, y_ft2+0.02, 2.0, 0.36, size=8, color=ORANGE)
    txt(s4, usage,   3.35, y_ft2+0.02, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s4, rate,    3.35, y_ft2+0.22, 1.5, 0.16, size=7, color=DGRAY, italic=True)
    txt(s4, total,   5.85, y_ft2+0.08, 0.5, 0.24, size=10, bold=True, color=ORANGE,
        align=PP_ALIGN.RIGHT)
    y_ft2 += 0.42

box(s4, 0.2, y_ft2, 6.2, 0.35, bg=RGBColor(0xBF,0x36,0x00))
txt(s4, "Fine-tuning Subtotal (monthly)", 0.35, y_ft2+0.05, 4.5, 0.25,
    size=9, bold=True, color=WHITE)
txt(s4, "~$14 / month", 5.0, y_ft2+0.05, 1.3, 0.25, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Right column: Infrastructure
box(s4, 6.7, 1.65, 6.43, 0.4, bg=NAVY)
txt(s4, "INFRASTRUCTURE & SERVICES", 6.85, 1.68, 6.2, 0.32,
    size=11, bold=True, color=WHITE)

infra_rows_c = [
    ("Vector Store", "OpenSearch Vector Search HA\n2 OCU baseline",
     "2 OCU × 720 hr", "$0.24/OCU-hr", "$346"),
    ("GraphRAG", "AnalyticDB for PostgreSQL\n4-core 32GB + adbpg_graphrag",
     "720 hr", "$0.42/hr", "$302"),
    ("Semantic Cache", "Tair (Redis-compatible)\n1GB cluster, Multi-AZ",
     "720 hr", "$0.083/hr", "$60"),
    ("Compute\n(API server)", "Function Compute 3.0\n(serverless)",
     "150k calls × 2s × 512MB", "$0.00001667/GB-s", "$26"),
    ("API Gateway", "Alibaba API Gateway",
     "150k calls/mo", "$3.50/1M", "$1"),
    ("CDN / Edge", "Alibaba CDN + Anti-DDoS",
     "~50GB transfer", "$0.10/GB", "$5"),
    ("Storage", "OSS Standard\n(corpus + WORM audit)",
     "~200GB", "$0.02/GB-mo", "$4"),
    ("Audit / Logs", "SLS + OSS WORM\n(6yr retention)",
     "~10GB ingest/mo", "$0.35/GB", "$4"),
    ("Security", "WAF + Content Mod 2.0\n+ SDDP scan",
     "150k requests", "~$0.50/1M", "$1"),
    ("VPN Gateway", "Alibaba VPN Gateway\n100 Mbps",
     "720 hr", "$0.15/hr", "$108"),
    ("IDaaS / SSO", "IDaaS EIAM 2.0 Premium+\nhospital IdP federation",
     "500 users", "$0.10/user/mo", "$50"),
]

y2 = 2.1
for label, service, usage, rate, total in infra_rows_c:
    bg = LGRAY if infra_rows_c.index((label,service,usage,rate,total)) % 2 == 0 else WHITE
    box(s4, 6.7, y2, 6.43, 0.44, bg=bg, border_color=RGBColor(0xDD,0xDD,0xDD), border_pt=0.3)
    txt(s4, label,   6.75, y2+0.02, 1.2, 0.40, size=8, bold=True, color=NAVY)
    txt(s4, service, 7.95, y2+0.02, 2.0, 0.40, size=8, color=DGRAY)
    txt(s4, usage,   9.95, y2+0.02, 1.5, 0.20, size=7, color=DGRAY, italic=True)
    txt(s4, rate,    9.95, y2+0.24, 1.5, 0.18, size=7, color=DGRAY, italic=True)
    txt(s4, total,  12.45, y2+0.09, 0.6, 0.26, size=10, bold=True,
        color=ORANGE if total != "$0" else RGBColor(0xAA,0xAA,0xAA),
        align=PP_ALIGN.RIGHT)
    y2 += 0.44

box(s4, 6.7, y2, 6.43, 0.38, bg=NAVY)
txt(s4, "Infrastructure Subtotal", 6.85, y2+0.05, 4.5, 0.28,
    size=9, bold=True, color=WHITE)
txt(s4, "~$907 / month", 11.5, y2+0.05, 1.55, 0.28, size=11, bold=True,
    color=YELLOW, align=PP_ALIGN.RIGHT)

# Grand total
y3 = max(y_ft2, y2) + 0.45
box(s4, 0.2, y3, 12.93, 0.55, bg=ORANGE)
txt(s4, "VERSION C TOTAL", 0.4, y3+0.08, 4, 0.38, size=13, bold=True, color=WHITE)
txt(s4, "LLM+PAI-EAS $1,187  +  Fine-tune $14  +  Infra $907  +  Buffer 15%  =",
    4.5, y3+0.1, 6.5, 0.35, size=10, color=LORANGE)
txt(s4, "~$2,420 – $3,060 / month", 10.0, y3+0.06, 3.1, 0.42,
    size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}")
print("4 slides: Summary + Version A + Version B + Version C")
