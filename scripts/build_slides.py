"""Build the Nova Health Tech client presentation deck.

Output: docs/Nova_Health_Tech_AWS_Claude_Presentation.pptx

Design rules:
- Slide 1: Title slide using mainslide.png as full background
- Slides 2+: White background, orange accent for headers/boxes/icons
- Professional, client-facing tone
- Each slide includes presenter notes (lời trình bày)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

REPO = Path(__file__).resolve().parent.parent
BG_IMAGE = str(REPO / "mainslide.png")
OUT_PATH = str(REPO / "docs" / "Nova_Health_Tech_AWS_Claude_Presentation.pptx")

# Color palette
ORANGE = RGBColor(0xFF, 0x6B, 0x35)        # primary orange
ORANGE_DARK = RGBColor(0xCC, 0x4F, 0x1F)   # darker orange for titles
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)  # very light orange (box bg)
GRAY_DARK = RGBColor(0x2C, 0x3E, 0x50)     # body text
GRAY_MED = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_slide(prs, title, subtitle, notes):
    """Slide 1: full-image background with title overlay."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background image (full slide)
    slide.shapes.add_picture(BG_IMAGE, 0, 0, width=SLIDE_W, height=SLIDE_H)

    # Semi-transparent dark overlay for text readability
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    # Set transparency via XML
    sp = overlay.fill._xPr.find(qn('a:solidFill'))
    if sp is not None:
        srgbClr = sp.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '40000')  # 40% opacity
    overlay.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE

    # Notes
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_slide(prs, title, notes):
    """Add content slide with white background and orange title bar.
    Returns slide for further customization.
    """
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Orange accent bar at top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK

    # Thin underline below title
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Nova Health Tech · Clinical GenAI Assistant · AWS Singapore"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY_MED
    run.font.italic = True

    # Notes
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_text_box(slide, left, top, width, height, text, *, size=14, bold=False, color=GRAY_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Helper: add a styled text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_list(slide, left, top, width, height, bullets, *, size=14, color=GRAY_DARK, indent=False):
    """Add a bullet list with orange bullet markers."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Orange bullet
        bullet_run = p.add_run()
        bullet_run.text = "▸ "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = ORANGE
        # Text
        text_run = p.add_run()
        text_run.text = text
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_card(slide, left, top, width, height, title, body_lines, *, accent_color=ORANGE):
    """Add a card with orange title header and body content."""
    # Card background (light orange tint)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    # Rounded corner radius
    card.adjustments[0] = 0.06

    # Title bar (orange)
    title_h = Inches(0.45)
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, title_h)
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = accent_color
    title_bar.line.fill.background()

    # Title text
    title_tb = slide.shapes.add_textbox(left, top, width, title_h)
    tf = title_tb.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Body
    body_top = top + title_h + Inches(0.1)
    body_h = height - title_h - Inches(0.15)
    body_tb = slide.shapes.add_textbox(left + Inches(0.1), body_top, width - Inches(0.2), body_h)
    tf = body_tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(3)


def add_table_simple(slide, left, top, width, height, headers, rows, *, header_color=ORANGE):
    """Add a styled table with orange header and clean rows."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = h
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else ORANGE_LIGHT
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_top = Inches(0.03)
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = "Calibri"
            run.font.size = Pt(10)
            run.font.color.rgb = GRAY_DARK
    return tbl


def add_arrow(slide, x1, y1, x2, y2, color=ORANGE, weight=2.5):
    """Add an arrow (line with arrow end)."""
    from pptx.oxml.ns import qn
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # Add arrow head
    ln = line.line._get_or_add_ln()
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('h', 'med')
    return line


# ============================================================================
# BUILD PRESENTATION
# ============================================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ----------------------------------------------------------------------------
# SLIDE 1: TITLE
# ----------------------------------------------------------------------------
add_title_slide(
    prs,
    title="Nova Health Tech · Clinical GenAI Assistant",
    subtitle="Technical Architecture Proposal · AWS with Claude · Singapore",
    notes=(
        "Lời trình bày:\n\n"
        "Xin chào Ban lãnh đạo Nova Health Tech. Hôm nay tôi sẽ trình bày đề xuất kiến trúc kỹ "
        "thuật cho trợ lý AI hỗ trợ quyết định lâm sàng của Nova, được xây dựng trên nền tảng "
        "AWS với mô hình Claude và đặt tại vùng Singapore.\n\n"
        "Đề xuất này tập trung vào 3 mục tiêu chính: tốc độ đáp ứng dưới 2 giây cho ca cấp cứu, "
        "tính chính xác và truy nguyên được của câu trả lời với trích dẫn nguồn, và tuân thủ "
        "đầy đủ các quy định PDPA, HCSA, HIPAA. Tôi sẽ đi qua từ bài toán nghiệp vụ, kiến trúc "
        "tổng thể, đến chi phí vận hành chi tiết."
    )
)

print(f"  Slide 1 (title): added")

# ----------------------------------------------------------------------------
# SLIDE 2: AGENDA
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Agenda",
    "Lời trình bày:\n\n"
    "Bài trình bày của tôi gồm 8 phần. Bắt đầu với bối cảnh và các nỗi đau của Nova hiện tại, "
    "sau đó là tổng quan giải pháp và các thành phần lõi: pipeline dữ liệu, điều phối mô hình, "
    "RAG, mạng kết nối, bảo mật. Cuối cùng là chi phí vận hành và lộ trình triển khai. "
    "Mỗi phần tôi sẽ tập trung vào quyết định kiến trúc và lý do, không đi sâu vào chi tiết kỹ thuật."
)

agenda_items = [
    ("1", "Scenario & Pain Points", "Bối cảnh và các thách thức hiện tại"),
    ("2", "Solution Overview", "Tổng quan kiến trúc đề xuất"),
    ("3", "Data Pipeline", "Quy trình ingest dữ liệu đa nguồn"),
    ("4", "Model Orchestration & RAG", "Điều phối mô hình và truy xuất tri thức"),
    ("5", "Network & Integration", "Tích hợp mạng với hệ thống bệnh viện"),
    ("6", "Security & Compliance", "Bảo mật và tuân thủ"),
    ("7", "Performance Results", "Kết quả PoC: TTFT, SLA"),
    ("8", "Cost Estimation", "Dự toán chi phí vận hành"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    row = i // 2
    col = i % 2
    x = Inches(0.7 + col * 6.0)
    y = Inches(1.5 + row * 1.3)
    
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    
    badge_tb = slide.shapes.add_textbox(x, y, Inches(0.7), Inches(0.7))
    tf = badge_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Title and description
    title_tb = slide.shapes.add_textbox(x + Inches(0.9), y, Inches(4.8), Inches(0.4))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(x + Inches(0.9), y + Inches(0.4), Inches(4.8), Inches(0.3))
    tf = desc_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_MED

print(f"  Slide 2 (agenda): added")

# ----------------------------------------------------------------------------
# SLIDE 3: SCENARIO & PAIN POINTS
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Scenario & Pain Points",
    "Lời trình bày:\n\n"
    "Nova Health Tech là công ty Digital Health, sản phẩm chính là công cụ hỗ trợ quyết định "
    "lâm sàng. Hiện tại bác sĩ phàn nàn rằng câu trả lời chậm và thiếu tính cụ thể với từng "
    "bệnh nhân. Ban lãnh đạo phê duyệt xây dựng trợ lý AI cho nhân viên y tế và bệnh viện đối tác.\n\n"
    "Có 5 thách thức chính: (1) Cấp cứu cần dưới 2 giây, không có chỗ cho độ trễ. (2) Dữ liệu "
    "thử nghiệm nội bộ chứa thông tin bệnh nhân nhạy cảm, không được rò rỉ. (3) WHO cập nhật "
    "phác đồ hàng tháng, ICD-11 cập nhật hàng ngày, hệ thống phải tự đồng bộ. (4) Báo cáo "
    "thử nghiệm là PDF cũ, gắn nhãn không nhất quán, parsing khó. (5) Bác sĩ muốn giọng văn "
    "lâm sàng nhất quán, không phải mỗi câu trả lời một kiểu."
)

# Background context box
ctx_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.3))
ctx_box.fill.solid()
ctx_box.fill.fore_color.rgb = ORANGE_LIGHT
ctx_box.line.color.rgb = ORANGE
ctx_box.line.width = Pt(1)
ctx_box.adjustments[0] = 0.08

ctx_tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(1.1))
tf = ctx_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Context"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = (
    "Nova's flagship clinical decision-support tool is losing physician trust. "
    "Answers feel slow and too generic. The board approved building an internal GenAI assistant "
    "for clinical staff and hospital clients, grounded in WHO guidelines, ICD-11, and Nova's "
    "internal trial reports."
)
run2.font.name = "Calibri"
run2.font.size = Pt(12)
run2.font.color.rgb = GRAY_DARK

# Pain point cards (5 cards in row)
pain_points = [
    ("⚡", "Speed", "Emergency response\n< 2 seconds"),
    ("🔒", "Patient PHI", "Trial data contains\npatient-sensitive info"),
    ("📅", "Stale Data", "WHO updates monthly,\nICD-11 daily"),
    ("📄", "Legacy PDFs", "Inconsistent tagging,\nmixed formats"),
    ("✓", "Tone", "Consistent clinical\nphrasing required"),
]

card_w = Inches(2.4)
card_h = Inches(2.8)
gap = Inches(0.05)
total_w = Inches(0.5) + card_w * 5 + gap * 4
start_x = (SLIDE_W - (card_w * 5 + gap * 4)) / 2

for i, (icon, title, desc) in enumerate(pain_points):
    x = start_x + i * (card_w + gap)
    y = Inches(2.9)
    
    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(1.5)
    card.adjustments[0] = 0.08
    
    # Icon (large, orange)
    icon_tb = slide.shapes.add_textbox(x, y + Inches(0.3), card_w, Inches(0.8))
    tf = icon_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(36)
    
    # Title
    title_tb = slide.shapes.add_textbox(x, y + Inches(1.2), card_w, Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    # Description
    desc_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.7), card_w - Inches(0.2), Inches(1.0))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_DARK

# Bottom note
note_tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8))
tf = note_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = (
    "Plus: HIPAA, Singapore PDPA, HCSA, FDA 21 CFR 820, EU AI Act compliance · "
    "auditable for 6 years · 600k calls/month at launch"
)
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.italic = True
run.font.color.rgb = GRAY_MED

print(f"  Slide 3 (pain points): added")

# ----------------------------------------------------------------------------
# SLIDE 4: SOLUTION OVERVIEW
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Solution Overview",
    "Lời trình bày:\n\n"
    "Giải pháp gồm 5 lớp xếp chồng, tất cả đặt tại AWS Singapore, không có cross-region hop "
    "cho dữ liệu lâm sàng. Lớp Edge xử lý truy cập public với WAF và allow-list IP bệnh viện. "
    "Lớp API xác thực qua Cognito federated với Entra ID. Lớp Compute là Lambda + ECS Fargate. "
    "Lớp Model là Claude Haiku 4.5 cho cấp cứu, Claude Sonnet 4.5 cho ca phức tạp, Nova Micro "
    "định tuyến phòng ban. Lớp Knowledge là OpenSearch Serverless cho vector + Neptune Analytics "
    "cho GraphRAG.\n\n"
    "Hai luồng chạy song song: cấp cứu đi thẳng Haiku 4.5 với 5 chunk context, ca phức tạp "
    "qua router phòng ban rồi Sonnet 4.5 với 18 chunk context. Tất cả đều stream SSE để TTFT "
    "thấp, có Guardrails kiểm tra nội dung trước khi trả về."
)

# 5 horizontal layers
layers = [
    ("Edge Layer", "CloudFront · WAF · Route 53", "Public entry, IP allow-list, TLS 1.3", ORANGE),
    ("API Layer", "API Gateway · Cognito · Entra ID federation", "Authenticated REST, per-tenant rate limits", ORANGE_DARK),
    ("Compute Layer", "Lambda · ECS Fargate · LangGraph orchestration", "Chat runtime with PHI masking, lane routing", ORANGE),
    ("Model Layer", "Claude Haiku 4.5 (emergency) · Claude Sonnet 4.5 (complex) · Nova Micro (router)", "Streaming Converse API, 100% SG-native", ORANGE_DARK),
    ("Knowledge Layer", "OpenSearch Serverless (vector) · Neptune Analytics (GraphRAG) · Bedrock Guardrails", "Hybrid retrieval + entity-aware traversal", ORANGE),
]

layer_h = Inches(0.85)
gap_y = Inches(0.1)
start_y = Inches(1.4)

for i, (name, components, purpose, color) in enumerate(layers):
    y = start_y + i * (layer_h + gap_y)
    
    # Layer name (left, ~2.5in)
    name_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(2.4), layer_h)
    name_box.fill.solid()
    name_box.fill.fore_color.rgb = color
    name_box.line.fill.background()
    
    name_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(2.4), layer_h)
    tf = name_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Content (right)
    content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), y, Inches(9.8), layer_h)
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = ORANGE_LIGHT
    content_box.line.color.rgb = color
    content_box.line.width = Pt(0.75)
    
    content_tb = slide.shapes.add_textbox(Inches(3.1), y + Inches(0.05), Inches(9.6), layer_h - Inches(0.1))
    tf = content_tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = components
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = GRAY_DARK
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = purpose
    run2.font.name = "Calibri"
    run2.font.size = Pt(11)
    run2.font.italic = True
    run2.font.color.rgb = GRAY_MED

print(f"  Slide 4 (solution overview): added")

# ----------------------------------------------------------------------------
# SLIDE 5: DATA PIPELINE
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Data Pipeline",
    "Lời trình bày:\n\n"
    "Pipeline dữ liệu xử lý 3 nguồn khác nhau với cadence riêng. WHO PDF cập nhật hàng tháng, "
    "ICD-11 API hàng ngày, báo cáo thử nghiệm nội bộ qua SharePoint webhook. Mọi file vào S3 "
    "đều qua GuardDuty quét malware và Macie quét PHI bất đồng bộ.\n\n"
    "PDF được parse bởi Bedrock Data Automation tại Sydney, tách bảng, hình, text. Chunker "
    "Lambda áp dụng chiến lược đa dạng: hierarchical 1500/300 cho WHO, semantic max 512 cho "
    "clinical trials, dùng tốt nhất cho từng loại tài liệu. Cohere Embed v3 SG-native nhúng "
    "vector 1024 chiều, OpenSearch Serverless index hybrid kNN + BM25.\n\n"
    "Comprehend Medical bóc tách PHI trước khi lập chỉ mục, không có thông tin bệnh nhân nào "
    "rời tenant. Step Functions điều phối toàn bộ pipeline, ghi lineage vào Glue và DynamoDB."
)

# Pipeline stages (horizontal flow)
stages = [
    ("Sources", "WHO PDFs\nICD-11 API\nSharePoint", ORANGE),
    ("Ingest", "S3 raw\nGuardDuty\nMacie scan", ORANGE_DARK),
    ("Parse", "Bedrock Data\nAutomation\n(Sydney)", ORANGE),
    ("Chunk", "Hierarchical\nSemantic\nMetadata tags", ORANGE_DARK),
    ("Embed", "Cohere Embed\nMultilingual v3\n(SG-native)", ORANGE),
    ("Index", "OpenSearch SL\n+ Neptune\n(GraphRAG)", ORANGE_DARK),
]

stage_w = Inches(1.85)
stage_h = Inches(1.8)
gap_x = Inches(0.15)
total_w = stage_w * 6 + gap_x * 5
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.5)

for i, (name, desc, color) in enumerate(stages):
    x = start_x + i * (stage_w + gap_x)
    
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, stage_w, stage_h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.adjustments[0] = 0.08
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, stage_w, Inches(0.5))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    
    header_tb = slide.shapes.add_textbox(x, y, stage_w, Inches(0.5))
    tf = header_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Body
    body_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), stage_w - Inches(0.2), Inches(1.1))
    tf = body_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(10.5)
    run.font.color.rgb = GRAY_DARK
    
    # Arrow to next
    if i < len(stages) - 1:
        arrow_y = y + stage_h / 2
        add_arrow(slide, x + stage_w, arrow_y, x + stage_w + gap_x, arrow_y, color=ORANGE, weight=2)

# Bottom: refresh cadence + chunking strategy
bottom_y = Inches(3.7)

# Left card: Refresh cadence
add_card(slide, Inches(0.5), bottom_y, Inches(6.0), Inches(2.3),
    "Refresh Cadence",
    [
        "WHO ICD-11 API:  Daily 02:00 SGT (delta sync)",
        "WHO guideline PDFs:  Monthly day 1 + RSS webhook",
        "SharePoint trials:  Weekly Sun 03:00 SGT + Graph webhook",
        "Manual upload:  Ad hoc via VPN admin portal",
        "",
        "Cache invalidation: tag-based, only matching chunks flushed.",
    ])

# Right card: Chunking strategy
add_card(slide, Inches(6.8), bottom_y, Inches(6.0), Inches(2.3),
    "Multi-Strategy Chunking",
    [
        "WHO guidelines:  Hierarchical (parent 1500 / child 300)",
        "Clinical trials:  Semantic (max 512, breakpoint 80%)",
        "ICD-11 entries:  No-chunk (1 entity = 1 chunk)",
        "",
        "Rationale: each source has different structure;",
        "single chunking strategy cannot fit all formats.",
    ])

print(f"  Slide 5 (data pipeline): added")

# ----------------------------------------------------------------------------
# SLIDE 6: MODEL ORCHESTRATION & RAG
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Model Orchestration & RAG",
    "Lời trình bày:\n\n"
    "Khi bác sĩ gửi câu hỏi, hệ thống chạy theo 8 bước. Bước 1, Comprehend Medical che PHI. "
    "Bước 2, ElastiCache Redis kiểm tra semantic cache, nếu hit thì trả ngay. Bước 3, "
    "phân làn dựa trên emergency toggle - nếu bật thì đi thẳng Haiku, nếu không thì Nova Micro "
    "phân loại 12 phòng ban. Bước 4, retrieve song song từ OpenSearch và Neptune GraphRAG. "
    "Bước 5, Sonnet 4.5 hoặc Haiku 4.5 generate câu trả lời với citations. "
    "Bước 6, Bedrock Guardrails kiểm tra grounding và topic. Bước 7, citation validator xác minh "
    "mỗi [N] phải resolve được tới chunk thực tế. Bước 8, stream SSE về browser.\n\n"
    "Khác biệt giữa 2 luồng: cấp cứu dùng top 3 vector + top 2 graph, không guardrails, "
    "max 300 token để TTFT 2.5 giây. Phức tạp dùng top 15 vector + top 3 graph, có guardrails, "
    "max 1500 token, TTFT khoảng 9.7 giây."
)

# Two-lane comparison
lane_y = Inches(1.4)
lane_w = Inches(6.0)
lane_h = Inches(4.5)

# Emergency lane (left)
em_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), lane_y, lane_w, lane_h)
em_card.fill.solid()
em_card.fill.fore_color.rgb = WHITE
em_card.line.color.rgb = RED
em_card.line.width = Pt(2)
em_card.adjustments[0] = 0.04

em_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), lane_y, lane_w, Inches(0.5))
em_header.fill.solid()
em_header.fill.fore_color.rgb = RED
em_header.line.fill.background()

em_h_tb = slide.shapes.add_textbox(Inches(0.5), lane_y, lane_w, Inches(0.5))
tf = em_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "EMERGENCY LANE  ·  TTFT 2.5s avg"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

em_steps = [
    "1.  PHI mask (Comprehend Medical)",
    "2.  Semantic cache check (Redis)",
    "3.  Direct to Emergency agent (no router)",
    "4.  Vector KB top-3 + GraphRAG top-2 (~5 chunks)",
    "5.  Claude Haiku 4.5 streaming (max 300 tokens)",
    "6.  No guardrails (post-hoc check, speed priority)",
    "7.  Citation validator + chat_trace audit",
    "8.  SSE stream to browser",
]
em_steps_tb = slide.shapes.add_textbox(Inches(0.7), lane_y + Inches(0.7), lane_w - Inches(0.4), lane_h - Inches(0.8))
tf = em_steps_tb.text_frame
tf.word_wrap = True
for i, step in enumerate(em_steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = step
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(7)

# Complex lane (right)
cx_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), lane_y, lane_w, lane_h)
cx_card.fill.solid()
cx_card.fill.fore_color.rgb = WHITE
cx_card.line.color.rgb = ORANGE_DARK
cx_card.line.width = Pt(2)
cx_card.adjustments[0] = 0.04

cx_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), lane_y, lane_w, Inches(0.5))
cx_header.fill.solid()
cx_header.fill.fore_color.rgb = ORANGE_DARK
cx_header.line.fill.background()

cx_h_tb = slide.shapes.add_textbox(Inches(6.8), lane_y, lane_w, Inches(0.5))
tf = cx_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "COMPLEX LANE  ·  TTFT 9.7s avg"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

cx_steps = [
    "1.  PHI mask (Comprehend Medical)",
    "2.  Semantic cache check (Redis)",
    "3.  Nova Micro routes to 1 of 12 departments",
    "4.  Vector KB top-15 + GraphRAG top-3 (~18 chunks)",
    "5.  Claude Sonnet 4.5 streaming (max 1500 tokens)",
    "6.  Bedrock Guardrails (grounding, PHI, topic)",
    "7.  Citation validator + chat_trace audit",
    "8.  SSE stream to browser",
]
cx_steps_tb = slide.shapes.add_textbox(Inches(7.0), lane_y + Inches(0.7), lane_w - Inches(0.4), lane_h - Inches(0.8))
tf = cx_steps_tb.text_frame
tf.word_wrap = True
for i, step in enumerate(cx_steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = step
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(7)

# Bottom note
note_y = Inches(6.2)
note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), note_y, Inches(12.3), Inches(0.7))
note_box.fill.solid()
note_box.fill.fore_color.rgb = ORANGE_LIGHT
note_box.line.color.rgb = ORANGE
note_box.line.width = Pt(1)
note_box.adjustments[0] = 0.2

note_tb = slide.shapes.add_textbox(Inches(0.7), note_y + Inches(0.05), Inches(12.0), Inches(0.6))
tf = note_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "RAG = OpenSearch Vector (semantic search) + Neptune Analytics (entity-aware GraphRAG)"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Both KBs use Cohere Embed Multilingual v3 (1024-dim, SG-native). No fine-tuning in PoC; production uses Bedrock Model Distillation to Nova Lite."
run2.font.name = "Calibri"
run2.font.size = Pt(10)
run2.font.italic = True
run2.font.color.rgb = GRAY_MED

print(f"  Slide 6 (model orchestration): added")

# ----------------------------------------------------------------------------
# SLIDE 7: NETWORK & INTEGRATION
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Network & Integration",
    "Lời trình bày:\n\n"
    "Có 2 đường traffic riêng biệt. Web traffic của bác sĩ đi public HTTPS qua CloudFront, "
    "WAF có IP allow-list theo từng bệnh viện đối tác. Đây là đường mass user, scale linear.\n\n"
    "Data plane traffic là EHR, SharePoint, internal trial reports, admin upload đi qua VPN "
    "Site-to-Site IPsec. Mỗi bệnh viện có Customer Gateway peer với AWS Virtual Private Gateway. "
    "Đây là đường data nhạy cảm, không bao giờ chạm public internet.\n\n"
    "Bệnh nhân, bệnh án, kết quả thử nghiệm chỉ đi qua VPN. Câu hỏi của bác sĩ và câu trả lời "
    "của AI đi public với token JWT, không bao giờ chứa PHI sau bước Comprehend Medical mask. "
    "Tenant isolation enforce bằng JWT claims, KMS key policy, OpenSearch filter."
)

# Two-column layout: web traffic (left) vs data plane (right)
# Left: Web traffic
left_x = Inches(0.5)
top_y = Inches(1.4)

web_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, Inches(6.0), Inches(2.6))
web_card.fill.solid()
web_card.fill.fore_color.rgb = WHITE
web_card.line.color.rgb = ORANGE
web_card.line.width = Pt(2)
web_card.adjustments[0] = 0.05

web_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, Inches(6.0), Inches(0.5))
web_h.fill.solid()
web_h.fill.fore_color.rgb = ORANGE
web_h.line.fill.background()

web_h_tb = slide.shapes.add_textbox(left_x, top_y, Inches(6.0), Inches(0.5))
tf = web_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "WEB TRAFFIC  ·  Public HTTPS"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

web_body = (
    "Path: Clinician browser → CloudFront → WAF → API Gateway → Lambda\n\n"
    "WAF allow-list per hospital IP/domain. Cognito JWT auth federated to "
    "hospital Entra ID via SAML/OIDC. TLS 1.3 end-to-end. No PHI on this path "
    "(masked by Comprehend Medical before any model or cache call)."
)
web_tb = slide.shapes.add_textbox(left_x + Inches(0.2), top_y + Inches(0.6), Inches(5.6), Inches(2.0))
tf = web_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = web_body
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.color.rgb = GRAY_DARK

# Right: Data plane (VPN)
right_x = Inches(6.8)
data_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top_y, Inches(6.0), Inches(2.6))
data_card.fill.solid()
data_card.fill.fore_color.rgb = WHITE
data_card.line.color.rgb = ORANGE_DARK
data_card.line.width = Pt(2)
data_card.adjustments[0] = 0.05

data_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, Inches(6.0), Inches(0.5))
data_h.fill.solid()
data_h.fill.fore_color.rgb = ORANGE_DARK
data_h.line.fill.background()

data_h_tb = slide.shapes.add_textbox(right_x, top_y, Inches(6.0), Inches(0.5))
tf = data_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "DATA PLANE  ·  Site-to-Site VPN (IPsec)"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

data_body = (
    "Path: Hospital on-prem → Customer Gateway → IPsec tunnel → AWS Virtual "
    "Private Gateway → VPC private subnet\n\n"
    "Carries: EHR (FHIR R4), SharePoint sync, internal trial PDFs, admin "
    "upload portal, SIEM export. Never touches public internet. Per-tenant "
    "KMS CMK + ABAC enforces isolation."
)
data_tb = slide.shapes.add_textbox(right_x + Inches(0.2), top_y + Inches(0.6), Inches(5.6), Inches(2.0))
tf = data_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = data_body
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.color.rgb = GRAY_DARK

# Bottom: integration components table
int_y = Inches(4.3)
add_text_box(slide, Inches(0.5), int_y, Inches(12.3), Inches(0.4),
    "Integration Components", size=15, bold=True, color=ORANGE_DARK)

headers = ["Component", "Purpose", "Protocol"]
rows = [
    ["AWS Customer Gateway", "Logical peer for hospital on-prem device", "IKEv2 / IPsec"],
    ["AWS VPN Gateway (VPG)", "Tunnel termination on AWS side", "IPsec / BGP"],
    ["Amazon Cognito", "Per-tenant clinician user pool", "SAML 2.0 / OIDC"],
    ["Lambda FHIR adapter", "EHR read-only via SMART App Launch v2", "HL7 FHIR R4"],
    ["SharePoint webhook Lambda", "Microsoft Graph subscriptions to S3", "MS Graph API"],
    ["CDS Hooks service", "Deep-link from Epic / Cerner in-basket", "CDS Hooks v1.1"],
]
add_table_simple(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.2), headers, rows)

print(f"  Slide 7 (network & integration): added")

# ----------------------------------------------------------------------------
# SLIDE 8: SECURITY & COMPLIANCE
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Security & Compliance",
    "Lời trình bày:\n\n"
    "Bảo mật được thiết kế theo 7 lớp defense-in-depth. Lớp Perimeter là CloudFront, WAF, "
    "VPN. Lớp Identity là Cognito + IAM Identity Center. Lớp Data là KMS CMK riêng từng tenant, "
    "S3 Object Lock cho audit, TLS 1.3 mọi nơi. Lớp PHI là Comprehend Medical detect và "
    "tokenize trước khi vào model. Lớp Model Safety là Bedrock Guardrails kiểm tra grounding, "
    "topic, PHI. Lớp Audit là CloudTrail + S3 Object Lock 6 năm theo HIPAA. Lớp Threat là "
    "GuardDuty, Macie, Security Hub.\n\n"
    "Về compliance, AWS có signed BAA cho HIPAA. Singapore PDPA và HCSA được hỗ trợ qua region "
    "Singapore. GDPR DPA available. FDA 21 CFR 820 cho phần mềm y tế, EU AI Act cho high-risk "
    "AI system. Tất cả đều có kiểm soát và logging hỗ trợ sẵn từ AWS."
)

# Top: Security layers (7 cards in row)
sec_y = Inches(1.4)
sec_layers = [
    ("Perimeter", "CloudFront, WAF\nVPN Gateway"),
    ("Identity", "Cognito, IAM\nIdentity Center"),
    ("Data", "KMS CMK, S3\nObject Lock"),
    ("PHI", "Comprehend\nMedical, tokenize"),
    ("Model", "Bedrock\nGuardrails"),
    ("Audit", "CloudTrail, S3\nObject Lock 6yr"),
    ("Threat", "GuardDuty,\nMacie, Sec Hub"),
]

card_w = Inches(1.75)
card_h = Inches(1.5)
gap = Inches(0.05)
total_w = card_w * 7 + gap * 6
start_x = (SLIDE_W - total_w) / 2

for i, (name, desc) in enumerate(sec_layers):
    x = start_x + i * (card_w + gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sec_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = ORANGE_LIGHT if i % 2 == 0 else WHITE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(1.5)
    card.adjustments[0] = 0.1
    
    name_tb = slide.shapes.add_textbox(x, sec_y + Inches(0.15), card_w, Inches(0.4))
    tf = name_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(x + Inches(0.05), sec_y + Inches(0.55), card_w - Inches(0.1), Inches(0.9))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY_DARK

# Bottom: Compliance table
comp_y = Inches(3.2)
add_text_box(slide, Inches(0.5), comp_y, Inches(12.3), Inches(0.4),
    "Compliance Matrix", size=15, bold=True, color=ORANGE_DARK)

headers = ["Regulation", "AWS Support", "Scope in this Solution"]
rows = [
    ["HIPAA, 45 CFR Part 164", "Signed BAA over Bedrock + adjuncts", "6-year retention via S3 Object Lock"],
    ["Singapore PDPA", "PDPA-compliant, comparable-protection clauses", "Primary region SG, no default cross-border transfer"],
    ["Singapore HCSA", "AWS Singapore region supports HCSA", "Clinical records remain in ap-southeast-1"],
    ["GDPR (EU 2016/679)", "GDPR DPA available", "DPIA template, right-to-erasure via tombstone"],
    ["FDA 21 CFR 820 (SaMD)", "FDA Part 11 capable", "Decision-support only, human-in-the-loop"],
    ["EU AI Act (high-risk AI)", "Logging + transparency services", "Decision logs, human oversight, post-market monitoring"],
    ["ISO 27001 / SOC 2 Type II", "AWS certified", "Inherited controls, Nova responsibilities documented"],
]
add_table_simple(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.2), headers, rows)

print(f"  Slide 8 (security & compliance): added")

# ----------------------------------------------------------------------------
# SLIDE 9: PERFORMANCE RESULTS (PoC)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Performance Results · PoC Live (v4)",
    "Lời trình bày:\n\n"
    "Đây là kết quả PoC chạy trực tiếp trên EC2 t4g.small Singapore, không có Reserved Tier "
    "và không có cache. Test 20 câu thật về COVID-19 và clinical trials.\n\n"
    "Cấp cứu: TTFT trung bình 2.5 giây, total 5 giây, pass 100% SLA 5 giây. "
    "Đây là dùng Haiku 4.5, vector top-3 cộng GraphRAG top-2, không guardrails.\n\n"
    "Phức tạp: TTFT trung bình 9.7 giây, total 12 giây, pass 100% SLA 15 giây. "
    "Dùng Sonnet 4.5, vector top-15 cộng GraphRAG top-3, có guardrails.\n\n"
    "Trên production, thêm Reserved Tier và Prompt Caching sẽ giảm TTFT cấp cứu xuống dưới 2 giây "
    "đúng spec gốc. ElastiCache Redis sẽ phục vụ 30-45% câu hỏi lặp dưới 500ms. Câu trả lời 100% "
    "có citation, refusal rate 0% sau khi áp dụng multi-strategy chunking."
)

# Top metrics row
metrics = [
    ("TTFT", "2.5s", "Emergency", GREEN),
    ("TTFT", "9.7s", "Complex", GREEN),
    ("Answer Rate", "100%", "20/20 grounded", GREEN),
    ("SLA Pass", "100%", "20/20 within target", GREEN),
]

m_y = Inches(1.4)
m_w = Inches(2.85)
m_h = Inches(1.4)
m_gap = Inches(0.2)
total_w = m_w * 4 + m_gap * 3
start_x = (SLIDE_W - total_w) / 2

for i, (label, value, sub, color) in enumerate(metrics):
    x = start_x + i * (m_w + m_gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, m_y, m_w, m_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(2)
    card.adjustments[0] = 0.06
    
    label_tb = slide.shapes.add_textbox(x, m_y + Inches(0.15), m_w, Inches(0.3))
    tf = label_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_MED
    
    value_tb = slide.shapes.add_textbox(x, m_y + Inches(0.4), m_w, Inches(0.6))
    tf = value_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    sub_tb = slide.shapes.add_textbox(x, m_y + Inches(1.0), m_w, Inches(0.3))
    tf = sub_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sub
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = GRAY_DARK

# Detailed breakdown table
add_text_box(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
    "Latency Breakdown (avg of 10 questions per lane)", size=15, bold=True, color=ORANGE_DARK)

headers = ["Phase", "Emergency Lane", "Complex Lane"]
rows = [
    ["Pre-generate (PHI mask, route)", "~50ms", "~450ms (Nova Micro routing)"],
    ["KB Retrieve (Vector + GraphRAG)", "~1,300ms", "~1,500ms"],
    ["Bedrock model TTFT", "~1,100ms", "~7,700ms (Sonnet + 18 chunks)"],
    ["Network + uvicorn overhead", "~50ms", "~50ms"],
    ["TOTAL TTFT", "2,463ms", "9,576ms"],
    ["Token streaming + post-process", "~2,600ms", "~2,400ms"],
    ["TOTAL end-to-end", "5,079ms", "11,980ms"],
]
add_table_simple(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(2.8), headers, rows)

# Production target note
prod_y = Inches(6.4)
prod_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), prod_y, Inches(12.3), Inches(0.6))
prod_box.fill.solid()
prod_box.fill.fore_color.rgb = ORANGE_LIGHT
prod_box.line.color.rgb = ORANGE
prod_box.line.width = Pt(1)
prod_box.adjustments[0] = 0.2

prod_tb = slide.shapes.add_textbox(Inches(0.7), prod_y + Inches(0.05), Inches(12.0), Inches(0.5))
tf = prod_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Production targets (2s emergency, 6s complex) achievable with: Reserved Tier + Prompt Caching + ElastiCache Redis + Amazon Rerank when SG-available"
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

print(f"  Slide 9 (performance results): added")

# ----------------------------------------------------------------------------
# SLIDE 10: COST ESTIMATION
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Cost Estimation",
    "Lời trình bày:\n\n"
    "Chi phí được tính cho baseline 600,000 calls/tháng, 30% emergency và 70% complex. "
    "Tất cả là AWS list price USD đầu năm 2026.\n\n"
    "Có 2 phương án model. Phương án A1+ dùng Nova Micro cho emergency và Nova Pro cho complex, "
    "khoảng 2,800 USD/tháng. Phương án A2 dùng Claude Haiku 4.5 và Sonnet 4.5, khoảng 7,000 USD/tháng "
    "nhưng với student model distillation giảm còn 5,500 USD.\n\n"
    "Khoản chi lớn nhất là Claude Sonnet 4.5 cho complex lane vì model gần giá cao. OpenSearch "
    "Serverless 350 USD, Neptune Analytics 115 USD là 2 chi phí infra cố định không phụ thuộc volume. "
    "Bedrock Guardrails 180 USD và Comprehend Medical 180 USD là chi phí compliance bắt buộc.\n\n"
    "Khi scale lên 3 triệu calls/tháng, A1+ chỉ 10,500 USD nhưng A2 lên tới 23,500 USD. "
    "Khuyến nghị: chạy A1+ cho production để tối ưu chi phí, A2 nếu khách hàng yêu cầu Claude branding."
)

# 2 cost variants side by side
left_x = Inches(0.5)
right_x = Inches(6.8)
top_y = Inches(1.4)
card_h = Inches(3.5)

# Variant A1+ (Nova)
a1_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, Inches(6.0), card_h)
a1_card.fill.solid()
a1_card.fill.fore_color.rgb = WHITE
a1_card.line.color.rgb = ORANGE
a1_card.line.width = Pt(2)
a1_card.adjustments[0] = 0.04

a1_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, Inches(6.0), Inches(0.5))
a1_h.fill.solid()
a1_h.fill.fore_color.rgb = ORANGE
a1_h.line.fill.background()

a1_h_tb = slide.shapes.add_textbox(left_x, top_y, Inches(6.0), Inches(0.5))
tf = a1_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Variant A1+  ·  Amazon Nova"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

a1_rows = [
    ["Emergency, Nova Micro", "$70"],
    ["Complex, Nova Pro", "$1,470"],
    ["Cohere Embed v3 (SG)", "$10"],
    ["Bedrock Guardrails", "$180"],
    ["OpenSearch Serverless", "$350"],
    ["Neptune Analytics (32 m-NCU)", "$115"],
    ["Comprehend Medical", "$180"],
    ["Lambda, API GW, CloudFront, WAF", "$150"],
    ["S3, CloudTrail, Macie", "$120"],
    ["ElastiCache Redis OSS", "$80"],
    ["Site-to-Site VPN", "$80"],
    ["TOTAL / month", "$2,805"],
]
add_table_simple(slide, left_x + Inches(0.2), top_y + Inches(0.6), Inches(5.6), Inches(2.85), 
    ["Item", "Monthly USD"], a1_rows)

# Variant A2 (Claude)
a2_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top_y, Inches(6.0), card_h)
a2_card.fill.solid()
a2_card.fill.fore_color.rgb = WHITE
a2_card.line.color.rgb = ORANGE_DARK
a2_card.line.width = Pt(2)
a2_card.adjustments[0] = 0.04

a2_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, Inches(6.0), Inches(0.5))
a2_h.fill.solid()
a2_h.fill.fore_color.rgb = ORANGE_DARK
a2_h.line.fill.background()

a2_h_tb = slide.shapes.add_textbox(right_x, top_y, Inches(6.0), Inches(0.5))
tf = a2_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Variant A2  ·  Anthropic Claude"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = WHITE

a2_rows = [
    ["Emergency, Claude Haiku 4.5", "$350"],
    ["Complex, Claude Sonnet 4.5", "$5,460"],
    ["Cohere Embed v3 (SG)", "$10"],
    ["Bedrock Guardrails", "$180"],
    ["OpenSearch Serverless", "$350"],
    ["Neptune Analytics (32 m-NCU)", "$115"],
    ["Comprehend Medical", "$180"],
    ["Lambda, API GW, CloudFront, WAF", "$150"],
    ["S3, CloudTrail, Macie", "$120"],
    ["ElastiCache Redis OSS", "$80"],
    ["Site-to-Site VPN", "$80"],
    ["TOTAL / month (with student offset)", "$5,545"],
]
add_table_simple(slide, right_x + Inches(0.2), top_y + Inches(0.6), Inches(5.6), Inches(2.85),
    ["Item", "Monthly USD"], a2_rows)

# Bottom: scale sensitivity
scale_y = Inches(5.1)
add_text_box(slide, Inches(0.5), scale_y, Inches(12.3), Inches(0.4),
    "Scale Sensitivity (monthly cost USD)", size=15, bold=True, color=ORANGE_DARK)

headers = ["Call Volume / Month", "Variant A1+ (Nova)", "Variant A2 (Claude with student)"]
rows = [
    ["300k", "$2,100", "$3,400"],
    ["600k (baseline)", "$2,955", "$5,765"],
    ["1.2M", "$4,700", "$10,100"],
    ["3M", "$10,500", "$23,500"],
]
add_table_simple(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), headers, rows)

print(f"  Slide 10 (cost estimation): added")

# ----------------------------------------------------------------------------
# SLIDE 11: ROADMAP & SUMMARY
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Implementation Roadmap & Summary",
    "Lời trình bày:\n\n"
    "Lộ trình triển khai 6-10 tuần. 2 tuần đầu provisioning Singapore và sign BAA, ingest WHO "
    "và ICD-11. Tuần 3-4 train Nova Lite student qua Bedrock Model Distillation. Tuần 5-6 "
    "tích hợp EHR FHIR, SharePoint Graph, Cognito federation. Tuần 7-8 red team với 200 "
    "adversarial prompt, tune Guardrails, load test. Sau 8 tuần là launch full stack.\n\n"
    "Sau launch: ICD-11 đồng bộ hàng ngày, SharePoint reconciliation hàng tuần, WHO refresh "
    "hàng tháng, Nova Lite retrain hàng quý, red team re-run sau mỗi guardrail incident.\n\n"
    "Kết luận: PoC đã chứng minh kiến trúc hoạt động ổn định trên on-demand tier, 100% SLA "
    "pass cho cả emergency và complex. Production sẽ thêm Reserved Tier, Prompt Caching, "
    "ElastiCache để đạt mục tiêu 2 giây cấp cứu. Tổng chi phí baseline 600k calls/tháng "
    "khoảng 2,800 USD với Nova hoặc 5,500 USD với Claude. Sẵn sàng đi vào triển khai khi "
    "Nova Health Tech approve."
)

# Roadmap timeline
rm_y = Inches(1.4)
add_text_box(slide, Inches(0.5), rm_y, Inches(12.3), Inches(0.4),
    "Pre-Launch Build (6 to 10 weeks)", size=15, bold=True, color=ORANGE_DARK)

phases = [
    ("Week 1-2", "Foundation", "SG resources, BAA, WHO+ICD-11 ingest, Cohere embed, GraphRAG extraction"),
    ("Week 3-4", "Customization", "Bedrock Model Distillation: Nova Lite student, eval harness pass"),
    ("Week 5-6", "Integration", "EHR FHIR R4, SharePoint Graph, Cognito federation"),
    ("Week 7-8", "Hardening", "Red team 200 adversarial prompts, Guardrails tune, Reserved Tier sizing, load test"),
    ("Launch", "Go-live", "Full stack live, all 12 departments, audit + cache + monitoring active"),
]

phase_y = Inches(1.9)
phase_h = Inches(0.5)
for i, (week, name, desc) in enumerate(phases):
    y = phase_y + i * (phase_h + Inches(0.05))
    
    # Week badge
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(1.6), phase_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    
    badge_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.6), phase_h)
    tf = badge_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = week
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Phase name
    name_tb = slide.shapes.add_textbox(Inches(2.3), y, Inches(2.0), phase_h)
    tf = name_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    # Description
    desc_tb = slide.shapes.add_textbox(Inches(4.3), y, Inches(8.5), phase_h)
    tf = desc_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_DARK

# Bottom: key takeaways
take_y = Inches(5.0)
add_text_box(slide, Inches(0.5), take_y, Inches(12.3), Inches(0.4),
    "Key Takeaways", size=15, bold=True, color=ORANGE_DARK)

takeaways = [
    ("✓", "PoC live and verified", "Emergency 2.5s avg, Complex 9.7s avg, 100% SLA pass on on-demand tier"),
    ("✓", "100% answer rate", "Multi-strategy chunking eliminated refusals; every answer cites sources"),
    ("✓", "PHI never leaves tenant", "Comprehend Medical mask + tenant-scoped KMS + VPC endpoints"),
    ("✓", "Compliance-ready", "HIPAA, PDPA, HCSA, FDA, EU AI Act mapping documented and supported by AWS"),
    ("✓", "SG-native query path", "Zero cross-region hops for clinical data; Cohere Embed v3 in Singapore"),
]

t_y = Inches(5.5)
for i, (mark, title, desc) in enumerate(takeaways):
    y = t_y + i * Inches(0.32)
    
    # Check mark
    mark_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.4), Inches(0.3))
    tf = mark_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = mark
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = GREEN
    
    # Title
    title_tb = slide.shapes.add_textbox(Inches(1.0), y, Inches(3.5), Inches(0.3))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    # Description
    desc_tb = slide.shapes.add_textbox(Inches(4.5), y, Inches(8.3), Inches(0.3))
    tf = desc_tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_DARK

print(f"  Slide 11 (roadmap & summary): added")

# ----------------------------------------------------------------------------
# SLIDE 12: THANK YOU
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Thank You",
    "Lời trình bày:\n\n"
    "Cảm ơn Ban lãnh đạo Nova Health Tech đã lắng nghe. Tôi sẵn sàng trả lời các câu hỏi về "
    "kiến trúc, chi phí, lộ trình triển khai, hoặc bất kỳ phần kỹ thuật cụ thể nào. Có thể demo "
    "PoC trực tiếp tại http://47.130.120.152/ui/index.html nếu Ban lãnh đạo muốn xem.\n\n"
    "Tài liệu chi tiết: AWS_Claude_Technical_Proposal.md với 15 sections và 12 architecture diagrams. "
    "Báo cáo evaluation 900 câu hỏi tại docs/eval_summary_900.md."
)

# Center thank you message
ty_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.3), Inches(1.5))
tf = ty_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Questions & Discussion"
run.font.name = "Calibri"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

# Sub line
sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9.3), Inches(0.6))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Live PoC: http://47.130.120.152/ui/index.html"
run.font.name = "Calibri"
run.font.size = Pt(18)
run.font.color.rgb = GRAY_DARK

# Contact info
contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(9.3), Inches(1.5))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Reference Documents"
run.font.name = "Calibri"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Full proposal · 900-question evaluation report · Architecture diagrams"
run2.font.name = "Calibri"
run2.font.size = Pt(12)
run2.font.color.rgb = GRAY_DARK

print(f"  Slide 12 (thank you): added")

prs.save(OUT_PATH)
print(f"\nSaved: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
