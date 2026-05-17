"""Build technology explanation slides + RAG vs Fine-tuning comparison slide."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme colors
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK  = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED     = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT   = RGBColor(0xCB, 0xD5, 0xE1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
RED          = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)
GREEN        = RGBColor(0x27, 0xAE, 0x60)
BLUE         = RGBColor(0x21, 0x96, 0xF3)
BLUE_LIGHT   = RGBColor(0xE3, 0xF2, 0xFD)
PURPLE       = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    ft = s.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.28))
    tf = ft.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Nova Health Tech  \u00b7  Clinical GenAI Assistant  \u00b7  AWS Singapore"
    run.font.name = "Calibri"; run.font.size = Pt(9)
    run.font.italic = True; run.font.color.rgb = GRAY_MED
    return s


def slide_title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.65))
    tf = tb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(28)
    run.font.bold = True; run.font.color.rgb = ORANGE_DARK
    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.92), Inches(len(text)*0.18), Inches(0.04))
    ul.fill.solid(); ul.fill.fore_color.rgb = ORANGE; ul.line.fill.background()


def add_text(slide, left, top, width, height, text, size=11, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or GRAY_DARK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb


def card(slide, left, top, width, height, fill, border, radius=0.06, border_pt=1.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(border_pt)
    s.adjustments[0] = radius
    return s


def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()


def bordered(shape, fill_color, line_color, line_pt=1.5, radius=0.06):
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color; shape.line.width = Pt(line_pt)
    shape.adjustments[0] = radius


def header_card(slide, left, top, width, height, title, fill_color, text_color=None):
    text_color = text_color or WHITE
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    h.fill.solid(); h.fill.fore_color.rgb = fill_color; h.line.fill.background()
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = title
    run.font.name = "Calibri"; run.font.size = Pt(13)
    run.font.bold = True; run.font.color.rgb = text_color


def bullet_list(slide, left, top, width, items, size=10.5, color=None,
                bullet_color=None, spacing=0.32):
    color = color or GRAY_DARK
    bullet_color = bullet_color or ORANGE
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(left, top + Inches(i * spacing), width, Inches(spacing))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        b = p.add_run(); b.text = "\u25b8 "
        b.font.name = "Calibri"; b.font.size = Pt(size)
        b.font.bold = True; b.font.color.rgb = bullet_color
        t = p.add_run(); t.text = item
        t.font.name = "Calibri"; t.font.size = Pt(size)
        t.font.color.rgb = color


# =============================================================================
# SLIDE 1: Knowledge Base Vector Store
# =============================================================================
slide = new_slide(prs)
slide_title(slide, "Knowledge Base: Vector Store")

# Concept box (left)
card(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5), BLUE_LIGHT, BLUE, radius=0.05)
header_card(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5), "What is a Vector Store?", BLUE)

add_text(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(1.2),
    "Text is converted into numbers (vectors) that capture meaning. "
    "Similar concepts end up close together in this mathematical space. "
    "When a doctor asks a question, the system finds the most similar chunks instantly.",
    size=11, color=GRAY_DARK)

# How it works steps
steps = [
    ("1", "Ingest", "PDF/API source parsed into text chunks"),
    ("2", "Embed", "Cohere Embed v3 converts each chunk to 1024-dim vector"),
    ("3", "Index", "OpenSearch Serverless stores vectors + metadata"),
    ("4", "Query", "Doctor's question embedded, nearest chunks retrieved"),
    ("5", "Rank", "Top-K chunks passed to LLM as grounding context"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.0 + i * 0.62)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.42), Inches(0.42))
    badge.fill.solid(); badge.fill.fore_color.rgb = BLUE; badge.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.7), y, Inches(0.42), Inches(0.42))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.name = "Calibri"; run.font.size = Pt(12)
    run.font.bold = True; run.font.color.rgb = WHITE
    add_text(slide, Inches(1.25), y, Inches(1.0), Inches(0.42), title,
             size=11, bold=True, color=BLUE)
    add_text(slide, Inches(2.3), y, Inches(3.8), Inches(0.42), desc, size=10.5, color=GRAY_DARK)

# Right: AWS vs Alibaba specs
card(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(2.5), ORANGE_LIGHT, ORANGE, radius=0.05)
header_card(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5), "AWS with Claude", ORANGE)
aws_items = [
    "OpenSearch Serverless (Singapore-native)",
    "Cohere Embed Multilingual v3 (1024-dim)",
    "Multi-strategy chunking: hierarchical (WHO), semantic (trials)",
    "Hybrid kNN + BM25 search",
    "No reranker in SG (production gap)",
]
bullet_list(slide, Inches(6.8), Inches(1.72), Inches(5.8), aws_items, size=10.5, bullet_color=ORANGE)

card(slide, Inches(6.6), Inches(3.75), Inches(6.2), Inches(2.85), PURPLE_LIGHT, PURPLE, radius=0.05)
header_card(slide, Inches(6.6), Inches(3.75), Inches(6.2), Inches(0.5), "Alibaba Cloud", PURPLE)
ali_items = [
    "OpenSearch Vector Search HA (dual-zone, Singapore)",
    "text-embedding-v4 (1024-dim) + tongyi-vision (1152-dim)",
    "Hierarchical chunking: parent 1500 / child 300, 15% overlap",
    "Hybrid BM25 + HNSW via Reciprocal Rank Fusion",
    "qwen3-rerank: top-20 \u2192 top-5 (available in SG)",
]
bullet_list(slide, Inches(6.8), Inches(4.37), Inches(5.8), ali_items, size=10.5, bullet_color=PURPLE)

print("  Slide 1 (Vector Store): done")

# =============================================================================
# SLIDE 2: GraphRAG
# =============================================================================
slide = new_slide(prs)
slide_title(slide, "GraphRAG: Knowledge Graph-Augmented Retrieval")

# Left: concept
card(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5), ORANGE_LIGHT, ORANGE, radius=0.05)
header_card(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5), "Why GraphRAG?", ORANGE)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.4), Inches(1.1),
    "Vector search finds similar text. But clinical reasoning often requires "
    "multi-hop connections: Drug A \u2192 treats Condition B \u2192 contraindicated with Drug C. "
    "GraphRAG captures these entity relationships explicitly.",
    size=11, color=GRAY_DARK)

# Graph concept visual (simple boxes + arrows)
# Entity nodes
entities = [
    (Inches(1.0), Inches(3.1), "Drug:\nCorticosteroid"),
    (Inches(3.2), Inches(2.7), "Condition:\nSevere COVID-19"),
    (Inches(3.2), Inches(3.8), "Guideline:\nWHO 2024"),
    (Inches(1.0), Inches(4.3), "Drug:\nBaricitinib"),
]
for ex, ey, elabel in entities:
    e = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ex, ey, Inches(1.8), Inches(0.65))
    e.fill.solid(); e.fill.fore_color.rgb = ORANGE
    e.line.fill.background(); e.adjustments[0] = 0.15
    tb = slide.shapes.add_textbox(ex, ey, Inches(1.8), Inches(0.65))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = elabel
    run.font.name = "Calibri"; run.font.size = Pt(9)
    run.font.bold = True; run.font.color.rgb = WHITE

# Relation labels
rels = [
    (Inches(2.0), Inches(2.95), "treats"),
    (Inches(2.0), Inches(4.15), "treats"),
    (Inches(3.4), Inches(3.35), "cited by"),
]
for rx, ry, rlabel in rels:
    add_text(slide, rx, ry, Inches(1.0), Inches(0.3), rlabel,
             size=9, italic=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.7), Inches(5.1), Inches(5.4), Inches(0.5),
    "Entity nodes + relation edges enable multi-hop reasoning across clinical knowledge.",
    size=10, italic=True, color=GRAY_MED)

# Right: AWS vs Alibaba
card(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(2.5), ORANGE_LIGHT, ORANGE, radius=0.05)
header_card(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5), "AWS with Claude", ORANGE)
aws_graph = [
    "Neptune Analytics (managed graph DB, Singapore)",
    "Claude 3 Haiku extracts entities + relations at ingest",
    "1,863 Entity nodes + 826 Chunk nodes (PoC, 1 WHO PDF)",
    "Bedrock KB GraphRAG: SEMANTIC search (HYBRID not available)",
    "Complex lane only: top-3 graph hits merged with top-15 vector",
]
bullet_list(slide, Inches(6.8), Inches(1.72), Inches(5.8), aws_graph, size=10.5, bullet_color=ORANGE)

card(slide, Inches(6.6), Inches(3.75), Inches(6.2), Inches(2.85), PURPLE_LIGHT, PURPLE, radius=0.05)
header_card(slide, Inches(6.6), Inches(3.75), Inches(6.2), Inches(0.5), "Alibaba Cloud", PURPLE)
ali_graph = [
    "AnalyticDB for PostgreSQL 7.0 + adbpg_graphrag extension",
    "3-zone HA in Singapore, 4-core 32GB vector-optimized",
    "Drug \u2192 condition \u2192 intervention \u2192 guideline nodes",
    "treats, contraindicates, cites relations",
    "Complex lane: parallel vector + graph retrieval",
]
bullet_list(slide, Inches(6.8), Inches(4.37), Inches(5.8), ali_graph, size=10.5, bullet_color=PURPLE)

print("  Slide 2 (GraphRAG): done")



# =============================================================================
# SLIDE 3: Multi-Agent Architecture (12 departments in PoC)
# =============================================================================
slide = new_slide(prs)
slide_title(slide, 'Multi-Agent Architecture: 12 Specialty Agents (PoC)')

RED_COLOR = RGBColor(0xC0, 0x39, 0x2B)

left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5))
bordered(left_card, ORANGE_LIGHT, ORANGE, line_pt=1.5)
header_card(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5), 'How It Works', ORANGE_DARK)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.4), Inches(0.55),
    'Each department is a specialized agent: its own system prompt, '
    'clinical scope, and model. The router picks the right one.',
    size=11, color=GRAY_DARK)

flow_items = [
    ('Doctors question', ORANGE, False),
    ('PHI mask + cache check', ORANGE, False),
    ('Emergency toggle ON?', RED_COLOR, True),
    ('YES: Emergency agent (Haiku 4.5)', RED_COLOR, False),
    ('NO: Router (Nova Micro, JSON mode)', ORANGE_DARK, False),
    ('Routes to 1 of 12 specialty agents', ORANGE, False),
    ('Specialty agent (Sonnet 4.5) generates answer', ORANGE_DARK, False),
    ('Citations validated + streamed to browser', GREEN, False),
]
for i, (text, color, is_decision) in enumerate(flow_items):
    y = Inches(2.38 + i * 0.46)
    if is_decision:
        shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(0.7), y, Inches(5.2), Inches(0.38))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(5.2), Inches(0.38))
        shape.adjustments[0] = 0.15
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.7), y, Inches(5.2), Inches(0.38))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = 'Calibri'; run.font.size = Pt(9.5)
    run.font.bold = True; run.font.color.rgb = WHITE

right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(5.5))
bordered(right_card, WHITE, ORANGE, line_pt=1.5)
header_card(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5),
            '12 Deployed Department Agents (PoC)', ORANGE)

depts = [
    ('Emergency', 'Haiku 4.5', RED_COLOR),
    ('Cardiology', 'Sonnet 4.5', ORANGE_DARK),
    ('Pulmonology', 'Sonnet 4.5', ORANGE_DARK),
    ('Gastroenterology', 'Sonnet 4.5', ORANGE_DARK),
    ('Nephrology', 'Sonnet 4.5', ORANGE_DARK),
    ('Endocrinology', 'Sonnet 4.5', ORANGE_DARK),
    ('Neurology', 'Sonnet 4.5', ORANGE_DARK),
    ('Infectious Disease', 'Sonnet 4.5', ORANGE_DARK),
    ('Oncology', 'Sonnet 4.5', ORANGE_DARK),
    ('Obstetrics & Gyn', 'Sonnet 4.5', ORANGE_DARK),
    ('Pediatrics', 'Sonnet 4.5', ORANGE_DARK),
    ('Radiology', 'Sonnet 4.5 + Vision', ORANGE),
]
col_w = Inches(1.9); row_h = Inches(0.72)
gap_x = Inches(0.1); gap_y = Inches(0.1)
start_x = Inches(6.7); start_y = Inches(1.72)

for i, (name, model, color) in enumerate(depts):
    col = i % 3; row = i // 3
    x = start_x + col * (col_w + gap_x)
    y = start_y + row * (row_h + gap_y)
    dc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
    bordered(dc, ORANGE_LIGHT if i % 2 == 0 else WHITE, color, line_pt=1.5, radius=0.1)
    add_text(slide, x + Inches(0.08), y + Inches(0.08), col_w - Inches(0.16), Inches(0.35),
             name, size=10, bold=True, color=color)
    add_text(slide, x + Inches(0.08), y + Inches(0.42), col_w - Inches(0.16), Inches(0.25),
             model, size=8.5, italic=True, color=GRAY_MED)

note_y = Inches(6.55)
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), note_y, Inches(12.3), Inches(0.42))
solid(note, ORANGE); note.adjustments[0] = 0.2
add_text(slide, Inches(0.7), note_y + Inches(0.06), Inches(11.9), Inches(0.32),
         'PoC: 12 departments, system prompt per agent, bedrock.converse() directly.  '
         'Production: 40 sub-specialties, Bedrock Agents service with tool calling.',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

print('  Slide 3 (Multi-Agent 12 depts): done')


# =============================================================================
# SLIDE 4: Agentic RAG - PoC vs Production
# =============================================================================
slide = new_slide(prs)
slide_title(slide, 'Agentic RAG: PoC vs Production Architecture')

GREEN_LIGHT  = RGBColor(0xE8, 0xF5, 0xE9)
AMBER        = RGBColor(0xFF, 0x8F, 0x00)
AMBER_LIGHT  = RGBColor(0xFF, 0xF8, 0xE1)

poc_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(5.8), Inches(5.5))
bordered(poc_card, GREEN_LIGHT, GREEN, line_pt=2)
poc_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5))
solid(poc_hdr, GREEN)
add_text(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(0.5),
         'PoC (Deployed Today)', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.4), Inches(0.45),
         'Fixed pipeline - NO Bedrock Agents, NO tool calling',
         size=11, bold=True, color=GREEN)

poc_steps = [
    ('1', 'PHI mask (regex)', GREEN),
    ('2', 'Cache lookup (Redis)', GREEN),
    ('3', 'Lane: emergency or complex (if/else)', GREEN),
    ('4', 'Router: Nova Micro via bedrock.converse()', GREEN),
    ('5', 'Retrieve: Bedrock KB Retrieve API (direct call)', GREEN),
    ('6', 'GraphRAG: Bedrock KB Retrieve API (direct call)', GREEN),
    ('7', 'Generate: bedrock.converse_stream() directly', GREEN),
    ('8', 'Cache write + SSE stream to browser', GREEN),
]
for i, (num, text, color) in enumerate(poc_steps):
    y = Inches(2.25 + i * 0.52)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.36), Inches(0.36))
    solid(badge, color)
    tb = slide.shapes.add_textbox(Inches(0.7), y, Inches(0.36), Inches(0.36))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.name = 'Calibri'; run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = WHITE
    add_text(slide, Inches(1.18), y + Inches(0.04), Inches(4.9), Inches(0.38),
             text, size=10.5, color=GRAY_DARK)

add_text(slide, Inches(0.7), Inches(6.35), Inches(5.4), Inches(0.25),
         'No icd11_lookup, no pubmed_search in PoC. PubMed = future feature.',
         size=9.5, italic=True, color=GRAY_MED)

prod_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.6), Inches(1.1), Inches(6.2), Inches(5.5))
bordered(prod_card, AMBER_LIGHT, AMBER, line_pt=2)
prod_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5))
solid(prod_hdr, AMBER)
add_text(slide, Inches(6.6), Inches(1.1), Inches(6.2), Inches(0.5),
         'Production Target (Bedrock Agents)', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.8), Inches(1.72), Inches(5.8), Inches(0.45),
         'Agentic loop - LLM decides which tools to call',
         size=11, bold=True, color=AMBER)

tools = [
    ('kb_retrieve', 'Vector + GraphRAG KB', 'Semantic search on WHO, trials, protocols'),
    ('graph_retrieve', 'Neptune Analytics', 'Multi-hop entity traversal'),
    ('icd11_lookup', 'WHO ICD-11 API', 'Live disease classification + synonyms'),
    ('pubmed_search', 'NCBI E-utilities', 'Real-time research literature'),
]
add_text(slide, Inches(6.8), Inches(2.25), Inches(5.8), Inches(0.35),
         'Agent tools (Bedrock action groups):', size=11, bold=True, color=AMBER)
for i, (tool_name, source, desc) in enumerate(tools):
    y = Inches(2.65 + i * 0.62)
    tc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y, Inches(5.8), Inches(0.55))
    bordered(tc, WHITE, AMBER, line_pt=1.2, radius=0.1)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), y, Inches(0.06), Inches(0.55))
    solid(acc, AMBER)
    add_text(slide, Inches(7.0), y + Inches(0.04), Inches(1.6), Inches(0.28),
             tool_name, size=10, bold=True, color=AMBER)
    add_text(slide, Inches(7.0), y + Inches(0.3), Inches(1.6), Inches(0.22),
             source, size=9, italic=True, color=GRAY_MED)
    add_text(slide, Inches(8.7), y + Inches(0.12), Inches(3.8), Inches(0.35),
             desc, size=10, color=GRAY_DARK)

add_text(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(0.45),
         'Why not in PoC?', size=11, bold=True, color=AMBER)
reasons = [
    'Bedrock Agent InvokeAgent blocked by IAM trust chain issue',
    'converse_stream() used directly instead (works, faster to build)',
    'Production: resolve IAM + enable full agentic loop',
]
bullet_list(slide, Inches(6.8), Inches(5.68), Inches(5.8), reasons,
            size=10, bullet_color=AMBER, spacing=0.28)

banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.38))
solid(banner, ORANGE); banner.adjustments[0] = 0.2
add_text(slide, Inches(0.7), Inches(6.69), Inches(11.9), Inches(0.32),
         'PoC proves the retrieval + generation quality. Production adds Bedrock Agents for dynamic tool selection.',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

print('  Slide 4 (Agentic RAG PoC vs Prod): done')

# =============================================================================
# SLIDE 5: RAG vs Fine-tuning Comparison (AWS vs Alibaba)
# =============================================================================
slide = new_slide(prs)
slide_title(slide, "RAG vs Fine-tuning: Knowledge Update Strategy")

# Subtitle
add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.35),
    "How each version handles periodic updates: WHO guidelines, PubMed, clinical trial reports",
    size=12, italic=True, color=GRAY_MED)

# ── Column headers ────────────────────────────────────────────────────────────
# Dimension column
add_text(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.45),
    "Dimension", size=13, bold=True, color=ORANGE_DARK)
ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.95), Inches(2.5), Inches(0.03))
ul.fill.solid(); ul.fill.fore_color.rgb = ORANGE; ul.line.fill.background()

# AWS header
aws_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3.2), Inches(1.45), Inches(4.5), Inches(0.55))
aws_hdr.fill.solid(); aws_hdr.fill.fore_color.rgb = ORANGE
aws_hdr.line.fill.background(); aws_hdr.adjustments[0] = 0.1
tb = slide.shapes.add_textbox(Inches(3.2), Inches(1.45), Inches(4.5), Inches(0.55))
tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "AWS with Claude"
run.font.name = "Calibri"; run.font.size = Pt(14); run.font.bold = True
run.font.color.rgb = WHITE

# Alibaba header
ali_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.0), Inches(1.45), Inches(4.8), Inches(0.55))
ali_hdr.fill.solid(); ali_hdr.fill.fore_color.rgb = PURPLE
ali_hdr.line.fill.background(); ali_hdr.adjustments[0] = 0.1
tb = slide.shapes.add_textbox(Inches(8.0), Inches(1.45), Inches(4.8), Inches(0.55))
tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "Alibaba Cloud"
run.font.name = "Calibri"; run.font.size = Pt(14); run.font.bold = True
run.font.color.rgb = WHITE

# ── Comparison rows ───────────────────────────────────────────────────────────
rows = [
    (
        "WHO Guidelines\n(monthly update)",
        "RAG only\nEventBridge cron \u2192 S3 \u2192 BDA parse \u2192 Cohere embed \u2192 OpenSearch upsert\nCache invalidated by source tag. No model retrain.",
        "RAG only\nCloudOps Scheduler \u2192 OSS \u2192 DocMind parse \u2192 text-embedding-v4 \u2192 OpenSearch upsert\nTair cache flushed by source tag. No model retrain.",
        GREEN,
    ),
    (
        "PubMed\n(real-time)",
        "RAG tool (query-time)\npubmed_search tool called by agent on demand\nNCBI E-utilities, 24h result cache\nNever ingested into KB",
        "RAG tool (query-time)\npubmed_search tool called by agent on demand\nNCBI E-utilities, 3 req/s free tier\nNever ingested into KB",
        GREEN,
    ),
    (
        "Clinical Trial\nReports (weekly)",
        "RAG only\nSharePoint Graph webhook \u2192 S3 \u2192 PHI mask \u2192 BDA parse\n\u2192 Cohere embed \u2192 OpenSearch upsert\nWeekly reconciliation + real-time webhook",
        "RAG only\nSharePoint Graph webhook \u2192 OSS \u2192 SDDP PHI mask \u2192 DocMind parse\n\u2192 text-embedding-v4 \u2192 OpenSearch upsert\nWeekly reconciliation + real-time webhook",
        GREEN,
    ),
    (
        "Clinical Tone\n& Phrasing",
        "Fine-tuning (quarterly)\nBedrock Model Distillation: Sonnet 4.5 \u2192 Nova Lite student\nSFT on de-identified invocation logs\nStudent serves 40% of complex traffic",
        "Fine-tuning (quarterly)\nPAI SFT + LoRA on Qwen3-8B\nDPO monthly on clinician preference pairs\nGRPO ad-hoc for tool-calling\nStudent serves 60% of complex traffic",
        ORANGE,
    ),
    (
        "ICD-11 Codes\n(daily)",
        "RAG (daily delta)\nEventBridge 02:00 SGT \u2192 WHO OAuth2 API \u2192 S3 \u2192 embed \u2192 upsert\nicd11_lookup tool for query-time expansion",
        "RAG (daily delta)\nCloudOps Scheduler 02:00 SGT \u2192 WHO OAuth2 API \u2192 OSS \u2192 embed \u2192 upsert\nicd11_lookup tool for query-time expansion",
        GREEN,
    ),
]

row_h = Inches(0.88)
start_y = Inches(2.1)

for i, (dim, aws_text, ali_text, indicator) in enumerate(rows):
    y = start_y + i * (row_h + Inches(0.04))
    bg_color = ORANGE_LIGHT if i % 2 == 0 else WHITE

    # Dimension cell
    dim_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(2.5), row_h)
    dim_bg.fill.solid(); dim_bg.fill.fore_color.rgb = bg_color
    dim_bg.line.color.rgb = GRAY_LIGHT; dim_bg.line.width = Pt(0.5)
    add_text(slide, Inches(0.6), y + Inches(0.1), Inches(2.3), row_h - Inches(0.2),
             dim, size=10.5, bold=True, color=GRAY_DARK)

    # AWS cell
    aws_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), y, Inches(4.5), row_h)
    aws_bg.fill.solid(); aws_bg.fill.fore_color.rgb = bg_color
    aws_bg.line.color.rgb = GRAY_LIGHT; aws_bg.line.width = Pt(0.5)
    add_text(slide, Inches(3.3), y + Inches(0.06), Inches(4.3), row_h - Inches(0.12),
             aws_text, size=9.5, color=GRAY_DARK)

    # Alibaba cell
    ali_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), y, Inches(4.8), row_h)
    ali_bg.fill.solid(); ali_bg.fill.fore_color.rgb = bg_color
    ali_bg.line.color.rgb = GRAY_LIGHT; ali_bg.line.width = Pt(0.5)
    add_text(slide, Inches(8.1), y + Inches(0.06), Inches(4.6), row_h - Inches(0.12),
             ali_text, size=9.5, color=GRAY_DARK)

    # Indicator dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(2.85), y + Inches(0.32), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = indicator; dot.line.fill.background()

# Legend
add_text(slide, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.35),
    "\u25cf Green = RAG (no model retrain, updates in minutes-hours)   "
    "\u25cf Orange = Fine-tuning (tone/format only, quarterly, ~$15-40/run)",
    size=10, color=GRAY_MED, align=PP_ALIGN.CENTER)

# Key insight box
insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.38))
insight.fill.solid(); insight.fill.fore_color.rgb = ORANGE
insight.line.fill.background(); insight.adjustments[0] = 0.2
add_text(slide, Inches(0.7), Inches(7.04), Inches(11.9), Inches(0.32),
    "Key insight: RAG handles all factual updates (WHO, PubMed, trials). "
    "Fine-tuning is ONLY for tone, phrasing, and clinical vocabulary. Never for facts.",
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

print("  Slide 5 (RAG vs Fine-tuning): done")

# =============================================================================
# Save
# =============================================================================
# =============================================================================
# SLIDE 6: LangChain + LangGraph in the PoC
# =============================================================================
slide = new_slide(prs)
slide_title(slide, 'LangChain + LangGraph in the PoC')

BLUE         = RGBColor(0x21, 0x96, 0xF3)
BLUE_LIGHT   = RGBColor(0xE3, 0xF2, 0xFD)
PURPLE       = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)

lg_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(6.0), Inches(3.2))
bordered(lg_card, BLUE_LIGHT, BLUE, line_pt=2)
lg_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.5))
solid(lg_hdr, BLUE)
add_text(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.5),
         'LangGraph  -  State Machine Orchestration', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.7), Inches(1.72), Inches(5.6), Inches(0.45),
         'Defines the request pipeline as a directed graph of nodes + edges. '
         'Each node is a processing step; edges are conditional transitions.',
         size=10.5, color=GRAY_DARK)

code_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.25), Inches(5.6), Inches(1.8))
code_bg.fill.solid(); code_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
code_bg.line.color.rgb = GRAY_LIGHT; code_bg.line.width = Pt(0.75)

code_lines = [
    'from langgraph.graph import END, StateGraph',
    '',
    'g = StateGraph(ChatState)',
    'g.add_node("phi_mask",  _node_phi_mask)',
    'g.add_node("pick_lane", _node_pick_lane)',
    'g.add_node("retrieve",  _node_retrieve)',
    'g.add_node("generate",  _node_generate)',
    'g.add_conditional_edges("cache_lookup",',
    '    _branch_on_lane,',
    '    {"emergency": "emergency_agent",',
    '     "complex":   "route_department"})',
    'graph = g.compile()',
]
for i, line in enumerate(code_lines):
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.32 + i * 0.13), Inches(5.3), Inches(0.14))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = line
    run.font.name = 'Courier New'; run.font.size = Pt(8.5)
    run.font.color.rgb = BLUE if (line.startswith('from') or line.startswith('g.') or line.startswith('graph')) else GRAY_DARK

lc_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8), Inches(1.1), Inches(6.0), Inches(3.2))
bordered(lc_card, PURPLE_LIGHT, PURPLE, line_pt=2)
lc_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.5))
solid(lc_hdr, PURPLE)
add_text(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.5),
         'LangChain  -  Text Splitting + Cache', size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(7.0), Inches(1.72), Inches(5.6), Inches(0.45),
         'Used for two specific utilities: semantic text splitting during ingestion, '
         'and the Redis semantic cache layer.',
         size=10.5, color=GRAY_DARK)

lc_items = [
    ('langchain-text-splitters',
     'SemanticChunker for clinical trial PDFs\n(max 512 tokens, 80th percentile breakpoint)'),
    ('langchain-community',
     'Redis cache integration (SHA-256 key in PoC,\nvector similarity in production)'),
    ('langchain-core',
     'Base types and interfaces used by LangGraph'),
]
for i, (pkg, desc) in enumerate(lc_items):
    y = Inches(2.28 + i * 0.65)
    pkg_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), y, Inches(5.6), Inches(0.58))
    bordered(pkg_bg, WHITE, PURPLE, line_pt=1.2, radius=0.08)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), y, Inches(0.06), Inches(0.58))
    solid(acc, PURPLE)
    add_text(slide, Inches(7.15), y + Inches(0.04), Inches(2.0), Inches(0.25),
             pkg, size=9.5, bold=True, color=PURPLE)
    add_text(slide, Inches(7.15), y + Inches(0.28), Inches(5.3), Inches(0.28),
             desc, size=9, italic=True, color=GRAY_DARK)

flow_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.1))
bordered(flow_bg, ORANGE_LIGHT, ORANGE, line_pt=1.5)
add_text(slide, Inches(0.7), Inches(4.52), Inches(12.0), Inches(0.38),
         'LangGraph State Machine: Full Node Graph (PoC)', size=13, bold=True, color=ORANGE_DARK)

nodes = [
    ('phi_mask',        Inches(0.65),  Inches(5.0),  ORANGE),
    ('pick_lane',       Inches(2.05),  Inches(5.0),  ORANGE),
    ('cache_lookup',    Inches(3.45),  Inches(5.0),  ORANGE),
    ('emergency_agent', Inches(2.05),  Inches(5.72), BLUE),
    ('route_dept',      Inches(4.85),  Inches(5.72), BLUE),
    ('retrieve',        Inches(6.25),  Inches(5.0),  GREEN),
    ('generate',        Inches(7.85),  Inches(5.0),  GREEN),
    ('cache_write',     Inches(9.45),  Inches(5.0),  ORANGE),
    ('END',             Inches(11.05), Inches(5.0),  GRAY_MED),
]
node_w = Inches(1.25); node_h = Inches(0.45)
for name, nx, ny, color in nodes:
    n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, node_w, node_h)
    n.fill.solid(); n.fill.fore_color.rgb = color; n.line.fill.background()
    n.adjustments[0] = 0.15
    tb = slide.shapes.add_textbox(nx, ny, node_w, node_h)
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = name
    run.font.name = 'Courier New'; run.font.size = Pt(8.5)
    run.font.bold = True; run.font.color.rgb = WHITE

add_text(slide, Inches(3.45), Inches(5.5), Inches(1.5), Inches(0.22),
         'emergency', size=8, italic=True, color=BLUE)
add_text(slide, Inches(3.45), Inches(5.65), Inches(1.5), Inches(0.22),
         'complex', size=8, italic=True, color=BLUE)

note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.38))
solid(note, ORANGE); note.adjustments[0] = 0.2
add_text(slide, Inches(0.7), Inches(6.69), Inches(11.9), Inches(0.32),
         'LangGraph = orchestration (state machine).  '
         'LangChain = utilities (text splitting, cache).  '
         'Neither replaces Bedrock Agents for tool calling.',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

print('  Slide 6 (LangChain + LangGraph): done')


# =============================================================================
# SLIDE 7: Compliance, Regulations & Security Standards
# =============================================================================
slide = new_slide(prs)
slide_title(slide, 'Compliance, Regulations & Security')

TEAL         = RGBColor(0x00, 0x89, 0x7B)
TEAL_LIGHT   = RGBColor(0xE0, 0xF2, 0xF1)
BLUE         = RGBColor(0x21, 0x96, 0xF3)
BLUE_LIGHT   = RGBColor(0xE3, 0xF2, 0xFD)
PURPLE       = RGBColor(0x7B, 0x1F, 0xA2)
PURPLE_LIGHT = RGBColor(0xF3, 0xE5, 0xF5)
AMBER        = RGBColor(0xFF, 0x8F, 0x00)
AMBER_LIGHT  = RGBColor(0xFF, 0xF8, 0xE1)
RED_COLOR    = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)

# ============================================================
# ROW 1: Singapore Regulations (top-left 3 cards)
# ============================================================
sg_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(7.8), Inches(0.38))
solid(sg_label, TEAL)
add_text(slide, Inches(0.5), Inches(1.1), Inches(7.8), Inches(0.38),
         'Singapore Healthcare Regulations', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

sg_regs = [
    ('PDPA', 'Personal Data Protection Act',
     'Patient data protection, 72h breach notification, DPO required',
     'PHI mask + KMS BYOK + audit trail'),
    ('HCSA 2020', 'Healthcare Services Act',
     'License for clinical decision support (replaces PHMC 1980)',
     'Decision support only, human-in-the-loop, audit ready'),
    ('HIA 2026', 'Health Information Act',
     'Mandatory NEHR contribution (effective early 2027)',
     'NEHR connector planned Year 2, consent + audit trail'),
    ('AIHGle 2.0', 'AI in Healthcare Guidelines\n(MOH+HSA, March 2026)',
     '7 principles: safety, fairness, transparency, explainability...',
     'All 7 principles addressed in architecture'),
]
sg_w = Inches(1.85); sg_h = Inches(1.55); sg_gap = Inches(0.08)
sg_start_x = Inches(0.5); sg_y = Inches(1.55)

for i, (code, name, scope, impl) in enumerate(sg_regs):
    x = sg_start_x + i * (sg_w + sg_gap)
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sg_y, sg_w, sg_h)
    bordered(c, TEAL_LIGHT, TEAL, line_pt=1.5, radius=0.06)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, sg_y, sg_w, Inches(0.38))
    solid(hdr, TEAL)
    add_text(slide, x, sg_y, sg_w, Inches(0.38),
             code, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.08), sg_y + Inches(0.42), sg_w - Inches(0.16), Inches(0.32),
             name, size=8.5, bold=True, color=TEAL)
    add_text(slide, x + Inches(0.08), sg_y + Inches(0.76), sg_w - Inches(0.16), Inches(0.38),
             scope, size=8, italic=True, color=GRAY_DARK)
    add_text(slide, x + Inches(0.08), sg_y + Inches(1.18), sg_w - Inches(0.16), Inches(0.32),
             impl, size=8, color=GREEN, bold=True)

# ============================================================
# ROW 1 RIGHT: Clinical Standards
# ============================================================
cl_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(8.55), Inches(1.1), Inches(4.3), Inches(0.38))
solid(cl_label, BLUE)
add_text(slide, Inches(8.55), Inches(1.1), Inches(4.3), Inches(0.38),
         'Clinical & International Standards', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cl_items = [
    ('HSA SaMD Class B', 'Software as Medical Device registration\nvia SHARE platform (July 2025)'),
    ('HIPAA 164.530(j)', '6-year audit retention via S3 Object Lock\nSigned BAA over Bedrock + S3 + Lambda'),
    ('ISO 27001 / SOC 2', 'Inherited from AWS Singapore\nNo separate Nova certification needed'),
    ('IMDA AI Verify', 'Responsible AI framework\nSelf-assessment planned Month 5'),
]
cl_y = Inches(1.55)
for i, (std, desc) in enumerate(cl_items):
    y = cl_y + i * Inches(0.38)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.55), y, Inches(4.3), Inches(0.35))
    bordered(bg, BLUE_LIGHT if i % 2 == 0 else WHITE, BLUE, line_pt=1, radius=0.08)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.55), y, Inches(0.05), Inches(0.35))
    solid(acc, BLUE)
    add_text(slide, Inches(8.65), y + Inches(0.02), Inches(1.3), Inches(0.18),
             std, size=9, bold=True, color=BLUE)
    add_text(slide, Inches(10.0), y + Inches(0.02), Inches(2.75), Inches(0.32),
             desc, size=8, italic=True, color=GRAY_DARK)

# ============================================================
# ROW 2: Shared Responsibility Model
# ============================================================
sr_y = Inches(3.28)
sr_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), sr_y, Inches(12.3), Inches(0.38))
solid(sr_label, ORANGE_DARK)
add_text(slide, Inches(0.5), sr_y, Inches(12.3), Inches(0.38),
         'AWS Shared Responsibility Model', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# AWS responsibility (left)
aws_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), sr_y + Inches(0.42), Inches(5.8), Inches(1.85))
bordered(aws_card, ORANGE_LIGHT, ORANGE, line_pt=1.5)
aws_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), sr_y + Inches(0.42), Inches(5.8), Inches(0.38))
solid(aws_hdr, ORANGE)
add_text(slide, Inches(0.5), sr_y + Inches(0.42), Inches(5.8), Inches(0.38),
         'AWS Responsibility  (Security OF the Cloud)', size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

aws_items = [
    'Physical data center security (Singapore ap-southeast-1)',
    'Hardware, network, hypervisor infrastructure',
    'Managed service availability (Bedrock, OpenSearch, Neptune)',
    'ISO 27001, SOC 2, PCI-DSS certifications',
    'Compliance with Singapore data residency requirements',
]
for i, item in enumerate(aws_items):
    tb = slide.shapes.add_textbox(Inches(0.65), sr_y + Inches(0.88 + i * 0.27), Inches(5.5), Inches(0.26))
    tf = tb.text_frame; p = tf.paragraphs[0]
    b = p.add_run(); b.text = '\u25b8 '
    b.font.name = 'Calibri'; b.font.size = Pt(10); b.font.bold = True; b.font.color.rgb = ORANGE
    t = p.add_run(); t.text = item
    t.font.name = 'Calibri'; t.font.size = Pt(9.5); t.font.color.rgb = GRAY_DARK

# Nova responsibility (middle)
nova_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.55), sr_y + Inches(0.42), Inches(3.0), Inches(1.85))
bordered(nova_card, PURPLE_LIGHT, PURPLE, line_pt=1.5)
nova_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.55), sr_y + Inches(0.42), Inches(3.0), Inches(0.38))
solid(nova_hdr, PURPLE)
add_text(slide, Inches(6.55), sr_y + Inches(0.42), Inches(3.0), Inches(0.38),
         'Nova Responsibility', size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

nova_items = [
    'PHI masking (Comprehend Medical)',
    'Application security + PDPA',
    'HCSA license + audit trail',
    'Guardrails + citation validator',
    'Incident response + DPO',
]
for i, item in enumerate(nova_items):
    tb = slide.shapes.add_textbox(Inches(6.7), sr_y + Inches(0.88 + i * 0.27), Inches(2.7), Inches(0.26))
    tf = tb.text_frame; p = tf.paragraphs[0]
    b = p.add_run(); b.text = '\u25b8 '
    b.font.name = 'Calibri'; b.font.size = Pt(10); b.font.bold = True; b.font.color.rgb = PURPLE
    t = p.add_run(); t.text = item
    t.font.name = 'Calibri'; t.font.size = Pt(9.5); t.font.color.rgb = GRAY_DARK

# Hospital responsibility (right)
hosp_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.8), sr_y + Inches(0.42), Inches(3.0), Inches(1.85))
bordered(hosp_card, AMBER_LIGHT, AMBER, line_pt=1.5)
hosp_hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(9.8), sr_y + Inches(0.42), Inches(3.0), Inches(0.38))
solid(hosp_hdr, AMBER)
add_text(slide, Inches(9.8), sr_y + Inches(0.42), Inches(3.0), Inches(0.38),
         'Hospital Responsibility', size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

hosp_items = [
    'Clinical decisions (physician)',
    'HCSA license (deployer)',
    'Staff training + consent',
    'EHR access controls',
    'Adverse event reporting',
]
for i, item in enumerate(hosp_items):
    tb = slide.shapes.add_textbox(Inches(9.95), sr_y + Inches(0.88 + i * 0.27), Inches(2.7), Inches(0.26))
    tf = tb.text_frame; p = tf.paragraphs[0]
    b = p.add_run(); b.text = '\u25b8 '
    b.font.name = 'Calibri'; b.font.size = Pt(10); b.font.bold = True; b.font.color.rgb = AMBER
    t = p.add_run(); t.text = item
    t.font.name = 'Calibri'; t.font.size = Pt(9.5); t.font.color.rgb = GRAY_DARK

# ============================================================
# ROW 3: Security Controls (bottom)
# ============================================================
sec_y = Inches(5.38)
sec_label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.5), sec_y, Inches(12.3), Inches(0.38))
solid(sec_label, GRAY_DARK)
add_text(slide, Inches(0.5), sec_y, Inches(12.3), Inches(0.38),
         'Security Controls in the Solution', size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

sec_controls = [
    ('\U0001f512', 'PHI Masking', 'Comprehend Medical\ntokenizes before model'),
    ('\U0001f5dd\ufe0f', 'KMS BYOK', 'Hospital controls\nencryption keys'),
    ('\U0001f4cb', 'Audit Trail', 'S3 Object Lock\n6-year WORM retention'),
    ('\U0001f6e1\ufe0f', 'Guardrails', 'Bedrock Guardrails\ngrounding + PHI filter'),
    ('\U0001f310', 'Data Residency', '100% ap-southeast-1\nzero cross-border'),
    ('\U0001f465', 'Access Control', 'Cognito + ABAC\ntenant isolation'),
]
ctrl_w = Inches(1.95); ctrl_h = Inches(0.88); ctrl_gap = Inches(0.1)
ctrl_start_x = Inches(0.5); ctrl_y = sec_y + Inches(0.42)

for i, (icon, title, desc) in enumerate(sec_controls):
    x = ctrl_start_x + i * (ctrl_w + ctrl_gap)
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, ctrl_y, ctrl_w, ctrl_h)
    bordered(c, ORANGE_LIGHT if i % 2 == 0 else WHITE, ORANGE, line_pt=1.2, radius=0.1)
    icon_tb = slide.shapes.add_textbox(x, ctrl_y + Inches(0.04), ctrl_w, Inches(0.38))
    tf = icon_tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = icon
    run.font.name = 'Segoe UI Emoji'; run.font.size = Pt(20)
    add_text(slide, x, ctrl_y + Inches(0.42), ctrl_w, Inches(0.2),
             title, size=9.5, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x, ctrl_y + Inches(0.62), ctrl_w, Inches(0.24),
             desc, size=8, italic=True, color=GRAY_MED, align=PP_ALIGN.CENTER)

# Bottom note
note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.38))
solid(note, ORANGE); note.adjustments[0] = 0.2
add_text(slide, Inches(0.7), Inches(6.69), Inches(11.9), Inches(0.32),
         'AWS secures the cloud infrastructure.  Nova secures the application + data.  Hospital secures clinical decisions + staff.',
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

print('  Slide 7 (Compliance + Security): done')


out = "docs/Technology_Explanation_Slides.pptx"
prs.save(out)
size = Path(out).stat().st_size
print(f"\nSaved: {out}  ({size:,} bytes)")
print(f"Total slides: {len(prs.slides)}")
