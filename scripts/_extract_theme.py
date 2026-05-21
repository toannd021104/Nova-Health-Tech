"""Extract theme colours and layout from Nova_Health_Tech_final.pptx."""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import pptx.enum.dml as dml_enum

prs = Presentation(r'docs/Nova_Health_Tech_final.pptx')
print(f"Slides: {len(prs.slides)}, Size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}")

for i, slide in enumerate(prs.slides):
    print(f"\n=== Slide {i+1} ===")
    for shape in slide.shapes:
        # Fill colour
        try:
            fill = shape.fill
            if fill.type and fill.type.name == 'SOLID':
                print(f"  SHAPE fill={fill.fore_color.rgb} name={shape.name[:30]}")
        except Exception:
            pass
        # Text
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text.strip()
                    if not t:
                        continue
                    try:
                        col = run.font.color.rgb if run.font.color.type else "theme"
                    except Exception:
                        col = "n/a"
                    sz = run.font.size
                    sz_pt = round(sz.pt, 1) if sz else None
                    print(f"  TEXT '{t[:50].encode('ascii','replace').decode()}' size={sz_pt} bold={run.font.bold} color={col}")
    if i >= 5:
        break
