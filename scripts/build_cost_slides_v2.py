"""Cost comparison slides matching Nova Health Tech theme.

Theme extracted from Nova_Health_Tech_final.pptx:
  Primary orange:  #FF6B35  (section header bars)
  Dark orange:     #CC4F1F  (sub-headers, card accents)
  Light orange bg: #FFE8DA  (card backgrounds)
  Purple title:    #5D1DC7  (main titles)
  Dark text:       #2C3E50  (body text)
  White:           #FFFFFF
  Light gray:      #F5F5F5  (alternating rows)
  Mid gray:        #7F8C8D  (muted text)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import copy

TEMPLATE = Path("docs/Nova_Health_Tech_final.pptx")
OUT      = Path("docs/Production_Cost_Comparison_3Versions.pptx")

# ── Theme colours ─────────────────────────────────────────────────────────────
ORANGE   = RGBColor(0xFF, 0x6B, 0x35)   # primary orange
DORANGE  = RGBColor(0xCC, 0x4F, 0x1F)   # dark orange
LORANGE  = RGBColor(0xFF, 0xE8, 0xDA)   # light orange bg
PURPLE   = RGBColor(0x5D, 0x1D, 0xC7)   # title purple
LPURPLE  = RGBColor(0xED, 0xE7, 0xF6)   # light purple bg
DARK     = RGBColor(0x2C, 0x3E, 0x50)   # body text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY    = RGBColor(0x7F, 0x8C, 0x8D)
DGRAY    = RGBColor(0x44, 0x44, 0x44)
GREEN    = RGBColor(0x27, 0xAE, 0x60)   # savings / positive
LGREEN   = RGBColor(0xE8, 0xF5, 0xE9)
RED      = RGBColor(0xC0, 0x00, 0x00)
YELLOW   = RGBColor(0xFF, 0xC1, 0x07)

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0):
    from pptx.util import Pt as _Pt
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color and line_pt:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h,
             size=10, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, fill=None, line_color=None, line_pt=0,
             wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
    if line_color and line_pt:
        tb.line.color.rgb = line_color; tb.line.width = Pt(line_pt)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def section_bar(slide, title, y=1.05):
    """Orange section header bar matching the deck style."""
    add_rect(slide, 0, y, 13.33, 0.42, fill=ORANGE)
    add_text(slide, title, 0.25, y+0.04, 12.8, 0.34,
             size=16, bold=True, color=WHITE)

def slide_title(slide, title, subtitle=""):
    """Purple title + orange underline bar."""
    add_rect(slide, 0, 0, 13.33, 0.08, fill=ORANGE)   # top accent strip
    add_text(slide, title, 0.3, 0.12, 12.5, 0.55,
             size=24, bold=True, color=PURPLE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.65, 12.5, 0.32,
                 size=11, italic=True, color=MGRAY)
    add_rect(slide, 0, 0.95, 13.33, 0.04, fill=ORANGE)  # divider

def footer_bar(slide, note=""):
    pass  # Footer removed per user request

def assumption_strip(slide):
    add_rect(slide, 0, 1.0, 13.33, 0.38, fill=LPURPLE)
    add_text(slide,
        "Assumptions: 500 clinicians  |  10 queries/day  |  150,000 queries/month  "
        "(30k emergency + 120k complex)  |  Avg 3,000 in / 600 out tokens (complex)  "
        "|  1,500 in / 300 out (emergency)  |  1 fine-tuning run/month  |  Singapore region",
        0.25, 1.02, 12.8, 0.32, size=8.5, italic=True, color=PURPLE)

# ── Open template and add blank slides ────────────────────────────────────────
prs = Presentation(TEMPLATE)
blank_layout = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(blank_layout)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — SUMMARY COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()
slide_title(s1, "Production Cost Comparison — 3 Versions",
            "Monthly estimate · 500 clinicians · 150,000 queries/month · Singapore region · All prices USD list price")
assumption_strip(s1)
footer_bar(s1, "Slide 1 of 4")

# Layout: label col 2.8" + 3 version cols 3.3" each = 12.7" total, starting at 0.3"
LX = 0.3    # left edge
LW = 2.8    # label column width
VW = 3.3    # version column width
V1 = LX + LW           # Version A starts
V2 = V1 + VW           # Version B starts
V3 = V2 + VW           # Version C starts

# Column headers row
add_rect(s1, LX, 1.45, LW, 0.96, fill=DARK)
add_text(s1, "Cost Category", LX+0.1, 1.55, LW-0.15, 0.36,
         size=11, bold=True, color=WHITE)
add_text(s1, "What it covers", LX+0.1, 1.88, LW-0.15, 0.28,
         size=8, italic=True, color=rgb(0xBB,0xBB,0xBB))

ver_cols = [
    (V1, PURPLE,  LPURPLE, "Version A",  "AWS + Claude",    "$2,350–$2,600"),
    (V2, DORANGE, LORANGE, "Version B",  "AWS + Qwen",      "$2,150–$2,350"),
    (V3, ORANGE,  rgb(0xFF,0xF3,0xE0), "Version C", "Alibaba Cloud", "$2,100–$2,300"),
]
for x, hc, bc, ver, stack, total in ver_cols:
    add_rect(s1, x, 1.45, VW, 0.52, fill=hc)
    add_text(s1, ver,   x+0.1, 1.47, VW-0.15, 0.26, size=12, bold=True, color=WHITE)
    add_text(s1, stack, x+0.1, 1.73, VW-0.15, 0.22, size=9,  color=WHITE)
    add_rect(s1, x, 1.97, VW, 0.44, fill=bc, line_color=hc, line_pt=1.5)
    add_text(s1, total, x+0.1, 2.00, VW-0.15, 0.38, size=14, bold=True,
             color=hc, align=PP_ALIGN.CENTER)

# Data rows: (label, description, ver_a, ver_b, ver_c)
rows = [
    ("LLM Inference",
     "Router + specialist agents\n+ emergency model",
     "$1,455 / mo", "$553 / mo", "$107 / mo"),
    ("Student Inference Endpoint",
     "Fine-tuned model (SageMaker scale-to-0\n/ PAI-EAS min=1 auto-scale)",
     "—\n(no student model)", "$710 / mo\n(360hr, scale-to-0)", "$918 / mo\n(min=1, 85% util)"),
    ("RAG / Knowledge Base",
     "Vector store + GraphRAG\n(OpenSearch + Neptune/AnalyticDB)",
     "$461 / mo", "$461 / mo", "$519 / mo"),
    ("Semantic Cache",
     "Redis/Tair in-memory cache\n(10min emergency / 24hr complex TTL)",
     "$120 / mo", "$120 / mo", "$60 / mo"),
    ("Fine-tuning (monthly)",
     "Data generation + GPU training\n+ model storage",
     "—\n(no fine-tuning)", "$17 / mo", "$9 / mo"),
    ("Compute / Serverless",
     "API server EC2 / Function Compute\n+ API Gateway",
     "$67 / mo", "$67 / mo", "$27 / mo"),
    ("Storage + Logs",
     "S3/OSS corpus + audit logs\n+ WORM archive",
     "$7 / mo", "$7 / mo", "$8 / mo"),
    ("Security / Observability",
     "WAF + Guardrails + IAM\n+ CloudWatch / ARMS",
     "$1 / mo", "$1 / mo", "$1 / mo"),
    ("Network / VPN / CDN",
     "VPN Gateway (hospital S2S)\n+ CDN + Anti-DDoS",
     "$42 / mo", "$42 / mo", "$158 / mo"),
    ("IdaaS / SSO",
     "Hospital IdP federation\n(SAML/OIDC)",
     "—", "—", "$50 / mo"),
]

y = 2.48
rh = 0.46
for i, (label, desc, a, b, c) in enumerate(rows):
    bg = LGRAY if i % 2 == 0 else WHITE
    # Label column
    add_rect(s1, LX, y, 0.08, rh, fill=ORANGE if i % 2 == 0 else DORANGE)
    add_rect(s1, LX+0.08, y, LW-0.08, rh, fill=bg,
             line_color=rgb(0xDD,0xDD,0xDD), line_pt=0.3)
    add_text(s1, label, LX+0.14, y+0.03, LW-0.22, 0.22,
             size=9, bold=True, color=DARK)
    add_text(s1, desc,  LX+0.14, y+0.24, LW-0.22, 0.20,
             size=7.5, italic=True, color=MGRAY)
    # Version columns
    for xi, val, hc in [(V1, a, PURPLE), (V2, b, DORANGE), (V3, c, ORANGE)]:
        is_na = "—" in val
        add_rect(s1, xi, y, VW, rh, fill=bg,
                 line_color=rgb(0xDD,0xDD,0xDD), line_pt=0.3)
        add_text(s1, val, xi+0.1, y+0.06, VW-0.18, rh-0.1,
                 size=10, bold=not is_na, italic=is_na,
                 color=MGRAY if is_na else hc,
                 align=PP_ALIGN.CENTER)
    y += rh

# Total row
add_rect(s1, LX, y, 0.08, 0.46, fill=DARK)
add_rect(s1, LX+0.08, y, LW-0.08, 0.46, fill=DARK)
add_text(s1, "TOTAL / MONTH", LX+0.14, y+0.10, LW-0.22, 0.28,
         size=10, bold=True, color=WHITE)
for xi, val, hc, bc in [
    (V1, "$2,350 – $2,600", PURPLE,  LPURPLE),
    (V2, "$2,150 – $2,350", DORANGE, LORANGE),
    (V3, "$2,100 – $2,300", ORANGE,  rgb(0xFF,0xF3,0xE0)),
]:
    add_rect(s1, xi, y, VW, 0.46, fill=bc, line_color=hc, line_pt=2)
    add_text(s1, val, xi+0.1, y+0.08, VW-0.18, 0.30, size=13, bold=True,
             color=hc, align=PP_ALIGN.CENTER)

# Insight box
y2 = y + 0.54
add_rect(s1, LX, y2, 12.73, 0.44, fill=LORANGE, line_color=ORANGE, line_pt=1)
add_text(s1,
    "Key insight:  Version B (AWS+Qwen) and C (Alibaba) are both cheaper than A due to elasticity.  "
    "SageMaker scale-to-0 (re:Invent 2024) cuts student endpoint cost 50% (360 active hr/mo).  "
    "PAI-EAS min=1 auto-scale saves 15%. Both B and C are 8-11% cheaper than Version A overall.",
    LX+0.15, y2+0.06, 12.4, 0.32, size=9, color=DORANGE)

# ── Shared table builder ──────────────────────────────────────────────────────
def cost_table(slide, x, y, w, rows, hdr_color, hdr_text, row_h=0.38):
    """Draw a cost breakdown table.
    Each row: (component, service, volume, unit_price, calculation, monthly_total)
    """
    add_rect(slide, x, y, w, 0.36, fill=hdr_color)
    add_text(slide, hdr_text, x+0.12, y+0.05, w-0.2, 0.26,
             size=10, bold=True, color=WHITE)
    y += 0.36
    # 6 columns: Component 16% | Service 22% | Volume 13% | Unit Price 13% | Calculation 24% | $/ mo 12%
    col_pct = [0.16, 0.22, 0.13, 0.13, 0.24, 0.12]
    col_w = [w * p for p in col_pct]
    sub_hdrs = ["Component", "Service / Model", "Volume", "Unit Price", "Calculation", "$/mo"]
    for ci, (sh, cw) in enumerate(zip(sub_hdrs, col_w)):
        cx = x + sum(col_w[:ci])
        add_rect(slide, cx, y, cw, 0.26, fill=DORANGE)
        add_text(slide, sh, cx+0.04, y+0.03, cw-0.06, 0.20,
                 size=7, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT if ci == 5 else PP_ALIGN.LEFT)
    y += 0.26
    for i, row in enumerate(rows):
        comp, svc, vol, price, calc, total = row
        bg = LGRAY if i % 2 == 0 else WHITE
        is_saving = str(total).startswith("-") or "\u2212" in str(total)
        is_na = str(total) in ("—", "\u2014")
        vals = [comp, svc, vol, price, calc, total]
        for ci, (val, cw) in enumerate(zip(vals, col_w)):
            cx = x + sum(col_w[:ci])
            cell_bg = LGREEN if is_saving else (rgb(0xF5,0xF5,0xF5) if is_na else bg)
            add_rect(slide, cx, y, cw, row_h, fill=cell_bg,
                     line_color=rgb(0xDD,0xDD,0xDD), line_pt=0.3)
            is_total_col = (ci == 5)
            tc = GREEN if (is_saving and is_total_col) else (MGRAY if is_na else DARK)
            add_text(slide, str(val), cx+0.04, y+0.03, cw-0.06, row_h-0.05,
                     size=7 if not is_total_col else 9,
                     bold=is_total_col and not is_na,
                     italic=is_na,
                     color=tc,
                     align=PP_ALIGN.RIGHT if is_total_col else PP_ALIGN.LEFT)
        y += row_h
    return y

def subtotal_row(slide, x, y, w, label, amount, hdr_color):
    add_rect(slide, x, y, w, 0.34, fill=hdr_color)
    add_text(slide, label, x+0.12, y+0.06, w*0.65, 0.24, size=9, bold=True, color=WHITE)
    add_text(slide, amount, x+w*0.65, y+0.04, w*0.33, 0.28, size=12, bold=True,
             color=YELLOW, align=PP_ALIGN.RIGHT)
    return y + 0.34

def grand_total(slide, x, y, w, label, breakdown, total, hdr_color):
    add_rect(slide, x, y, w, 0.52, fill=hdr_color)
    add_text(slide, label, x+0.15, y+0.06, w*0.25, 0.40, size=13, bold=True, color=WHITE)
    add_text(slide, breakdown, x+w*0.27, y+0.10, w*0.45, 0.32, size=9, color=WHITE, italic=True)
    add_text(slide, total, x+w*0.73, y+0.06, w*0.25, 0.40, size=14, bold=True,
             color=YELLOW, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — VERSION A: AWS + CLAUDE
# ═══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
slide_title(s2, "Version A — AWS + Claude  |  Cost Breakdown",
            "Amazon Bedrock · ap-southeast-1 Singapore · 150,000 queries/month · No fine-tuning")
assumption_strip(s2)
footer_bar(s2, "Version A total: ~$2,480 – $2,700 / month")

# Left: LLM
llm_a = [
    ("Emergency lane",    "Claude Haiku 4.5 (Bedrock)",    "30k calls",  "$1/$5 per 1M",    "30k×1.5k×$1 + 30k×0.3k×$5",   "$90"),
    ("Router",            "Nova Micro (Bedrock)",           "120k calls", "$0.035/$0.14/1M", "120k×0.5k×$0.035 + 120k×0.04k×$0.14", "$3"),
    ("Specialist agents", "Claude Sonnet 4.5 (Bedrock)",   "120k calls", "$3/$15 per 1M",   "120k×3k×$3 + 120k×0.6k×$15",  "$2,160"),
    ("Embeddings",        "Cohere Embed Multilingual v3",  "150k calls", "$0.10 per 1M",    "150k×0.08k×$0.10",             "$1"),
    ("Guardrails",        "Bedrock Guardrails",             "120k calls", "$0.15/1k units",  "120k×3.6 units×$0.15/1k",      "$65"),
    ("Prompt cache",      "Auto (Claude 4.x, ~60% hit)",   "—",          "−90% cached tok", "~40% reduction on Sonnet cost", "−$864"),
]
y = cost_table(s2, 0.25, 1.45, 6.3, llm_a, PURPLE, "LLM INFERENCE  (Bedrock on-demand, Singapore)")
y = subtotal_row(s2, 0.25, y+0.04, 6.3, "LLM Subtotal (after prompt cache savings)", "~$1,455 / mo", PURPLE)

# Left: Fine-tuning
add_rect(s2, 0.25, y+0.08, 6.3, 0.30, fill=MGRAY)
add_text(s2, "FINE-TUNING  —  Not applicable in Version A (base Claude only)",
         0.37, y+0.10, 6.1, 0.24, size=9, italic=True, color=WHITE)

# Right: Infrastructure
infra_a = [
    ("Vector KB",       "Bedrock KB OpenSearch Serverless", "2 OCU × 720hr", "$0.24/OCU-hr", "2×720×$0.24",    "$346"),
    ("GraphRAG KB",     "Neptune Analytics",                "1 m-NCU × 720hr","$0.16/m-NCU-hr","1×720×$0.16",  "$115"),
    ("Semantic cache",  "ElastiCache Redis OSS r7g.large",  "720 hr",        "$0.166/hr",    "720×$0.166",     "$120"),
    ("API server",      "EC2 t4g.medium × 2 (HA)",          "720hr × 2",     "$0.046/hr",    "2×720×$0.046",   "$66"),
    ("API Gateway",     "AWS API Gateway REST",              "150k calls",    "$3.50/1M",     "0.15×$3.50",     "$1"),
    ("CDN",             "CloudFront",                        "~50GB transfer","$0.12/GB",     "50×$0.12",       "$6"),
    ("Storage",         "S3 Standard (corpus + logs)",       "~100GB",        "$0.023/GB-mo", "100×$0.023",     "$2"),
    ("Audit / Logs",    "CloudWatch Logs + S3 archive",      "~10GB/mo",      "$0.50/GB",     "10×$0.50",       "$5"),
    ("Security",        "AWS WAF + Bedrock Guardrails",      "150k req",      "$0.60/1M",     "0.15×$0.60",     "$1"),
    ("VPN Gateway",     "AWS VPN (1 IPsec tunnel)",          "720 hr",        "$0.05/hr",     "720×$0.05",      "$36"),
]
y2 = cost_table(s2, 6.78, 1.45, 6.3, infra_a, DORANGE, "INFRASTRUCTURE & SERVICES  (Singapore)")
y2 = subtotal_row(s2, 6.78, y2+0.04, 6.3, "Infrastructure Subtotal", "~$696 / mo", DORANGE)

# Grand total
yg = max(y + 0.52, y2 + 0.12)
grand_total(s2, 0.25, yg, 12.83,
            "VERSION A TOTAL",
            "LLM $1,455  +  Infra $696  +  10–20% buffer",
            "~$2,350 – $2,600 / month", PURPLE)
# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VERSION B: AWS + QWEN
# ═══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
slide_title(s3, "Version B — AWS + Qwen  |  Cost Breakdown",
            "Amazon Bedrock Nova + Qwen3-4B SFT student · ap-southeast-1 Singapore · Monthly fine-tuning included")
assumption_strip(s3)
footer_bar(s3, "Version B total: ~$3,120 – $3,680 / month")

llm_b = [
    ("Emergency lane",    "Nova Lite (Bedrock SG)",          "30k calls",  "$0.06/$0.24/1M", "30k×1.5k×$0.06 + 30k×0.3k×$0.24", "$5"),
    ("Router",            "Nova Micro (Bedrock SG)",         "120k calls", "$0.035/$0.14/1M","120k×0.5k×$0.035 + 120k×0.04k×$0.14","$3"),
    ("Specialist agents", "Nova Pro (Bedrock SG)",           "120k calls", "$0.80/$3.20/1M", "120k×3k×$0.80 + 120k×0.6k×$3.20", "$518"),
    ("Student model",     "Qwen3-4B SFT+LoRA\nSageMaker g5.xlarge\nscale-to-0 (re:Invent 2024)",
     "360hr/mo\n(12hr/day active)", "$1.97/hr",
     "360hr × $1.97\n(50% vs always-on $1,419)", "$710"),
    ("Embeddings",        "Cohere Embed Multilingual v3",    "150k calls", "$0.10 per 1M",   "150k×0.08k×$0.10",                 "$1"),
    ("Guardrails",        "Bedrock Guardrails",               "48k calls",  "$0.15/1k units", "48k×3.6×$0.15/1k",                 "$26"),
]
y = cost_table(s3, 0.25, 1.45, 6.3, llm_b, PURPLE, "LLM + STUDENT ENDPOINT  (all Singapore)")
y = subtotal_row(s3, 0.25, y+0.04, 6.3, "LLM + Student Endpoint Subtotal", "~$1,263 / mo", PURPLE)

# Fine-tuning
ft_b = [
    ("Data generation", "Qwen3.5-397B via DeepInfra (HF)", "10k Q&A/mo", "$0.54/$3.40/1M", "10k×0.5k×$0.54 + 10k×0.3k×$3.40", "$4"),
    ("Training job",    "SageMaker ml.g4dn.2xlarge",        "~3 hr/run",  "$1.32/hr",       "3×$1.32",                           "$4"),
    ("Model storage",   "S3 checkpoint (~8GB)",              "720 hr",     "$0.023/GB-mo",   "8×$0.023",                          "$0"),
]
y = cost_table(s3, 0.25, y+0.08, 6.3, ft_b, DORANGE, "MONTHLY FINE-TUNING  (SageMaker + DeepInfra)", row_h=0.36)
y = subtotal_row(s3, 0.25, y+0.04, 6.3, "Fine-tuning Subtotal", "~$17 / mo", DORANGE)

infra_b = [
    ("Vector KB",       "Bedrock KB OpenSearch Serverless", "2 OCU × 720hr", "$0.24/OCU-hr", "2×720×$0.24",    "$346"),
    ("GraphRAG KB",     "Neptune Analytics",                "1 m-NCU × 720hr","$0.16/m-NCU-hr","1×720×$0.16",  "$115"),
    ("Semantic cache",  "ElastiCache Redis OSS r7g.large",  "720 hr",        "$0.166/hr",    "720×$0.166",     "$120"),
    ("API server",      "EC2 t4g.medium × 2 (HA)",          "720hr × 2",     "$0.046/hr",    "2×720×$0.046",   "$66"),
    ("API Gateway",     "AWS API Gateway REST",              "150k calls",    "$3.50/1M",     "0.15×$3.50",     "$1"),
    ("CDN",             "CloudFront",                        "~50GB",         "$0.12/GB",     "50×$0.12",       "$6"),
    ("Storage",         "S3 Standard",                       "~100GB",        "$0.023/GB-mo", "100×$0.023",     "$2"),
    ("Audit / Logs",    "CloudWatch Logs + S3",              "~10GB/mo",      "$0.50/GB",     "10×$0.50",       "$5"),
    ("Security",        "AWS WAF + Guardrails",              "150k req",      "$0.60/1M",     "0.15×$0.60",     "$1"),
    ("VPN Gateway",     "AWS VPN (1 tunnel)",                "720 hr",        "$0.05/hr",     "720×$0.05",      "$36"),
    ("Scale-to-zero",   "SageMaker Serverless (alt.)",       "72k × 3s × 6GB","$0.00002/GB-s","72k×3×6×$0.00002","$26"),
]
y2 = cost_table(s3, 6.78, 1.45, 6.3, infra_b, ORANGE, "INFRASTRUCTURE & SERVICES  (Singapore)")
y2 = subtotal_row(s3, 6.78, y2+0.04, 6.3, "Infrastructure Subtotal", "~$696 / mo", ORANGE)

yg = max(y + 0.12, y2 + 0.12)
grand_total(s3, 0.25, yg, 12.83,
            "VERSION B TOTAL",
            "LLM+Student $1,263  +  Fine-tune $17  +  Infra $696  +  10–20% buffer",
            "~$2,150 – $2,350 / month", DORANGE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — VERSION C: ALIBABA CLOUD
# ═══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
slide_title(s4, "Version C — Alibaba Cloud  |  Cost Breakdown",
            "Model Studio (DashScope Intl) + PAI-EAS · ap-southeast-1 Singapore International · Monthly fine-tuning included")
assumption_strip(s4)
footer_bar(s4, "Version C total: ~$2,420 – $3,060 / month")

llm_c = [
    ("Emergency lane",    "Qwen3.5-Flash (Model Studio)",   "30k calls",  "$0.05/$0.20/1M",  "30k×1.5k×$0.05 + 30k×0.3k×$0.20",  "$4"),
    ("Router",            "Qwen3.5-Flash JSON mode",         "120k calls", "$0.05/$0.20/1M",  "120k×0.5k×$0.05 + 120k×0.04k×$0.20","$4"),
    ("Specialist agents", "Qwen3.5-Plus (Model Studio)",    "120k calls", "$0.115/$0.688/1M","120k×3k×$0.115 + 120k×0.6k×$0.688", "$91"),
    ("Student model",     "Qwen3-8B SFT+LoRA\nPAI-EAS A10 GPU\nmin=1 replica, auto-scale",
     "720hr/mo\n(min=1, 85% util)", "$1.50/hr avg",
     "720hr × $1.50 × 0.85\n(15% saving vs always-on $1,080)", "$918"),
    ("Embeddings",        "text-embedding-v4",               "150k calls", "$0.07 per 1M",    "150k×0.08k×$0.07",                   "$1"),
    ("Re-rank",           "qwen3-rerank",                    "120k calls", "$0.10 per 1M",    "120k×3.6k×$0.10",                    "$43"),
    ("Context cache",     "Qwen Context Cache (L2, ~60%)",   "90k hits",   "−90% cached tok", "~40% reduction on Plus cost",        "−$36"),
]
y = cost_table(s4, 0.25, 1.45, 6.3, llm_c, PURPLE, "LLM + PAI-EAS STUDENT  (Model Studio DashScope Intl)")
y = subtotal_row(s4, 0.25, y+0.04, 6.3, "LLM + PAI-EAS Subtotal", "~$1,025 / mo", PURPLE)

ft_c = [
    ("Data generation", "Qwen3.5-Plus as teacher (Model Studio)", "10k Q&A/mo", "$0.115/$0.688/1M","10k×0.5k×$0.115 + 10k×0.3k×$0.688","$8"),
    ("Training job",    "PAI-DLC A10 GPU, Qwen3-8B SFT+LoRA",    "~4 hr/run",  "$1.50/hr",        "4×$1.50",                            "$6"),
    ("Model storage",   "OSS checkpoint (~16GB)",                  "720 hr",     "$0.02/GB-mo",     "16×$0.02",                           "$0"),
]
y = cost_table(s4, 0.25, y+0.08, 6.3, ft_c, DORANGE, "MONTHLY FINE-TUNING  (PAI-DLC + Model Studio)", row_h=0.36)
y = subtotal_row(s4, 0.25, y+0.04, 6.3, "Fine-tuning Subtotal", "~$9 / mo", DORANGE)

infra_c = [
    ("Vector Store",    "OpenSearch Vector Search HA (2 OCU)", "2 OCU × 720hr","$0.24/OCU-hr",    "2×720×$0.24",          "$346"),
    ("GraphRAG",        "AnalyticDB for PostgreSQL 4c/32GB",   "720 hr",       "$0.42/hr",         "720×$0.42",            "$302"),
    ("Semantic cache",  "Tair (Redis-compatible, Multi-AZ)",   "720 hr",       "$0.083/hr",        "720×$0.083",           "$60"),
    ("Compute",         "Function Compute 3.0 (serverless)",   "150k×2s×512MB","$0.0000167/GB-s",  "150k×2×0.5×$0.0000167","$26"),
    ("API Gateway",     "Alibaba API Gateway",                  "150k calls",   "$3.50/1M",         "0.15×$3.50",           "$1"),
    ("CDN / Edge",      "Alibaba CDN + Anti-DDoS",              "~50GB",        "$0.10/GB",         "50×$0.10",             "$5"),
    ("Storage",         "OSS Standard + WORM audit",            "~200GB",       "$0.02/GB-mo",      "200×$0.02",            "$4"),
    ("Audit / Logs",    "SLS + OSS WORM (6yr retention)",       "~10GB/mo",     "$0.35/GB",         "10×$0.35",             "$4"),
    ("Security",        "WAF + Content Mod 2.0 + SDDP",         "150k req",     "~$0.50/1M",        "0.15×$0.50",           "$1"),
    ("VPN Gateway",     "Alibaba VPN Gateway 100Mbps",          "720 hr",       "$0.15/hr",         "720×$0.15",            "$108"),
    ("IdaaS / SSO",     "IDaaS EIAM 2.0 Premium+",              "500 users",    "$0.10/user/mo",    "500×$0.10",            "$50"),
]
y2 = cost_table(s4, 6.78, 1.45, 6.3, infra_c, ORANGE, "INFRASTRUCTURE & SERVICES  (Singapore International)")
y2 = subtotal_row(s4, 6.78, y2+0.04, 6.3, "Infrastructure Subtotal", "~$881 / mo", ORANGE)

yg = max(y + 0.12, y2 + 0.12)
grand_total(s4, 0.25, yg, 12.83,
            "VERSION C TOTAL",
            "LLM+PAI-EAS $1,025  +  Fine-tune $9  +  Infra $881  +  10–20% buffer",
            "~$2,100 – $2,300 / month", ORANGE)

# ── Save ──────────────────────────────────────────────────────────────────────
# Save with all slides (template slides + our 4 new ones)
# The 4 new slides are at the end
prs.save(str(OUT))
print(f"Saved: {OUT}")
total_slides = len(prs.slides)
print(f"Total slides: {total_slides} (first {total_slides-4} are template, last 4 are cost slides)")
print("4 cost slides: Summary + Version A + Version B + Version C")
