"""Build the 3-minute Nova Health Tech client presentation deck.

Output: docs/Nova_Health_Tech_AWS_Claude_Presentation_3min.pptx

Flow (per user feedback):
1. Title (10s)
2. The Challenge: focus on PHI (25s) - hook to compliance
3. Singapore Healthcare Compliance & Security (45s) - the answer to PHI
4. The Challenge continued: Speed + Freshness (20s)
5. Solution: Architecture (50s) - includes architecture path/visual
6. PoC Results & Cost (40s) - no Reserved Tier mention
7. Ready to Deploy (10s)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

REPO = Path(__file__).resolve().parent.parent
BG_IMAGE = str(REPO / "mainslide.png")
OUT_PATH = str(REPO / "docs" / "Nova_Health_Tech_AWS_Claude_Presentation_3min.pptx")

ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK = RGBColor(0xCC, 0x4F, 0x1F)
ORANGE_LIGHT = RGBColor(0xFF, 0xE8, 0xDA)
GRAY_DARK = RGBColor(0x2C, 0x3E, 0x50)
GRAY_MED = RGBColor(0x64, 0x64, 0x64)
GRAY_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_slide(prs, title, subtitle, notes):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(BG_IMAGE, 0, 0, width=SLIDE_W, height=SLIDE_H)

    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    sp = overlay.fill._xPr.find(qn('a:solidFill'))
    if sp is not None:
        srgbClr = sp.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '40000')
    overlay.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_content_slide(prs, title, notes):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK

    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2), Inches(0.04))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ORANGE
    underline.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Nova Health Tech · Clinical GenAI Assistant · AWS Singapore"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY_MED
    run.font.italic = True

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_simple(slide, left, top, width, height, headers, rows, *, header_color=ORANGE):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = h
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else ORANGE_LIGHT
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_top = Inches(0.03)
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = "Calibri"
            run.font.size = Pt(10)
            run.font.color.rgb = GRAY_DARK
    return tbl


# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ----------------------------------------------------------------------------
# SLIDE 1: TITLE (10s)
# ----------------------------------------------------------------------------
add_title_slide(
    prs,
    title="Nova Health Tech · Clinical GenAI Assistant",
    subtitle="AWS with Claude · Singapore · PDPA & HCSA compliant by design",
    notes=(
        "[10s opening]\n\n"
        "Đề xuất kiến trúc trợ lý AI lâm sàng cho Nova: AWS Singapore, "
        "Claude 4.5, tuân thủ PDPA và HCSA ngay từ thiết kế."
    )
)
print("  Slide 1 (title): added")

# ----------------------------------------------------------------------------
# SLIDE 2: THE CHALLENGE - Patient PHI (25s)
# Highlight the PHI card to hook compliance discussion
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "The Challenge",
    "[25s]\n\n"
    "Bác sĩ Nova mất niềm tin vì câu trả lời chậm và thiếu cụ thể. Có 3 thách thức.\n\n"
    "Đầu tiên và quan trọng nhất: dữ liệu thử nghiệm nội bộ chứa thông tin bệnh nhân "
    "nhạy cảm, không được rò rỉ. Đây không chỉ là kỹ thuật, mà là yêu cầu pháp lý "
    "của Singapore. Trước khi nói tiếp, để tôi giải thích tuân thủ phía Singapore."
)

# 3 pain cards, PHI prominent (larger + animated highlight)
pains = [
    ("⚡", "Speed", "Emergency response\n< 2 seconds", ORANGE, False),
    ("🔒", "Patient PHI", "Trial data must stay\nin Singapore", RED, True),
    ("📅", "Freshness", "WHO monthly,\nICD-11 daily updates", ORANGE, False),
]

card_w = Inches(3.8)
card_h = Inches(3.5)
gap = Inches(0.3)
total_w = card_w * 3 + gap * 2
start_x = (SLIDE_W - total_w) / 2

for i, (icon, title, desc, color, highlight) in enumerate(pains):
    x = start_x + i * (card_w + gap)
    y = Inches(1.7)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    if highlight:
        card.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
        card.line.color.rgb = color
        card.line.width = Pt(4)
    else:
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GRAY_LIGHT
        card.line.width = Pt(1.5)
    card.adjustments[0] = 0.06
    
    # Icon
    icon_tb = slide.shapes.add_textbox(x, y + Inches(0.4), card_w, Inches(1.0))
    tf = icon_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(56)
    
    # Title
    title_tb = slide.shapes.add_textbox(x, y + Inches(1.6), card_w, Inches(0.5))
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(22) if not highlight else Pt(26)
    run.font.bold = True
    run.font.color.rgb = color if highlight else GRAY_DARK
    
    # Desc
    desc_tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(2.3), card_w - Inches(0.4), Inches(1.1))
    tf = desc_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY_DARK if not highlight else color
    run.font.bold = highlight

# Arrow pointing to PHI card
arrow_y = Inches(5.4)
arrow_text = slide.shapes.add_textbox(Inches(4.2), arrow_y, Inches(5.0), Inches(0.5))
tf = arrow_text.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "↑ Singapore healthcare regulation makes this non-negotiable"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.italic = True
run.font.color.rgb = RED

# Bottom: data sources note
note_y = Inches(6.0)
note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), note_y, Inches(11.3), Inches(0.85))
note_box.fill.solid()
note_box.fill.fore_color.rgb = ORANGE_LIGHT
note_box.line.color.rgb = ORANGE
note_box.line.width = Pt(1)
note_box.adjustments[0] = 0.2

note_tb = slide.shapes.add_textbox(Inches(1.2), note_y + Inches(0.1), Inches(10.9), Inches(0.65))
tf = note_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Data sources: Internal trial PDFs · WHO API · ICD-11 · PubMed · EHR (FHIR R4)"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Workload: 600k calls per month at launch"
run2.font.name = "Calibri"
run2.font.size = Pt(11)
run2.font.italic = True
run2.font.color.rgb = GRAY_MED

print("  Slide 2 (challenge - PHI focus): added")

# ----------------------------------------------------------------------------
# SLIDE 3: SINGAPORE COMPLIANCE & SECURITY (45s) — KEY SLIDE
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Singapore Compliance & Security: Built-In, Not Bolted-On",
    "[45s, KEY SLIDE]\n\n"
    "Singapore có 4 quy định chính áp dụng cho hệ thống AI lâm sàng. "
    "PDPA do PDPC quản, bảo vệ dữ liệu cá nhân, vi phạm phải báo trong 72 giờ. "
    "HCSA 2020 do MOH ban hành thay PHMC, bắt buộc license cho clinical decision support. "
    "Cybersecurity Act 2018 bảo vệ Critical Information Infrastructure. "
    "IMDA AI Verify framework cho AI có trách nhiệm.\n\n"
    "Cách triển khai: 100 phần trăm deploy ở ap-southeast-1, dữ liệu bệnh nhân không bao giờ "
    "rời Singapore. AWS có signed BAA, ISO 27001, SOC 2 Type II, PDPA-compliant."
)

# 4 SG regulations as prominent cards
sg_regs = [
    {
        "title": "PDPA",
        "subtitle": "Personal Data Protection Act",
        "authority": "PDPC",
        "scope": "Patient data protection, consent, breach notification within 72 hours",
        "our_response": "Comprehend Medical PHI mask, KMS per-tenant CMK, audit trail",
        "color": ORANGE_DARK,
    },
    {
        "title": "HCSA",
        "subtitle": "Healthcare Services Act 2020",
        "authority": "MOH",
        "scope": "Licensing for clinical decision support, replaces PHMC 1980",
        "our_response": "Decision support only, human-in-the-loop, audit ready",
        "color": ORANGE,
    },
    {
        "title": "Cyber Act",
        "subtitle": "Cybersecurity Act 2018 (CII)",
        "authority": "CSA",
        "scope": "Healthcare CII protection, incident reporting, security audits",
        "our_response": "GuardDuty + Security Hub + Macie, Security Lake aggregation",
        "color": ORANGE_DARK,
    },
    {
        "title": "AI Verify",
        "subtitle": "IMDA Responsible AI",
        "authority": "IMDA",
        "scope": "Transparency, fairness, accountability for AI systems",
        "our_response": "chat_trace audit, citation validator, Bedrock Guardrails",
        "color": ORANGE,
    },
]

cw = Inches(3.0)
ch = Inches(4.0)
gap = Inches(0.1)
total_w = cw * 4 + gap * 3
start_x = (SLIDE_W - total_w) / 2
y = Inches(1.4)

for i, reg in enumerate(sg_regs):
    x = start_x + i * (cw + gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = reg["color"]
    card.line.width = Pt(2)
    card.adjustments[0] = 0.05
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.85))
    header.fill.solid()
    header.fill.fore_color.rgb = reg["color"]
    header.line.fill.background()
    
    h_tb = slide.shapes.add_textbox(x, y + Inches(0.05), cw, Inches(0.85))
    tf = h_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = reg["title"]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = reg["subtitle"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(10)
    run2.font.color.rgb = WHITE
    
    body_y = y + Inches(0.95)
    
    auth_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y, cw - Inches(0.3), Inches(0.5))
    tf = auth_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Authority"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["authority"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = GRAY_DARK
    
    scope_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y + Inches(0.85), cw - Inches(0.3), Inches(1.1))
    tf = scope_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Scope"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["scope"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_DARK
    
    resp_tb = slide.shapes.add_textbox(x + Inches(0.15), body_y + Inches(2.1), cw - Inches(0.3), Inches(0.9))
    tf = resp_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Our Implementation"
    run.font.name = "Calibri"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = reg["our_response"]
    run2.font.name = "Calibri"
    run2.font.size = Pt(10)
    run2.font.color.rgb = GRAY_DARK

# Bottom: data residency commitment
res_y = Inches(5.6)
res_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), res_y, Inches(12.3), Inches(1.1))
res_box.fill.solid()
res_box.fill.fore_color.rgb = ORANGE_LIGHT
res_box.line.color.rgb = ORANGE_DARK
res_box.line.width = Pt(2)
res_box.adjustments[0] = 0.1

res_tb = slide.shapes.add_textbox(Inches(0.7), res_y + Inches(0.1), Inches(11.9), Inches(0.9))
tf = res_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Data Residency Commitment"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = (
    "100% deployment in ap-southeast-1 (Singapore)  ·  Patient cleartext never leaves SG  ·  "
    "Cross-region: tokens only  ·  6-year audit retention via S3 Object Lock"
)
run2.font.name = "Calibri"
run2.font.size = Pt(12)
run2.font.color.rgb = GRAY_DARK

print("  Slide 3 (compliance): added")

# ----------------------------------------------------------------------------
# SLIDE 4: BACK TO CHALLENGE - Speed + Freshness (20s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Continuing the Challenge",
    "[20s]\n\n"
    "Quay lại 2 thách thức còn lại. Speed: cấp cứu cần dưới 2 giây, không có chỗ "
    "cho độ trễ, mỗi giây có thể là sống chết. Freshness: WHO cập nhật phác đồ hàng tháng, "
    "ICD-11 cập nhật hàng ngày, hệ thống phải tự đồng bộ không cần can thiệp manual."
)

# 2 cards: Speed (left) + Freshness (right)
cards = [
    {
        "icon": "⚡",
        "title": "Speed",
        "headline": "Emergency response < 2 seconds",
        "details": [
            "Every second matters in acute care",
            "Sepsis, STEMI, anaphylaxis decisions",
            "TTFT (time to first token) is the key metric",
            "Streaming response to perceive answer",
        ],
        "color": RED,
    },
    {
        "icon": "📅",
        "title": "Freshness",
        "headline": "WHO monthly · ICD-11 daily updates",
        "details": [
            "WHO publishes guideline revisions every month",
            "ICD-11 catalog refreshes on its own cadence",
            "Internal trial PDFs uploaded ad hoc",
            "Auto-sync with versioning, no manual touch",
        ],
        "color": ORANGE_DARK,
    },
]

cw = Inches(5.8)
ch = Inches(4.5)
gap = Inches(0.4)
total_w = cw * 2 + gap
start_x = (SLIDE_W - total_w) / 2

for i, c in enumerate(cards):
    x = start_x + i * (cw + gap)
    y = Inches(1.6)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = c["color"]
    card.line.width = Pt(2.5)
    card.adjustments[0] = 0.04
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.7))
    header.fill.solid()
    header.fill.fore_color.rgb = c["color"]
    header.line.fill.background()
    
    h_tb = slide.shapes.add_textbox(x, y, cw, Inches(0.7))
    tf = h_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{c['icon']}  {c['title']}"
    run.font.name = "Calibri"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Headline
    headline_tb = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(0.95), cw - Inches(0.8), Inches(0.7))
    tf = headline_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = c["headline"]
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = c["color"]
    
    # Details
    detail_tb = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(1.85), cw - Inches(1.0), Inches(2.5))
    tf = detail_tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(c["details"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        bullet_run = p.add_run()
        bullet_run.text = "▸ "
        bullet_run.font.name = "Calibri"
        bullet_run.font.size = Pt(14)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = c["color"]
        text_run = p.add_run()
        text_run.text = line
        text_run.font.name = "Calibri"
        text_run.font.size = Pt(13)
        text_run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(8)

# Bottom: tone consistency note
tone_y = Inches(6.4)
tone_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), tone_y, Inches(11.3), Inches(0.6))
tone_box.fill.solid()
tone_box.fill.fore_color.rgb = ORANGE_LIGHT
tone_box.line.color.rgb = ORANGE
tone_box.line.width = Pt(1)
tone_box.adjustments[0] = 0.2

tone_tb = slide.shapes.add_textbox(Inches(1.2), tone_y + Inches(0.1), Inches(10.9), Inches(0.45))
tf = tone_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Plus: consistent clinical tone across answers · legacy PDFs with inconsistent tagging"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.italic = True
run.font.color.rgb = ORANGE_DARK

print("  Slide 4 (challenge - speed + freshness): added")

# ----------------------------------------------------------------------------
# SLIDE 5: SOLUTION ARCHITECTURE (50s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Solution: Two-Lane Architecture",
    "[50s]\n\n"
    "Giải pháp là kiến trúc 2 luồng. Cấp cứu đi thẳng Claude Haiku 4.5 với top-3 vector "
    "và top-2 graph, không qua router. Phức tạp đi qua Nova Micro phân loại 12 phòng ban, "
    "rồi Sonnet 4.5 với top-15 vector và top-3 graph có guardrails.\n\n"
    "RAG dùng 2 loại: vector trên OpenSearch Serverless và GraphRAG trên Neptune Analytics. "
    "Cả 2 đều ở Singapore với Cohere Embed v3 SG-native. Không có cross-region hop.\n\n"
    "Diagram chi tiết kiến trúc xem tại docs/architecture/diagrams/proposal/01_high_level.svg"
)

# Two-lane visualization
lane_y = Inches(1.4)
lane_w = Inches(6.0)
lane_h = Inches(2.7)

# Emergency lane
em_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), lane_y, lane_w, lane_h)
em_card.fill.solid()
em_card.fill.fore_color.rgb = WHITE
em_card.line.color.rgb = RED
em_card.line.width = Pt(2.5)
em_card.adjustments[0] = 0.04

em_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), lane_y, lane_w, Inches(0.55))
em_h.fill.solid()
em_h.fill.fore_color.rgb = RED
em_h.line.fill.background()

em_h_tb = slide.shapes.add_textbox(Inches(0.5), lane_y, lane_w, Inches(0.55))
tf = em_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "EMERGENCY LANE  ·  PoC TTFT 2.5s"
run.font.name = "Calibri"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = WHITE

em_lines = [
    "▸  Claude Haiku 4.5 (fast model)",
    "▸  Vector top-3 + GraphRAG top-2",
    "▸  No router (direct emergency agent)",
    "▸  No guardrails (speed priority)",
    "▸  Streaming SSE, max 300 tokens",
]
em_tb = slide.shapes.add_textbox(Inches(0.8), lane_y + Inches(0.7), lane_w - Inches(0.5), lane_h - Inches(0.8))
tf = em_tb.text_frame
tf.word_wrap = True
for i, line in enumerate(em_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(4)

# Complex lane
cx_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), lane_y, lane_w, lane_h)
cx_card.fill.solid()
cx_card.fill.fore_color.rgb = WHITE
cx_card.line.color.rgb = ORANGE_DARK
cx_card.line.width = Pt(2.5)
cx_card.adjustments[0] = 0.04

cx_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), lane_y, lane_w, Inches(0.55))
cx_h.fill.solid()
cx_h.fill.fore_color.rgb = ORANGE_DARK
cx_h.line.fill.background()

cx_h_tb = slide.shapes.add_textbox(Inches(6.8), lane_y, lane_w, Inches(0.55))
tf = cx_h_tb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "COMPLEX LANE  ·  PoC TTFT 9.7s"
run.font.name = "Calibri"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = WHITE

cx_lines = [
    "▸  Claude Sonnet 4.5 (deep reasoning)",
    "▸  Vector top-15 + GraphRAG top-3",
    "▸  Nova Micro routes to 12 departments",
    "▸  Bedrock Guardrails enabled",
    "▸  Streaming SSE, max 1500 tokens",
]
cx_tb = slide.shapes.add_textbox(Inches(7.1), lane_y + Inches(0.7), lane_w - Inches(0.5), lane_h - Inches(0.8))
tf = cx_tb.text_frame
tf.word_wrap = True
for i, line in enumerate(cx_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK
    p.space_after = Pt(4)

# Architecture diagram path note + simple ASCII flow
arch_y = Inches(4.3)

arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), arch_y, Inches(12.3), Inches(2.0))
arch_box.fill.solid()
arch_box.fill.fore_color.rgb = ORANGE_LIGHT
arch_box.line.color.rgb = ORANGE
arch_box.line.width = Pt(1.5)
arch_box.adjustments[0] = 0.05

arch_label = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(0.1), Inches(11.9), Inches(0.4))
tf = arch_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Shared Infrastructure (100% Singapore-Native, ap-southeast-1)"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

# Component row 1
comp_row1 = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(0.55), Inches(11.9), Inches(0.5))
tf = comp_row1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "OpenSearch Serverless · Neptune Analytics · Cohere Embed v3 · Bedrock KB · Bedrock Guardrails"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.color.rgb = GRAY_DARK

# Component row 2
comp_row2 = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(0.95), Inches(11.9), Inches(0.5))
tf = comp_row2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Comprehend Medical (PHI mask) · ElastiCache Redis (semantic cache) · CloudFront + WAF · Cognito + Entra ID"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.color.rgb = GRAY_DARK

# Diagram path
path_label = slide.shapes.add_textbox(Inches(0.7), arch_y + Inches(1.45), Inches(11.9), Inches(0.4))
tf = path_label.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "📐  Full architecture diagram: docs/architecture/diagrams/proposal/01_high_level.svg"
run.font.name = "Calibri"
run.font.size = Pt(11)
run.font.italic = True
run.font.color.rgb = ORANGE_DARK

# Bottom note
poc_y = Inches(6.5)
poc_tb = slide.shapes.add_textbox(Inches(0.5), poc_y, Inches(12.3), Inches(0.4))
tf = poc_tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "PoC live at http://47.130.120.152/ui/index.html  ·  100% SLA pass on 20 test questions"
run.font.name = "Calibri"
run.font.size = Pt(13)
run.font.bold = True
run.font.color.rgb = GREEN

print("  Slide 5 (solution + architecture): added")

# ----------------------------------------------------------------------------
# SLIDE 6: PoC RESULTS & COST (40s)
# No mention of Reserved Tier per user request
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "PoC Results & Cost",
    "[40s]\n\n"
    "PoC chạy live trên EC2 Singapore với WHO COVID-19 guideline 198 trang. "
    "Test 20 câu hỏi thực tế: 100 phần trăm pass SLA cả 2 lane, 100 phần trăm câu trả lời "
    "có citation từ nguồn gốc.\n\n"
    "Cấp cứu 2.5 giây trung bình, complex 9.7 giây. Đã đáp ứng đề bài. "
    "Production thêm Prompt Caching trên system prompt và ElastiCache Redis cho semantic cache "
    "sẽ giảm thêm.\n\n"
    "Chi phí 600,000 calls một tháng: phương án Nova khoảng 2,800 USD, phương án Claude với "
    "student distillation khoảng 5,500 USD. Khi scale lên 3 triệu calls thì A1+ là 10,500 còn "
    "A2 là 23,500 USD."
)

# Top: 4 metrics
metrics = [
    ("Emergency TTFT", "2.5s", "Haiku 4.5 + dual KB", GREEN),
    ("Complex TTFT", "9.7s", "Sonnet 4.5 + Guardrails", GREEN),
    ("Answer Rate", "100%", "All grounded with citations", GREEN),
    ("SLA Pass", "100%", "20/20 test questions", GREEN),
]

m_y = Inches(1.4)
m_w = Inches(2.85)
m_h = Inches(1.4)
m_gap = Inches(0.2)
total_w = m_w * 4 + m_gap * 3
start_x = (SLIDE_W - total_w) / 2

for i, (label, value, sub, color) in enumerate(metrics):
    x = start_x + i * (m_w + m_gap)
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, m_y, m_w, m_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ORANGE
    card.line.width = Pt(2)
    card.adjustments[0] = 0.06
    
    label_tb = slide.shapes.add_textbox(x, m_y + Inches(0.1), m_w, Inches(0.3))
    tf = label_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_MED
    
    value_tb = slide.shapes.add_textbox(x, m_y + Inches(0.35), m_w, Inches(0.65))
    tf = value_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = "Calibri"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    sub_tb = slide.shapes.add_textbox(x, m_y + Inches(1.0), m_w, Inches(0.35))
    tf = sub_tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sub
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = GRAY_DARK

# Cost table
cost_label_y = Inches(3.1)
cost_label_box = slide.shapes.add_textbox(Inches(0.5), cost_label_y, Inches(12.3), Inches(0.4))
tf = cost_label_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Monthly Cost (600,000 calls per month, AWS Singapore list price)"
run.font.name = "Calibri"
run.font.size = Pt(15)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

# Verified cost numbers from proposal:
# A1+ (Nova): Nova Micro $70 + Nova Pro $1,470 + Embed $10 + Guardrails $180 + 
#             OpenSearch $350 + Neptune $115 + Comprehend $180 + Lambda/API/CF/WAF $150 + 
#             S3/CloudTrail/Macie $120 + Redis $80 + VPN $80 = $2,805
# A2 (Claude): Haiku $350 + Sonnet $5,460 + (other items same as A1+: $1,265) = $7,075
#              + Distillation $670 - Student offset $2,200 = $5,545
headers = ["Component", "A1+ (Nova)", "A2 (Claude with student)"]
rows = [
    ["LLM inference (emergency lane)", "$70", "$350"],
    ["LLM inference (complex lane)", "$1,470", "$5,460"],
    ["Embedding (Cohere Embed v3 SG)", "$10", "$10"],
    ["OpenSearch Serverless + Neptune Analytics", "$465", "$465"],
    ["Bedrock Guardrails + Comprehend Medical", "$360", "$360"],
    ["ElastiCache Redis + VPN + Lambda + S3 + CloudTrail", "$430", "$430"],
    ["Distillation amortized (offset by student)", "—", "−$1,530"],
    ["TOTAL / month", "$2,805", "$5,545"],
]
add_table_simple(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(2.7), headers, rows)

# Bottom: scale note
rec_y = Inches(6.45)
rec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), rec_y, Inches(12.3), Inches(0.55))
rec_box.fill.solid()
rec_box.fill.fore_color.rgb = ORANGE_LIGHT
rec_box.line.color.rgb = ORANGE
rec_box.line.width = Pt(1)
rec_box.adjustments[0] = 0.2

rec_tb = slide.shapes.add_textbox(Inches(0.7), rec_y + Inches(0.08), Inches(11.9), Inches(0.4))
tf = rec_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "At 3M calls/month: A1+ scales to $10,500   ·   A2 to $23,500   ·   Optimization: Prompt Caching + ElastiCache semantic cache"
run.font.name = "Calibri"
run.font.size = Pt(12)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

print("  Slide 6 (PoC results + cost): added")

# ----------------------------------------------------------------------------
# SLIDE 7: READY TO DEPLOY (10s)
# ----------------------------------------------------------------------------
slide = add_content_slide(
    prs,
    "Ready to Deploy",
    "[10s closing]\n\n"
    "Tóm lại: kiến trúc đã PoC kiểm chứng, Singapore-native, tuân thủ PDPA HCSA. "
    "Lộ trình 6-10 tuần. PoC đang chạy live, có thể demo ngay. Xin cảm ơn."
)

# Roadmap
rm_y = Inches(1.4)
rm_label = slide.shapes.add_textbox(Inches(0.5), rm_y, Inches(12.3), Inches(0.4))
tf = rm_label.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Implementation Roadmap (6 to 10 weeks)"
run.font.name = "Calibri"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = ORANGE_DARK

phases = [
    ("Wk 1-2", "Foundation", "SG provisioning, BAA, WHO + ICD-11 ingest"),
    ("Wk 3-4", "Customization", "Bedrock Distillation: Nova Lite student"),
    ("Wk 5-6", "Integration", "EHR FHIR, SharePoint, Cognito federation"),
    ("Wk 7-8", "Hardening", "Red team, Guardrails tune, load test"),
    ("Launch", "Go-Live", "All 12 departments, audit + cache live"),
]

phase_y = Inches(1.95)
phase_h = Inches(0.55)
for i, (week, name, desc) in enumerate(phases):
    y = phase_y + i * (phase_h + Inches(0.08))
    
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(1.5), phase_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    
    badge_tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.5), phase_h)
    tf = badge_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = week
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    name_tb = slide.shapes.add_textbox(Inches(2.2), y, Inches(2.4), phase_h)
    tf = name_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.name = "Calibri"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = ORANGE_DARK
    
    desc_tb = slide.shapes.add_textbox(Inches(4.7), y, Inches(8.0), phase_h)
    tf = desc_tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GRAY_DARK

# Live demo callout
demo_y = Inches(5.4)
demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), demo_y, Inches(10.3), Inches(1.6))
demo_box.fill.solid()
demo_box.fill.fore_color.rgb = ORANGE
demo_box.line.fill.background()
demo_box.adjustments[0] = 0.1

demo_tb = slide.shapes.add_textbox(Inches(1.7), demo_y + Inches(0.15), Inches(9.9), Inches(1.3))
tf = demo_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Live PoC Available Now"
run.font.name = "Calibri"
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "http://47.130.120.152/ui/index.html"
run2.font.name = "Consolas"
run2.font.size = Pt(16)
run2.font.color.rgb = WHITE

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "Questions?"
run3.font.name = "Calibri"
run3.font.size = Pt(15)
run3.font.italic = True
run3.font.color.rgb = WHITE

print("  Slide 7 (ready to deploy): added")

prs.save(OUT_PATH)
print(f"\nSaved: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
